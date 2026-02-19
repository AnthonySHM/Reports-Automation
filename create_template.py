"""
Generate the reusable Mantix4-style PPTX template.

Produces a single .pptx file with 18 slides that mirror the structure of the
RDLR vulnerability report but contain NO client-specific data.  Every "(cont.)"
slide is collapsed into the first slide of its section, and all numbers / IPs /
CVEs / names are replaced with placeholders.

The visual style faithfully reproduces the Mantix4 Managed Security Service
report: dark content panels, green accent lines, branded header/footer bars,
coloured stat-card rings, and alternating-blue tables.

Run:
    python create_template.py

Output:
    templates/mantix4_report_template.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Dimensions & paths
# ---------------------------------------------------------------------------
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

LOGO_PATH = Path("assets/mantix4_logo.png")
OUTPUT_DIR = Path("templates")

# ---------------------------------------------------------------------------
# Colour palette  (extracted from the reference PDF)
# ---------------------------------------------------------------------------
CLR_BG        = RGBColor(0xFF, 0xFF, 0xFF)   # slide background
CLR_DARK      = RGBColor(0x1B, 0x21, 0x2C)   # dark charcoal
CLR_BLACK     = RGBColor(0x00, 0x00, 0x00)
CLR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

CLR_GREEN     = RGBColor(0x38, 0x76, 0x1D)   # accent line
CLR_GREEN_BR  = RGBColor(0x18, 0x80, 0x3B)   # cover border

CLR_PANEL     = RGBColor(0xD6, 0xDB, 0xDF)   # light grey content panel (was dark 0x1B, 0x21, 0x2C)
CLR_PANEL_LT  = RGBColor(0xE8, 0xEB, 0xEE)   # even lighter panel (was 0x26, 0x2E, 0x3D)

CLR_STAT_GRN  = RGBColor(0x2E, 0x8B, 0x57)   # Raw Telemetry (darker green for visibility on light bg)
CLR_STAT_PCH  = RGBColor(0xD2, 0x69, 0x1E)   # Key Detections (darker orange)
CLR_STAT_PNK  = RGBColor(0x8B, 0x00, 0x8B)   # Analyst Investigations (darker purple)
CLR_STAT_BLU  = RGBColor(0x41, 0x69, 0xE1)   # SOC Tickets (darker blue)

CLR_TBL_DARK  = RGBColor(0x0B, 0x53, 0x94)   # table row dark blue
CLR_TBL_LIGHT = RGBColor(0x3D, 0x85, 0xC6)   # table row light blue
CLR_TBL_BDR   = RGBColor(0x00, 0x34, 0x49)   # table border

CLR_GREY      = RGBColor(0xBF, 0xBF, 0xBF)
CLR_MUTED     = RGBColor(0x99, 0x99, 0x99)

# Layout constants (in inches, properly centered)
MARGIN_L      = 0.38   # ~27 pts
MARGIN_R      = 0.49   # ~35 pts
LINE_Y        = 0.94   # green accent line Y
TITLE_BAR_Y   = 0.13   # title banner top
TITLE_BAR_H   = 0.58   # title banner height
BODY_TOP      = 1.10   # content start Y
BODY_H        = 3.98   # content height (extends to footer top)
FOOTER_Y      = 5.08   # footer strip top
FOOTER_H      = 0.55   # footer strip height
CONTENT_W     = 7.12   # content panel width
CONTENT_L     = (10 - CONTENT_W) / 2   # content panel left (centered) = 1.44"

# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, *, fill=None, line_clr=None, line_w=None, name=None):
    """Add a plain rectangle shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                               Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_clr:
        s.line.color.rgb = line_clr
        s.line.width = Pt(line_w or 1)
    else:
        s.line.fill.background()
    if name:
        s.name = name
    return s


def _line(slide, x1, y1, x2, y2, *, color=CLR_GREEN, width=3):
    """Horizontal/vertical line via a connector-like thin rectangle."""
    # Use a freeform is complex; a thin rectangle is simpler
    thickness = Pt(width)
    if y1 == y2:  # horizontal
        r = _rect(slide, x1, y1, x2 - x1, 0.01, fill=color, name="accent_line")
        r.height = thickness
    else:
        r = _rect(slide, x1, y1, 0.01, y2 - y1, fill=color)
        r.width = thickness
    r.line.fill.background()
    return r


def _tb(slide, l, t, w, h, text, *, sz=14, bold=False, clr=CLR_BLACK,
        align=PP_ALIGN.LEFT, font="Calibri", name=None, anchor=MSO_ANCHOR.TOP):
    """Add a text box and return it."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Set vertical alignment
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        anchor_map = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr",
                      MSO_ANCHOR.BOTTOM: "b"}
        bodyPr.set("anchor", anchor_map.get(anchor, "t"))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = clr
    p.font.name = font
    p.alignment = align
    return box


def _add_para(tf, text, *, sz=12, bold=False, clr=CLR_BLACK, font="Calibri",
              align=PP_ALIGN.LEFT, level=0, spc_before=4, spc_after=2,
              bullet_char=None):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = clr
    p.font.name = font
    p.alignment = align
    p.level = level
    p.space_before = Pt(spc_before)
    p.space_after = Pt(spc_after)
    if bullet_char:
        pPr = p._p.get_or_add_pPr()
        # Set proper hanging-indent margins for the bullet level
        indent_emu = 228600                         # 0.25" hanging indent
        margin_emu = indent_emu + (level * 457200)  # +0.5" per level
        pPr.set("marL", str(margin_emu))
        pPr.set("indent", str(-indent_emu))
        # Bullet font and character
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "+mj-lt")
        bc = etree.SubElement(pPr, qn("a:buChar"))
        bc.set("char", bullet_char)
    elif level > 0:
        # Indented sub-item without visible bullet
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(level * 457200))
    return p


def _oval(slide, l, t, w, h, *, stroke_clr, stroke_w=3, name=None):
    """Add a circle/oval outline (no fill)."""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t),
                               Inches(w), Inches(h))
    s.fill.background()
    s.line.color.rgb = stroke_clr
    s.line.width = Pt(stroke_w)
    if name:
        s.name = name
    return s


def _add_arrow(slide, x1, y1, x2, y2, color, width=2, name=None):
    """Add a horizontal arrow connector between two points with arrowhead."""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Pt
    from lxml import etree
    
    # Add a straight connector
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    
    # Style the line
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    
    # Add arrowhead at the end by manipulating XML directly
    # python-pptx doesn't expose end_arrow_type, so we add it via XML
    ln = connector.line._ln
    
    # Define namespace
    a_ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    
    # Create tailEnd element (the arrowhead at the end of the line)
    tail_end = etree.Element(f'{a_ns}tailEnd')
    tail_end.set('type', 'triangle')  # triangle is the standard arrow shape
    
    # Add the tailEnd element to the line XML
    ln.append(tail_end)
    
    if name:
        connector.name = name
    
    return connector


# ---------------------------------------------------------------------------
# Composite decorators
# ---------------------------------------------------------------------------

def _add_footer(slide, slide_num=""):
    """Footer bar: dark strip with logo + page number (every content slide)."""
    # Dark footer strip
    _rect(slide, 0, FOOTER_Y, 10, FOOTER_H, fill=CLR_DARK, name="footer_bar")
    # Logo
    if LOGO_PATH.exists():
        logo_w = Inches(1.3)
        logo_h = Inches(0.34)
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(0.38), Inches(FOOTER_Y + 0.10),
            logo_w, logo_h,
        )
    # Page number
    _tb(slide, 9.0, FOOTER_Y + 0.05, 0.7, 0.4,
        slide_num, sz=10, clr=CLR_WHITE, align=PP_ALIGN.RIGHT,
        name="page_number")


def _add_header(slide, title, *, wide=False):
    """Title banner: centred dark bar with white title text + green accent line."""
    bw = 5.5 if not wide else 7.5
    bl = (10 - bw) / 2
    # Title background bar
    _rect(slide, bl, TITLE_BAR_Y, bw, TITLE_BAR_H, fill=CLR_BLACK,
          name="title_bar")
    # Title text
    _tb(slide, bl, TITLE_BAR_Y + 0.05, bw, TITLE_BAR_H - 0.05,
        title, sz=20, bold=True, clr=CLR_WHITE, align=PP_ALIGN.CENTER,
        name="slide_title", anchor=MSO_ANCHOR.MIDDLE)
    # Green accent line
    _line(slide, MARGIN_L, LINE_Y, 10 - MARGIN_R, LINE_Y,
          color=CLR_GREEN, width=3)


def _content_panel(slide, *, left=CONTENT_L, top=BODY_TOP,
                   width=CONTENT_W, height=BODY_H, name="content_panel"):
    """Dark rounded-corner content panel for body text/images."""
    return _rect(slide, left, top, width, height, fill=CLR_PANEL, name=name)


def _placeholder_area(slide, l, t, w, h, label, *, name=None):
    """Dashed-border placeholder rectangle for charts/images."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.background()
    s.line.color.rgb = CLR_GREY
    s.line.width = Pt(1)
    # Set dashed line
    ln = s._element.find(qn("a:ln") if s._element.find(qn("a:ln")) is not None
                         else ".//{%s}ln" % "http://schemas.openxmlformats.org/drawingml/2006/main")
    if ln is None:
        spPr = s._element.find(qn("p:spPr"))
        if spPr is None:
            spPr = s._element.find(".//{%s}spPr" % "http://schemas.openxmlformats.org/drawingml/2006/main")
        if spPr is not None:
            ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.set("dash", "dash")
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    if name:
        s.name = name
    tf = s.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Vertical centre the label
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "ctr")
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = CLR_GREY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return s


# =====================================================================
#  SLIDE BUILDERS
# =====================================================================

def slide_cover(prs):
    """Slide 1 — Cover / title."""
    slide = _blank(prs)

    # Full-page white
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)

    # Large centred light panel with green border
    _rect(slide, 0.93, 1.27, 8.14, 3.54,
          fill=CLR_PANEL, line_clr=CLR_GREEN_BR, line_w=2,
          name="cover_border")

    # Mantix4 logo centred inside the panel
    if LOGO_PATH.exists():
        logo_w = Inches(3.1)
        logo_h = Inches(0.80)
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(3.45), Inches(1.55),
            logo_w, logo_h,
        )

    # "Managed Security Service"
    _tb(slide, 1.5, 2.55, 7.0, 0.45,
        "Managed Security Service",
        sz=20, bold=True, clr=CLR_DARK, align=PP_ALIGN.CENTER,
        name="cover_title")

    # Client name
    _tb(slide, 1.5, 3.05, 7.0, 0.45,
        "[Client Name]",
        sz=20, bold=True, clr=CLR_DARK, align=PP_ALIGN.CENTER,
        name="client_name")

    # Reporting period
    _tb(slide, 1.5, 3.55, 7.0, 0.38,
        "Reporting Period: [Start Date] to [End Date]",
        sz=14, clr=CLR_DARK, align=PP_ALIGN.CENTER,
        name="reporting_period")

    # Confidential
    _tb(slide, 1.5, 4.10, 7.0, 0.38,
        "CLIENT CONFIDENTIAL",
        sz=14, clr=CLR_DARK, align=PP_ALIGN.CENTER,
        name="confidential_label")

    # Page number bottom-right (small dark bar)
    _rect(slide, 9.27, 5.10, 0.60, 0.43, fill=CLR_BLACK)
    _tb(slide, 9.27, 5.12, 0.55, 0.35,
        "1", sz=10, clr=CLR_WHITE, align=PP_ALIGN.CENTER,
        name="page_number")


def slide_intro(prs):
    """Slide 2 — Report Introduction."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Report Introduction")

    # Content panel (centered)
    _content_panel(slide)

    # Body text (dark on light) - centered with padding
    text_padding = 0.20
    box = _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.15, 
              CONTENT_W - (2 * text_padding), BODY_H - 0.30, "",
              sz=13, clr=CLR_DARK, name="intro_body")
    tf = box.text_frame

    bullets = [
        "\u2022  This monthly cyber security report highlights areas that require attention "
        "and / or mitigation actions that can help to advance the security posture of "
        "your environment.",
        "\u2022  The data presented herein is pulled from Mantix4\u2019s Managed Security Service "
        "sources (SIEM, SentinelOne, NDR, and SOC), offering a comprehensive view of "
        "your organization\u2019s defenses.",
        "\u2022  Please note that certain risk mitigations (e.g., software patching) must be "
        "done with caution to ensure they do not impact service availability or integrity.",
        "\u2022  The raw statistics accompanying this report give deeper insights into the "
        "current security posture and ongoing improvements.",
        "\u2022  The NDR reports reflect filtered network data in an attempt to reduce noise "
        "and highlight anomalies.",
    ]

    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = CLR_DARK
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        p.space_after = Pt(4)

    _add_footer(slide, "2")


def slide_cyber_insight(prs):
    """Slide 3 — Cyber Insight Summary."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Cyber Insight Summary")

    _content_panel(slide)

    text_padding = 0.20
    box = _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.12, 
              CONTENT_W - (2 * text_padding), BODY_H - 0.25, "",
              sz=12, clr=CLR_DARK, name="cyber_insight_body")
    tf = box.text_frame

    # Intro
    p0 = tf.paragraphs[0]
    p0.text = ("The following key host-based observations and actions have been "
               "identified during the reporting period:")
    p0.font.size = Pt(12)
    p0.font.color.rgb = CLR_DARK
    p0.font.name = "Calibri"
    p0.space_after = Pt(10)

    items = [
        "\u2022  Critical Patch Management Findings:",
        "     \u2013  Windows Systems: [Count] missing Microsoft KB patches identified "
         "across the network.",
        "     \u2013  Third-Party Software: [Count] applications require updates to "
         "address potential vulnerabilities.",
        "     \u2013  Recommendation: Prioritize the deployment of these updates (with "
         "caution) to mitigate security risks and maintain system integrity.",
        "",
        "\u2022  [Additional observations and recommendations to be populated]",
    ]
    for i, text in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(12)
            if i == 0 or i == 5:  # Main bullets are bold
                p.font.bold = True
            else:
                p.font.bold = False
            p.font.color.rgb = CLR_DARK
            p.font.name = "Calibri"
            p.space_before = Pt(3)
            p.space_after = Pt(2)
        else:
            para = tf.add_paragraph()
            para.text = text
            para.font.size = Pt(12)
            if i == 5:  # Last main bullet is bold
                para.font.bold = True
            else:
                para.font.bold = False
            para.font.color.rgb = CLR_DARK
            para.font.name = "Calibri"
            para.space_before = Pt(3)
            para.space_after = Pt(2)

    _add_footer(slide, "3")


def slide_service_coverage(prs):
    """Slide 4 — Service Coverage & Deployment."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Service Coverage & Deployment", wide=True)

    _content_panel(slide)

    text_padding = 0.22
    box = _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.08, 
              CONTENT_W - (2 * text_padding), BODY_H - 0.20, "",
              sz=11, clr=CLR_DARK, name="service_coverage_body")
    tf = box.text_frame

    lines = [
        "\u2022  Endpoint Deployment ([Count] systems detected):",
        "     a.  EPIC SIEM:  [Count] agents",
        "     b.  SentinelOne XDR:  [Count] agents",
        "     c.  DFIR Agent:  [Count] agents",
        "\u2022  Network Detection & Response (NDR):  Active",
        "\u2022  MS365 Cloud Coverage:",
        "     a.  Active user accounts:  [Count]",
        "\u2022  Vulnerability Management:",
        "     a.  Host-Based (OS / 3rd Party SW):  [Count] systems reporting",
        "     b.  Network-Based (Perimeter):  Client Managed",
        "     c.  Internal Network:  Subscribed",
        "\u2022  Security Operations Center (SOC) 24x7:  Active",
    ]
    for i, text in enumerate(lines):
        # Main bullets start with bullet character
        is_main_bullet = text.startswith("\u2022")
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # For lines with status values (Active, Client Managed, Subscribed),
        # split into bold label and non-bold value
        if is_main_bullet and ("Active" in text or "Client Managed" in text or "Subscribed" in text):
            # Split on the colon to separate label from value
            parts = text.split(":", 1)
            if len(parts) == 2:
                label = parts[0] + ":"
                value = parts[1]
                
                # Add bold label
                run1 = p.add_run()
                run1.text = label
                run1.font.size = Pt(11)
                run1.font.bold = True
                run1.font.color.rgb = CLR_DARK
                run1.font.name = "Calibri"
                
                # Add non-bold value
                run2 = p.add_run()
                run2.text = value
                run2.font.size = Pt(11)
                run2.font.bold = False
                run2.font.color.rgb = CLR_DARK
                run2.font.name = "Calibri"
            else:
                # Fallback: use original behavior
                p.text = text
                p.font.size = Pt(11)
                p.font.bold = is_main_bullet
                p.font.color.rgb = CLR_DARK
                p.font.name = "Calibri"
        else:
            # Normal behavior for other lines
            p.text = text
            p.font.size = Pt(11)
            p.font.bold = is_main_bullet
            p.font.color.rgb = CLR_DARK
            p.font.name = "Calibri"
        
        p.space_before = Pt(4)
        p.space_after = Pt(2)

    _add_footer(slide, "4")


def slide_service_stats(prs):
    """Slide 5 — Service Statistics with four stat rings."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Service Statistics")

    # -- Telemetry breakdown (top-right panel) --
    _rect(slide, 3.40, 1.10, 4.25, 1.65, fill=CLR_PANEL, name="telemetry_panel")

    box = _tb(slide, 3.50, 1.12, 4.05, 1.55, "",
              sz=10, clr=CLR_DARK, name="telemetry_breakdown")
    tf = box.text_frame
    p0 = tf.paragraphs[0]
    p0.text = "** Raw Telemetry Source Events:"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = CLR_DARK
    p0.font.name = "Calibri"

    for line in ["\u2022  Endpoint (SIEM):  [Count]",
                 "\u2022  Network (NDR):  [Count]*",
                 "\u2022  Cloud (MS365):  [Count]"]:
        _add_para(tf, line, sz=10, clr=CLR_DARK, spc_before=3, spc_after=1)

    _add_para(tf, "*The NDR value is an estimate based on typical traffic "
              "rates in this time period", sz=7, clr=CLR_MUTED,
              spc_before=6)

    # -- Four stat cards --
    card_data = [
        ("Raw Telemetry", CLR_STAT_GRN,
         "Total number of raw events ingested into the SIEM and NDR "
         "during the reporting period.", 0.55),
        ("Key Detections", CLR_STAT_PCH,
         "Detections involving critical or high severity events, issues "
         "and behaviors.", 2.95),
        ("Analyst Investigations", CLR_STAT_PNK,
         "Investigative actions, including deeper analysis of contextual "
         "sources (e.g. DFIR).", 5.35),
        ("SOC Tickets", CLR_STAT_BLU,
         "Tickets generated by the 24x7 SOC involving notable activities "
         "requiring confirmation from the client.", 7.75),
    ]

    for label, accent, desc, left in card_data:
        key = label.lower().replace(" ", "_")
        card_top = 3.0
        card_w = 2.0

        # Background card
        _rect(slide, left, card_top, card_w, 2.0, fill=CLR_PANEL)

        # Coloured ring (oval)
        ring_sz = 0.95
        ring_l = left + (card_w - ring_sz) / 2
        _oval(slide, ring_l, card_top + 0.12, ring_sz, ring_sz,
              stroke_clr=accent, stroke_w=3, name=f"stat_{key}_ring")

        # Value inside ring
        _tb(slide, left, card_top + 0.30, card_w, 0.55,
            "[Count]", sz=13, bold=True, clr=CLR_DARK,
            align=PP_ALIGN.CENTER, name=f"stat_{key}_value",
            anchor=MSO_ANCHOR.MIDDLE)

        # Coloured label
        _tb(slide, left, card_top + 1.10, card_w, 0.28,
            label, sz=9, bold=True, clr=accent, align=PP_ALIGN.CENTER,
            name=f"stat_{key}_label")

        # Description
        _tb(slide, left + 0.08, card_top + 1.38, card_w - 0.16, 0.55,
            desc, sz=7, clr=CLR_DARK, align=PP_ALIGN.CENTER,
            name=f"stat_{key}_desc")

    # -- Add arrows between the stat cards --
    # Vertical center of the rings (card_top + ring offset + ring size/2)
    arrow_y = 3.0 + 0.12 + 0.95 / 2  # = 3.595
    
    # Arrow 1: Raw Telemetry → Key Detections (green)
    # From right edge of first card to left edge of second card
    arrow1_start = 0.55 + 2.0  # = 2.55 (right edge of Raw Telemetry card)
    arrow1_end = 2.95  # left edge of Key Detections card
    _add_arrow(slide, arrow1_start, arrow_y, arrow1_end, arrow_y, CLR_STAT_GRN, width=3)
    
    # Arrow 2: Key Detections → SOC Tickets (orange)
    arrow2_start = 2.95 + 2.0  # = 4.95 (right edge of Key Detections card)
    arrow2_end = 5.35  # left edge of SOC Tickets card
    _add_arrow(slide, arrow2_start, arrow_y, arrow2_end, arrow_y, CLR_STAT_PCH, width=3)
    
    # Arrow 3: SOC Tickets → Analyst Investigations (blue)
    arrow3_start = 5.35 + 2.0  # = 7.35 (right edge of SOC Tickets card)
    arrow3_end = 7.75  # left edge of Analyst Investigations card
    _add_arrow(slide, arrow3_start, arrow_y, arrow3_end, arrow_y, CLR_STAT_BLU, width=3)

    _add_footer(slide, "5")


def _ndr_chart_slide(prs, title, shape_name, slide_num):
    """NDR chart / image slide with grey content area below green line."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, title, wide=True)

    # Grey panel behind the chart - must start BELOW green line (0.94")
    panel_top = 1.00  # Starts below the green line at 0.94"
    panel_height = 4.00  # Reduced to fit above footer
    
    _content_panel(slide, left=CONTENT_L, width=CONTENT_W, top=panel_top, height=panel_height,
                   name="chart_panel")

    # Placeholder for chart image (well within panel to avoid edge overlap)
    image_padding = 0.12
    image_top = panel_top + 0.15  # More offset from panel top
    image_height = panel_height - 0.30  # Reduced height for top/bottom padding
    _placeholder_area(slide, CONTENT_L + image_padding, image_top, 
                      CONTENT_W - (2 * image_padding), image_height,
                      "[Chart / Image placeholder]", name=shape_name)

    _add_footer(slide, slide_num)


def slide_required_patches(prs):
    """Slide 14 — Required Software Patches."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Required Software Patches")

    _content_panel(slide)

    text_padding = 0.25
    box = _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.12,
              CONTENT_W - (2 * text_padding), 0.75, "",
              sz=13, clr=CLR_DARK, name="required_patches_body")
    tf = box.text_frame

    for i, text in enumerate([
        "\u2022  There are currently [Count] Microsoft KB patches that need to "
        "be deployed.",
        "\u2022  There are [Count] software packages that require updating. "
        "Patching should be done with caution to ensure systems and "
        "applications are not negatively impacted.",
    ]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = CLR_DARK
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        p.space_after = Pt(6)

    # Chart placeholder centered in the remaining content area
    # Text ends at ~BODY_TOP + 0.87, content panel ends at BODY_TOP + BODY_H
    # Center the chart vertically between text bottom and panel bottom
    chart_h = 2.2
    text_bottom = BODY_TOP + 0.87
    available = (BODY_TOP + BODY_H) - text_bottom
    chart_top = text_bottom + (available - chart_h) / 2
    _placeholder_area(slide, CONTENT_L + 0.25, chart_top,
                      CONTENT_W - 0.5, chart_h,
                      "[Patch trend chart placeholder]",
                      name="required_patches_chart")

    _add_footer(slide, "14")


def _table_slide(prs, title, columns, note, shape_name, slide_num):
    """Table slide with alternating-blue rows and a note bar at the bottom."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, title)

    # -- Note bar at bottom (just above footer) --
    # Footer starts at FOOTER_Y (5.08"), note should be above it
    note_height = 0.50
    note_top = FOOTER_Y - note_height - 0.08  # 8pt gap from footer
    
    _content_panel(slide, left=CONTENT_L, top=note_top, width=CONTENT_W, height=note_height,
                   name=f"{shape_name}_note_panel")
    note_padding = 0.12
    _tb(slide, CONTENT_L + note_padding, note_top + 0.05, 
        CONTENT_W - (2 * note_padding), note_height - 0.10,
        note, sz=8, clr=CLR_DARK, name=f"{shape_name}_note")

    # -- Table area (constrained to never reach the note) --
    # Table starts after the header and must end before the note with some padding
    n_cols = len(columns)
    n_rows = 6  # header + 5 placeholder rows
    
    tbl_left = Inches(0.38)
    tbl_top = Inches(1.10)
    tbl_w = Inches(9.14)
    
    # Calculate maximum table height: from table top to note panel minus padding
    max_table_height = note_top - 1.10 - 0.15  # 15pt padding before note
    tbl_h = Inches(max_table_height)
    row_h = Emu(tbl_h / n_rows)  # Distribute height evenly among rows

    shape = slide.shapes.add_table(n_rows, n_cols,
                                   tbl_left, tbl_top, tbl_w, tbl_h)
    shape.name = shape_name
    tbl = shape.table

    # -- Header row --
    for ci, col in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = col
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = CLR_WHITE
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_BLACK

    # -- Data rows (alternating blue, placeholder text) --
    for ri in range(1, n_rows):
        row_clr = CLR_TBL_DARK if ri % 2 == 1 else CLR_TBL_LIGHT
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            cell.text = "\u2014"  # em-dash placeholder
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = CLR_WHITE
                p.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_clr

    _add_footer(slide, slide_num)


def slide_internal_vulns(prs):
    """Slide 17 — Internal Network Vulnerabilities."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Internal Network Vulnerabilities", wide=True)

    _content_panel(slide)

    text_padding = 0.20
    _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.10, 
        CONTENT_W - (2 * text_padding), 0.70,
        "The internal network subnets were scanned for network "
        "vulnerabilities (not host-based software vulnerabilities) from the "
        "Mantix4 Sensor. The following were identified (Reported as "
        "[Category] [Protocol] [Severity] Issue Description):",
        sz=11, clr=CLR_DARK, name="internal_vulns_intro")

    placeholder_padding = 0.20
    _placeholder_area(slide, CONTENT_L + placeholder_padding, BODY_TOP + 0.90, 
                      CONTENT_W - (2 * placeholder_padding), 2.85,
                      "[Vulnerability findings to be populated]",
                      name="internal_vulns_content")

    _add_footer(slide, "17")


def slide_internal_mitigation(prs):
    """Slide 18 — Internal Network Mitigation."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Internal Network Mitigation", wide=True)

    _content_panel(slide)

    text_padding = 0.20
    _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.10, 
        CONTENT_W - (2 * text_padding), 0.70,
        "The following mitigation actions are recommended based on the "
        "network vulnerability findings identified during the reporting "
        "period:",
        sz=11, clr=CLR_DARK, name="internal_mitigation_intro")

    placeholder_padding = 0.20
    _placeholder_area(slide, CONTENT_L + placeholder_padding, BODY_TOP + 0.90, 
                      CONTENT_W - (2 * placeholder_padding), 2.85,
                      "[Mitigation recommendations to be populated]",
                      name="internal_mitigation_content")

    _add_footer(slide, "18")


# =====================================================================
#  BUILD
# =====================================================================

def build_template() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)                                               # 1
    slide_intro(prs)                                               # 2
    slide_cyber_insight(prs)                                       # 3
    slide_service_coverage(prs)                                    # 4
    slide_service_stats(prs)                                       # 5

    ndr = [                                                        # 6-13
        ("NDR - Top Outbound Data Destination (1st)",  "ndr_outbound_data_1",  "6"),
        ("NDR - Top Outbound Data Destination (2nd)",  "ndr_outbound_data_2",  "7"),
        ("NDR - Top IP Destinations by Connection Count", "ndr_top_ip",        "8"),
        ("NDR - Top URLs by Connection Count",         "ndr_top_urls",         "9"),
        ("NDR - Top External Destination Receiving Data", "ndr_ext_dest",     "10"),
        ("NDR - Country by Connection Count",          "ndr_country",         "11"),
        ("NDR - Beaconing Score",                      "ndr_beaconing",       "12"),
        ("NDR - Sensitive Data",                       "ndr_sensitive_data",  "13"),
    ]
    for title, name, num in ndr:
        _ndr_chart_slide(prs, title, name, num)

    slide_required_patches(prs)                                    # 14

    note_text = ("Note: All patching, particularly on servers, should be done "
                 "with caution. Patches should be tested, and they should be "
                 "assessed for compatibility prior to deployment.")

    _table_slide(prs, "Top Software Patches",                     # 15
                 ["Package Name", "Recommended Action"],
                 note_text, "software_patches_table", "15")

    _table_slide(prs, "Top Microsoft Patches",                    # 16
                 ["Package Name", "KB Patch"],
                 note_text, "ms_patches_table", "16")

    slide_internal_vulns(prs)                                      # 17
    slide_internal_mitigation(prs)                                 # 18

    OUTPUT_DIR.mkdir(exist_ok=True)
    out = OUTPUT_DIR / "mantix4_report_template.pptx"
    prs.save(str(out))
    return out


def build_multi_sensor_template(sensor_ids: list[str]) -> Path:
    """Generate a template with separate NDR slide sections per sensor.

    Slide layout:
        Slides 1-5:  Fixed intro slides (cover, intro, cyber insight,
                     service coverage, service statistics)
        Per sensor:  8 NDR slides each (outbound 1, outbound 2, top IP,
                     top URLs, ext dest, country, beaconing, sensitive data)
        After NDR:   Patches & vulnerability slides (required patches,
                     software patches table, MS patches table, internal
                     vulns, internal mitigation)

    Parameters
    ----------
    sensor_ids : list[str]
        Sorted list of sensor identifiers (e.g. ``["GAPRD", "VAHQ", "VAPRD"]``).

    Returns
    -------
    Path
        Path to the saved ``.pptx`` file.
    """
    import logging
    logger = logging.getLogger("shm.create_template")

    if not sensor_ids or sensor_ids == ["DEFAULT"]:
        logger.info("build_multi_sensor_template called with DEFAULT; using standard template")
        return build_template()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- Fixed intro slides (1-5) ---
    slide_cover(prs)                    # 1
    slide_intro(prs)                    # 2
    slide_cyber_insight(prs)            # 3
    slide_service_coverage(prs)         # 4
    slide_service_stats(prs)            # 5

    # --- NDR sections: 8 slides per sensor ---
    ndr_base_slides = [
        ("Top Outbound Data Destination (1st)",    "ndr_outbound_data_1"),
        ("Top Outbound Data Destination (2nd)",    "ndr_outbound_data_2"),
        ("Top IP Destinations by Connection Count", "ndr_top_ip"),
        ("Top URLs by Connection Count",           "ndr_top_urls"),
        ("Top External Destination Receiving Data", "ndr_ext_dest"),
        ("Country by Connection Count",            "ndr_country"),
        ("Beaconing Score",                        "ndr_beaconing"),
        ("Sensitive Data",                         "ndr_sensitive_data"),
    ]

    slide_num = 6  # next slide number (1-based display)

    for sensor_id in sensor_ids:
        suffix = sensor_id.lower()
        for base_title, base_shape_name in ndr_base_slides:
            title = f"NDR - {sensor_id} - {base_title}"
            shape_name = f"{base_shape_name}_{suffix}"
            _ndr_chart_slide(prs, title, shape_name, str(slide_num))
            slide_num += 1

        # Internal Network Vulnerabilities (per sensor)
        slide = _blank(prs)
        _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
        _add_header(slide, f"Internal Network Vulnerabilities - {sensor_id}", wide=True)
        _content_panel(slide)
        text_padding = 0.20
        _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.10,
            CONTENT_W - (2 * text_padding), 0.70,
            "The internal network subnets were scanned for network "
            "vulnerabilities (not host-based software vulnerabilities) from the "
            "Mantix4 Sensor. The following were identified (Reported as "
            "[Category] [Protocol] [Severity] Issue Description):",
            sz=11, clr=CLR_DARK, name=f"internal_vulns_intro_{suffix}")
        placeholder_padding = 0.20
        _placeholder_area(slide, CONTENT_L + placeholder_padding, BODY_TOP + 0.90,
                          CONTENT_W - (2 * placeholder_padding), 2.85,
                          "[Vulnerability findings to be populated]",
                          name=f"internal_vulns_content_{suffix}")
        _add_footer(slide, str(slide_num))
        slide_num += 1

        # Internal Network Mitigation (per sensor)
        slide = _blank(prs)
        _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
        _add_header(slide, f"Internal Network Mitigation - {sensor_id}", wide=True)
        _content_panel(slide)
        text_padding = 0.20
        _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.10,
            CONTENT_W - (2 * text_padding), 0.70,
            "The following mitigation actions are recommended based on the "
            "network vulnerability findings identified during the reporting "
            "period:",
            sz=11, clr=CLR_DARK, name=f"internal_mitigation_intro_{suffix}")
        placeholder_padding = 0.20
        _placeholder_area(slide, CONTENT_L + placeholder_padding, BODY_TOP + 0.90,
                          CONTENT_W - (2 * placeholder_padding), 2.85,
                          "[Mitigation recommendations to be populated]",
                          name=f"internal_mitigation_content_{suffix}")
        _add_footer(slide, str(slide_num))
        slide_num += 1

    # --- Patches slides ---
    # Required Software Patches (update page number)
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Required Software Patches")
    _content_panel(slide)
    text_padding = 0.25
    box = _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.12,
              CONTENT_W - (2 * text_padding), 0.75, "",
              sz=13, clr=CLR_DARK, name="required_patches_body")
    tf = box.text_frame
    for i, text in enumerate([
        "\u2022  There are currently [Count] Microsoft KB patches that need to "
        "be deployed.",
        "\u2022  There are [Count] software packages that require updating. "
        "Patching should be done with caution to ensure systems and "
        "applications are not negatively impacted.",
    ]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = CLR_DARK
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        p.space_after = Pt(6)
    # Chart placeholder centered in the remaining content area
    chart_h = 2.2
    text_bottom = BODY_TOP + 0.87
    available = (BODY_TOP + BODY_H) - text_bottom
    chart_top = text_bottom + (available - chart_h) / 2
    _placeholder_area(slide, CONTENT_L + 0.25, chart_top,
                      CONTENT_W - 0.5, chart_h,
                      "[Patch trend chart placeholder]",
                      name="required_patches_chart")
    _add_footer(slide, str(slide_num))
    slide_num += 1

    note_text = ("Note: All patching, particularly on servers, should be done "
                 "with caution. Patches should be tested, and they should be "
                 "assessed for compatibility prior to deployment.")

    _table_slide(prs, "Top Software Patches",
                 ["Package Name", "Recommended Action"],
                 note_text, "software_patches_table", str(slide_num))
    slide_num += 1

    _table_slide(prs, "Top Microsoft Patches",
                 ["Package Name", "KB Patch"],
                 note_text, "ms_patches_table", str(slide_num))
    slide_num += 1

    # Internal Network Vulnerabilities
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Internal Network Vulnerabilities", wide=True)
    _content_panel(slide)
    text_padding = 0.20
    _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.10,
        CONTENT_W - (2 * text_padding), 0.70,
        "The internal network subnets were scanned for network "
        "vulnerabilities (not host-based software vulnerabilities) from the "
        "Mantix4 Sensor. The following were identified (Reported as "
        "[Category] [Protocol] [Severity] Issue Description):",
        sz=11, clr=CLR_DARK, name="internal_vulns_intro")
    placeholder_padding = 0.20
    _placeholder_area(slide, CONTENT_L + placeholder_padding, BODY_TOP + 0.90,
                      CONTENT_W - (2 * placeholder_padding), 2.85,
                      "[Vulnerability findings to be populated]",
                      name="internal_vulns_content")
    _add_footer(slide, str(slide_num))
    slide_num += 1

    # Internal Network Mitigation
    slide = _blank(prs)
    _rect(slide, 0, 0, 10, 5.625, fill=CLR_BG)
    _add_header(slide, "Internal Network Mitigation", wide=True)
    _content_panel(slide)
    text_padding = 0.20
    _tb(slide, CONTENT_L + text_padding, BODY_TOP + 0.10,
        CONTENT_W - (2 * text_padding), 0.70,
        "The following mitigation actions are recommended based on the "
        "network vulnerability findings identified during the reporting "
        "period:",
        sz=11, clr=CLR_DARK, name="internal_mitigation_intro")
    placeholder_padding = 0.20
    _placeholder_area(slide, CONTENT_L + placeholder_padding, BODY_TOP + 0.90,
                      CONTENT_W - (2 * placeholder_padding), 2.85,
                      "[Mitigation recommendations to be populated]",
                      name="internal_mitigation_content")
    _add_footer(slide, str(slide_num))

    total_slides = slide_num
    OUTPUT_DIR.mkdir(exist_ok=True)
    # Use a distinct filename to avoid overwriting the standard template
    filename = "mantix4_multi_sensor_template.pptx"
    out = OUTPUT_DIR / filename
    prs.save(str(out))

    logger.info(
        "Multi-sensor template saved: %s (%d sensors, %d total slides)",
        out, len(sensor_ids), total_slides,
    )
    return out


if __name__ == "__main__":
    path = build_template()
    print(f"Template saved to: {path}")
    print("Total slides: 18")
