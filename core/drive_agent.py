"""
Google Drive Agent — Fetch NDR chart images from a shared Drive folder.

Navigates the shared Drive structure:

    Root (shared folder)
    +-- <Client Name>
        +-- NDR-<client>-<date>   (most recent by date)
            +-- image1.png
            +-- image2.png
            +-- ...

Downloads images and maps them to the correct NDR slides in the
Mantix4 PPTX template.

Authentication
--------------
Uses a Google Cloud service-account JSON key.  Set the path via the
``SHM_GOOGLE_CREDENTIALS`` environment variable, or pass it directly
to the ``DriveAgent`` constructor.
"""

from __future__ import annotations

import copy
import csv
import io
import logging
import math
import os
import re
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional

from google.oauth2.service_account import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from PIL import Image as PILImage
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from core.slide_utils import decorate_slide, insert_slide_at
from core.transform import Rect, compute_contain_transform

logger = logging.getLogger("shm.drive_agent")

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
ROOT_FOLDER_ID = "1pLBfcoLHtksS53KCRpFX0jekqNtxm_82"
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff"}
CSV_EXTENSIONS = {".csv"}
TXT_EXTENSIONS = {".txt"}

# Frontend slug -> search terms for matching Drive client folders
CLIENT_FOLDER_MAP: dict[str, list[str]] = {
    "montereymech":  ["monterey mechanical", "monterey mech", "montereymech"],
    "elephant":      ["elephant"],
    "cic_xerox":     ["cic"],
    "wiegmann":      ["wiegmann"],
    "saints":        ["saints joachim", "saints"],
    "rdlr":          ["rdlr"],
    "jorgensen_labs": ["jorgensen"],
    "gatedrentals":  ["gated rentals", "gatedrentals", "gated"],
    "cjfig":         ["cjfig"],
    "packard":       ["packard"],
}

# Display names (from frontend dropdown) -> slug for Drive lookup
DISPLAY_NAME_TO_SLUG: dict[str, str] = {
    "monterey mechanical - wt": "montereymech",
    "monterey mechanical": "montereymech",
    "elephant - 1st of the month": "elephant",
    "elephant": "elephant",
    "cic - x (but x is not providing it)": "cic_xerox",
    "cic - x": "cic_xerox",
    "wiegmann - wt": "wiegmann",
    "wiegmann": "wiegmann",
    "the saints joachim and anne - wt": "saints",
    "saints": "saints",
    "rdlr - direct": "rdlr",
    "rdlr": "rdlr",
    "jorgensen labs - x": "jorgensen_labs",
    "jorgensen labs": "jorgensen_labs",
    "gated rentals - direct": "gatedrentals",
    "gated rentals": "gatedrentals",
    "cjfig - direct": "cjfig",
    "cjfig": "cjfig",
    "packard - direct (you got it)": "packard",
    "packard - direct": "packard",
    "packard": "packard",
}

# NDR image filename keywords -> (slide_index, shape_name)
NDR_FILE_MAP: dict[str, tuple[int, str]] = {
    "outbound_data_1":  (5, "ndr_outbound_data_1"),
    "outbound_data_2":  (6, "ndr_outbound_data_2"),
    "top_outbound_1":   (5, "ndr_outbound_data_1"),
    "top_outbound_2":   (6, "ndr_outbound_data_2"),
    "heatmap":          (5, "ndr_outbound_data_1"),
    "top_ip":           (7, "ndr_top_ip"),
    "ip_dest":          (7, "ndr_top_ip"),
    "top_url":          (8, "ndr_top_urls"),
    "url":              (8, "ndr_top_urls"),
    "ext_dest":         (9, "ndr_ext_dest"),
    "external_dest":    (9, "ndr_ext_dest"),
    "country":          (10, "ndr_country"),
    "beacon":           (11, "ndr_beaconing"),
    "sensitive":        (12, "ndr_sensitive_data"),
}

# Ordered slots for positional fallback
NDR_SLIDES_ORDERED: list[tuple[int, str]] = [
    (5,  "ndr_outbound_data_1"),
    (6,  "ndr_outbound_data_2"),
    (7,  "ndr_top_ip"),
    (8,  "ndr_top_urls"),
    (9,  "ndr_ext_dest"),
    (10, "ndr_country"),
    (11, "ndr_beaconing"),
    (12, "ndr_sensitive_data"),
]

# Heatmap / outbound images are assigned to these slides in file order (1st → slide 5, 2nd → slide 6).
OUTBOUND_SLIDES: list[tuple[int, str]] = [
    (5, "ndr_outbound_data_1"),
    (6, "ndr_outbound_data_2"),
]
# Match filenames with/without spaces/underscores: "heatmap", "heat map", "Outbound 1", etc.
OUTBOUND_STEM_KEYWORDS = (
    "heatmap", "heat_map", "heat_map_1", "heat_map_2",
    "outbound_data_1", "outbound_data_2", "top_outbound_1", "top_outbound_2",
    "outbound_1", "outbound_2", "outbound_1st", "outbound_2nd",
    "outbound_data", "top_outbound",
)

# NDR placeholder box (inches) — centered on slide, positioned below green line
CONTENT_WIDTH = 7.12
SLIDE_WIDTH = 10.0
CONTENT_LEFT = (SLIDE_WIDTH - CONTENT_WIDTH) / 2  # 1.44"
IMAGE_PADDING = 0.12  # padding inside content panel
PANEL_TOP = 1.00  # Panel starts below green line at 0.94"
PANEL_HEIGHT = 4.00

NDR_BOX_INCHES = {
    "left": CONTENT_LEFT + IMAGE_PADDING,   # 1.56"
    "top": PANEL_TOP + 0.15,  # 1.15" - clearly within the grey panel
    "width": CONTENT_WIDTH - (2 * IMAGE_PADDING),  # 6.88"
    "height": PANEL_HEIGHT - 0.30  # 3.70" - fits comfortably within panel with top/bottom padding
}

# Table box: positioned below green line with increased padding to stay within grey panel
TABLE_BOX_INCHES = {
    "left": CONTENT_LEFT + IMAGE_PADDING,   # 1.56"
    "top": PANEL_TOP + 0.20,  # 1.20" - more padding from panel top (panel at 1.00")
    "width": CONTENT_WIDTH - (2 * IMAGE_PADDING),  # 6.88"
    "height": 2.8  # Reduced to ensure it stays well within grey background with padding
}

# Pagination
# Maximum data rows per table page (reduced to fit in 2.8" height)
MAX_ROWS_PER_PAGE = 7

# Table styling colours (must match create_template.py)
CLR_BLACK     = RGBColor(0x00, 0x00, 0x00)
CLR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CLR_TBL_DARK  = RGBColor(0x0B, 0x53, 0x94)
CLR_TBL_LIGHT = RGBColor(0x3D, 0x85, 0xC6)
CLR_MUTED     = RGBColor(0x99, 0x99, 0x99)
EMU_PER_INCH  = 914_400

# Table-slide mapping: slide_title -> (slide_index, shape_name).
# Template: slide 14 = Top Software Patches, slide 15 = Top Microsoft Patches (0-based).
# Keys are normalised (lowercase, spaces/punctuation -> underscores).
TABLE_SLIDE_MAP: dict[str, tuple[int, str]] = {
    "sensitive_data":                  (12, "ndr_sensitive_data"),
    "beaconing_score":                 (11, "ndr_beaconing"),
    "external_receivers":              (9,  "ndr_ext_dest"),
    "domain_by_connection_count":      (8,  "ndr_top_urls"),
    "destination_by_connection_count": (7,  "ndr_top_ip"),
    "country_by_connection_count":     (10, "ndr_country"),
    "top_software_patches":            (14, "software_patches_table"),
    "top_microsoft_patches":           (15, "ms_patches_table"),
    "software_patches":                (14, "software_patches_table"),
    "microsoft_patches":               (15, "ms_patches_table"),
    "software_patch":                 (14, "software_patches_table"),
    "microsoft_patch":                (15, "ms_patches_table"),
    "patches_software":               (14, "software_patches_table"),
    "patches_microsoft":              (15, "ms_patches_table"),
}

# Column names that indicate "source" — case-insensitive substring match.
SOURCE_COLUMN_PATTERNS: list[str] = [
    "source", "sources", "data source", "data_source",
    "reference", "references", "origin",
]


# ---------------------------------------------------------------------------
# Data & Exceptions
# ---------------------------------------------------------------------------
@dataclass
class NDRAsset:
    """A downloaded NDR image and its target slide."""
    local_path: Path
    slide_index: int
    shape_name: str
    original_filename: str
    sensor_id: Optional[str] = None  # Sensor identifier (e.g., "VAPRD", "VAHQ")


@dataclass
class CSVTableAsset:
    """A downloaded CSV file and its target table slide."""
    local_path: Path
    slide_index: int
    shape_name: str
    original_filename: str
    sensor_id: Optional[str] = None  # Sensor identifier (e.g., "VAPRD", "VAHQ")


@dataclass
class SensorAssetGroup:
    """Assets for a single sensor."""
    sensor_id: str  # e.g., "VAPRD", "DEFAULT" for single-sensor clients
    ndr_images: list[NDRAsset] = field(default_factory=list)
    csv_tables: list[CSVTableAsset] = field(default_factory=list)


class DriveAgentError(Exception):
    """Base error for Drive agent operations."""


class ClientFolderNotFound(DriveAgentError):
    """No matching client folder found in Drive."""


class NDRFolderNotFound(DriveAgentError):
    """No NDR subfolder found for the client."""


class VNFolderNotFound(DriveAgentError):
    """No VN subfolder found for the client."""


# ---------------------------------------------------------------------------
# Agent
# ---------------------------------------------------------------------------
class DriveAgent:
    """Fetches NDR chart images from a shared Google Drive folder.

    Parameters
    ----------
    credentials_path : str | Path | None
        Path to the service-account JSON key file.
        Falls back to ``SHM_GOOGLE_CREDENTIALS`` env var.
    root_folder_id : str
        The shared Drive folder ID.
    """

    def __init__(
        self,
        credentials_path: Optional[str | Path] = None,
        root_folder_id: str = ROOT_FOLDER_ID,
    ):
        creds_path = credentials_path or os.getenv("SHM_GOOGLE_CREDENTIALS")
        if not creds_path:
            raise DriveAgentError(
                "No credentials path provided. Set SHM_GOOGLE_CREDENTIALS "
                "env var or pass credentials_path."
            )
        creds_path = Path(creds_path)
        if not creds_path.exists():
            raise DriveAgentError(f"Credentials file not found: {creds_path}")

        self.credentials = Credentials.from_service_account_file(
            str(creds_path), scopes=SCOPES
        )
        self.service = build("drive", "v3", credentials=self.credentials)
        self.root_folder_id = root_folder_id
        logger.info("Drive agent initialised (root=%s)", root_folder_id)

    # ------------------------------------------------------------------ #
    # Low-level Drive helpers
    # ------------------------------------------------------------------ #
    def _list_folders(self, parent_id: str) -> list[dict]:
        """List all sub-folders in *parent_id*."""
        results: list[dict] = []
        page_token = None
        query = (
            f"'{parent_id}' in parents "
            f"and mimeType='application/vnd.google-apps.folder' "
            f"and trashed=false"
        )
        while True:
            resp = self.service.files().list(
                q=query,
                fields="nextPageToken, files(id, name)",
                pageSize=100,
                pageToken=page_token,
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            ).execute()
            results.extend(resp.get("files", []))
            page_token = resp.get("nextPageToken")
            if not page_token:
                break
        return results

    def _list_image_files(self, folder_id: str) -> list[dict]:
        """List all image files in *folder_id*."""
        results: list[dict] = []
        page_token = None
        query = (
            f"'{folder_id}' in parents "
            f"and mimeType!='application/vnd.google-apps.folder' "
            f"and trashed=false"
        )
        while True:
            resp = self.service.files().list(
                q=query,
                fields="nextPageToken, files(id, name, mimeType, modifiedTime)",
                pageSize=100,
                pageToken=page_token,
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            ).execute()
            results.extend(resp.get("files", []))
            page_token = resp.get("nextPageToken")
            if not page_token:
                break

        # Keep only actual image files (by extension)
        return [
            f for f in results
            if Path(f["name"]).suffix.lower() in IMAGE_EXTENSIONS
        ]

    def _list_csv_files(self, folder_id: str) -> list[dict]:
        """List all CSV files in *folder_id* and in common data subfolders."""
        seen_ids: set[str] = set()
        csv_files: list[dict] = []

        def add_unique(files: list[dict]) -> None:
            for f in files:
                if f["id"] not in seen_ids:
                    seen_ids.add(f["id"])
                    csv_files.append(f)

        add_unique(self._list_files_by_ext(folder_id, CSV_EXTENSIONS))

        data_subfolders = (
            "tables", "table", "csvs", "csv", "data", "reports", "exports",
        )
        for sub in self._list_folders(folder_id):
            if sub["name"].lower() in data_subfolders:
                add_unique(self._list_files_by_ext(sub["id"], CSV_EXTENSIONS))

        return csv_files

    def _list_files_by_ext(
        self, folder_id: str, extensions: set[str]
    ) -> list[dict]:
        """List files in *folder_id* whose extension is in *extensions*."""
        results: list[dict] = []
        page_token = None
        query = (
            f"'{folder_id}' in parents "
            f"and mimeType!='application/vnd.google-apps.folder' "
            f"and trashed=false"
        )
        while True:
            resp = self.service.files().list(
                q=query,
                fields="nextPageToken, files(id, name, mimeType, modifiedTime)",
                pageSize=100,
                pageToken=page_token,
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            ).execute()
            results.extend(resp.get("files", []))
            page_token = resp.get("nextPageToken")
            if not page_token:
                break

        return [
            f for f in results
            if Path(f["name"]).suffix.lower() in extensions
        ]

    def _download_file(self, file_id: str, dest_path: Path) -> Path:
        """Download a Drive file to *dest_path*."""
        request = self.service.files().get_media(fileId=file_id)
        dest_path.parent.mkdir(parents=True, exist_ok=True)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        dest_path.write_bytes(fh.getvalue())
        logger.info("Downloaded %s -> %s", file_id, dest_path)
        return dest_path

    # ------------------------------------------------------------------ #
    # Folder navigation
    # ------------------------------------------------------------------ #
    def find_client_folder(self, client_slug: str) -> dict:
        """Find the Drive folder matching a frontend client slug or display name.

        Raises ``ClientFolderNotFound`` if no match.
        """
        raw = client_slug.strip()
        normalized = raw.lower()
        slug = DISPLAY_NAME_TO_SLUG.get(normalized)
        if slug is None:
            for display_name, s in DISPLAY_NAME_TO_SLUG.items():
                if display_name in normalized or normalized in display_name:
                    slug = s
                    break
        if slug is not None:
            search_terms = CLIENT_FOLDER_MAP.get(slug, [slug])
            logger.info("Resolved display name '%s' -> slug '%s'", raw, slug)
        else:
            search_terms = CLIENT_FOLDER_MAP.get(normalized, [normalized, raw])
        folders = self._list_folders(self.root_folder_id)
        logger.info(
            "Found %d client folders in root. Searching for '%s' ...",
            len(folders), client_slug,
        )

        for folder in folders:
            name_lower = folder["name"].lower()
            for term in search_terms:
                if term.lower() in name_lower:
                    logger.info(
                        "Matched client '%s' -> folder '%s' (%s)",
                        client_slug, folder["name"], folder["id"],
                    )
                    return folder

        available = [f["name"] for f in folders]
        raise ClientFolderNotFound(
            f"No folder found for client '{client_slug}'. "
            f"Available folders: {available}"
        )

    def find_client_pptx(self, client_folder_id: str, client_slug: str,
                         dest_dir: Path | str = "assets/tmp") -> Path | None:
        """Find and download the most recently modified PPTX from a client folder.

        Searches the client folder (top-level only) for ``.pptx`` files and
        downloads the most recently modified one to a local temp path.

        Parameters
        ----------
        client_folder_id : str
            Google Drive folder ID for the client.
        client_slug : str
            Client identifier used for the local filename.
        dest_dir : Path | str
            Directory to download the reference PPTX into.

        Returns
        -------
        Path | None
            Path to the downloaded PPTX, or ``None`` if no PPTX found.
        """
        dest_dir = Path(dest_dir)
        pptx_files = self._list_files_by_ext(client_folder_id, {".pptx"})
        if not pptx_files:
            logger.warning(
                "No PPTX files found in client folder %s", client_folder_id,
            )
            return None

        # Sort by modifiedTime descending to pick the most recent
        pptx_files.sort(
            key=lambda f: f.get("modifiedTime", ""), reverse=True,
        )
        latest = pptx_files[0]
        logger.info(
            "Found %d PPTX file(s) in client folder; using '%s' (modified %s)",
            len(pptx_files), latest["name"], latest.get("modifiedTime", "?"),
        )

        dest_path = dest_dir / f"{client_slug}_reference.pptx"
        return self._download_file(latest["id"], dest_path)

    def find_client_pptx_files(
        self,
        client_folder_id: str,
        client_slug: str,
        dest_dir: Path | str = "assets/tmp",
        count: int = 2,
    ) -> list[dict]:
        """Download the *count* most recent PPTX files from a client folder.

        Searches only the top-level client folder (Archive is excluded).

        Parameters
        ----------
        client_folder_id : str
            Google Drive folder ID for the client.
        client_slug : str
            Client identifier used for local filenames.
        dest_dir : Path | str
            Directory to download files into.
        count : int
            Maximum number of files to download (default: 2).

        Returns
        -------
        list[dict]
            List of dicts with keys ``local_path``, ``name``, ``modified_time``,
            sorted most-recent-first.
        """
        dest_dir = Path(dest_dir)

        # Collect PPTX from top-level client folder only (no Archive)
        pptx_files: list[dict] = self._list_files_by_ext(
            client_folder_id, {".pptx"},
        )

        if not pptx_files:
            logger.warning(
                "No PPTX files found in client folder %s", client_folder_id,
            )
            return []

        # Sort by modifiedTime descending (most recent first)
        pptx_files.sort(
            key=lambda f: f.get("modifiedTime", ""), reverse=True,
        )

        # Take the most recent `count` files — these are previous reports
        # (the report being generated now has NOT been uploaded to Drive yet).
        previous = pptx_files[:count]

        results = []
        for i, entry in enumerate(previous):
            dest_path = dest_dir / f"{client_slug}_prev_{i}.pptx"
            try:
                local = self._download_file(entry["id"], dest_path)
                results.append({
                    "local_path": local,
                    "name": entry["name"],
                    "modified_time": entry.get("modifiedTime"),
                })
            except Exception as exc:
                logger.warning(
                    "Failed to download previous PPTX '%s': %s",
                    entry["name"], exc,
                )

        logger.info(
            "Downloaded %d previous PPTX file(s) for '%s' (from %d candidates)",
            len(results), client_slug, len(pptx_files),
        )
        return results

    def fetch_reference_pptx(
        self, client_slug: str, dest_dir: Path | str = "assets/tmp",
    ) -> Path | None:
        """Download the most recent PPTX from the client's Drive folder.

        Convenience wrapper around :meth:`find_client_folder` +
        :meth:`find_client_pptx`.

        Returns
        -------
        Path | None
            Local path to the downloaded reference PPTX, or ``None`` if the
            client folder or PPTX could not be found.
        """
        try:
            folder = self.find_client_folder(client_slug)
        except ClientFolderNotFound:
            logger.warning(
                "fetch_reference_pptx: no Drive folder for '%s'", client_slug,
            )
            return None
        return self.find_client_pptx(folder["id"], client_slug, dest_dir)

    def find_latest_ndr_folder(self, client_folder_id: str) -> dict:
        """Find the most recent NDR-* subfolder by date.

        Expects folder names like ``NDR-ClientName-2026-01-22``,
        ``NDR-ClientName-Jan2026``, etc.

        Raises ``NDRFolderNotFound`` if none exist.
        """
        folders = self._list_folders(client_folder_id)
        ndr_folders = [
            f for f in folders
            if f["name"].upper().startswith("NDR")
        ]
        if not ndr_folders:
            raise NDRFolderNotFound(
                f"No NDR folders found in folder {client_folder_id}. "
                f"Available subfolders: {[f['name'] for f in folders]}"
            )

        def _parse_date(name: str) -> datetime:
            """Try to extract a date from an NDR folder name."""
            patterns = [
                (r"(\d{4}-\d{2}-\d{2})", "%Y-%m-%d"),
                (r"(\d{2}-\d{2}-\d{4})", "%m-%d-%Y"),
                (r"(\d{1,2}[A-Za-z]{3}\d{4})", "%d%b%Y"),
                (r"([A-Za-z]{3}\d{4})", "%b%Y"),
                (r"(\d{8})", "%m%d%Y"),
            ]
            for regex, fmt in patterns:
                match = re.search(regex, name)
                if match:
                    try:
                        return datetime.strptime(match.group(1), fmt)
                    except ValueError:
                        continue
            return datetime.min

        ndr_folders.sort(key=lambda f: _parse_date(f["name"]), reverse=True)
        latest = ndr_folders[0]
        logger.info(
            "Selected NDR folder: '%s' (%s) from %d candidates",
            latest["name"], latest["id"], len(ndr_folders),
        )
        return latest

    def find_latest_vn_folder(self, client_folder_id: str) -> dict:
        """Find the most recent VN-* subfolder by date.

        Expects folder names like ``VN-ClientName-2026-01-22``,
        ``VN-ClientName-Jan2026``, etc.

        Raises ``VNFolderNotFound`` if none exist.
        """
        folders = self._list_folders(client_folder_id)
        vn_folders = [
            f for f in folders
            if f["name"].upper().startswith("VN")
        ]
        if not vn_folders:
            raise VNFolderNotFound(
                f"No VN folders found in folder {client_folder_id}. "
                f"Available subfolders: {[f['name'] for f in folders]}"
            )

        def _parse_date(name: str) -> datetime:
            """Try to extract a date from a VN folder name."""
            patterns = [
                (r"(\d{4}-\d{2}-\d{2})", "%Y-%m-%d"),
                (r"(\d{2}-\d{2}-\d{4})", "%m-%d-%Y"),
                (r"(\d{1,2}[A-Za-z]{3}\d{4})", "%d%b%Y"),
                (r"([A-Za-z]{3}\d{4})", "%b%Y"),
                (r"(\d{8})", "%m%d%Y"),
            ]
            for regex, fmt in patterns:
                match = re.search(regex, name)
                if match:
                    try:
                        return datetime.strptime(match.group(1), fmt)
                    except ValueError:
                        continue
            return datetime.min

        vn_folders.sort(key=lambda f: _parse_date(f["name"]), reverse=True)
        latest = vn_folders[0]
        logger.info(
            "Selected VN folder: '%s' (%s) from %d candidates",
            latest["name"], latest["id"], len(vn_folders),
        )
        return latest

    def find_ndr_folders_by_sensor(
        self,
        client_folder_id: str,
        sensor_ids: list[str],
    ) -> dict[str, dict]:
        """Find the latest NDR folder for each sensor.

        For each sensor ID, looks for folders matching
        ``NDR-*-<sensor_id>-*`` and picks the most recent by date.

        Parameters
        ----------
        client_folder_id : str
            Google Drive folder ID for the client.
        sensor_ids : list[str]
            Sensor identifiers to search for (e.g. ``["GAPRD", "VAHQ"]``).

        Returns
        -------
        dict[str, dict]
            Mapping of sensor_id -> Drive folder dict (``{id, name}``).
            Sensors with no matching folder are omitted.
        """
        folders = self._list_folders(client_folder_id)
        ndr_folders = [
            f for f in folders
            if f["name"].upper().startswith("NDR")
        ]

        def _parse_date(name: str) -> datetime:
            patterns = [
                (r"(\d{4}-\d{2}-\d{2})", "%Y-%m-%d"),
                (r"(\d{2}-\d{2}-\d{4})", "%m-%d-%Y"),
                (r"(\d{1,2}[A-Za-z]{3}\d{4})", "%d%b%Y"),
                (r"([A-Za-z]{3}\d{4})", "%b%Y"),
                (r"(\d{8})", "%m%d%Y"),
            ]
            for regex, fmt in patterns:
                match = re.search(regex, name)
                if match:
                    try:
                        return datetime.strptime(match.group(1), fmt)
                    except ValueError:
                        continue
            return datetime.min

        result: dict[str, dict] = {}

        for sensor_id in sensor_ids:
            # Match folders containing the sensor ID (case-insensitive)
            matching = [
                f for f in ndr_folders
                if sensor_id.lower() in f["name"].lower()
            ]
            if not matching:
                logger.warning(
                    "No NDR folder found for sensor '%s'", sensor_id,
                )
                continue

            # Pick the most recent by date
            matching.sort(key=lambda f: _parse_date(f["name"]), reverse=True)
            result[sensor_id] = matching[0]
            logger.info(
                "Sensor '%s' -> folder '%s' (%s)",
                sensor_id, matching[0]["name"], matching[0]["id"],
            )

        return result

    # ------------------------------------------------------------------ #
    # Sensor detection & grouping
    # ------------------------------------------------------------------ #
    def detect_sensors(self, client_slug: str) -> list[str]:
        """Scan a client's NDR folders to detect all available sensor IDs.

        Looks for folders matching patterns like:
            NDR-<ClientName>-<SensorID>-<Date>
            NDR-<ClientName>-<SensorID>-<Date>-<Extra>

        Sensor IDs are uppercase alphanumeric identifiers (e.g. "VAPRD",
        "VAHQ", "DC01").

        Returns
        -------
        list[str]
            Sorted list of unique sensor IDs found, or ``["DEFAULT"]`` if
            no sensor-specific folders are detected.
        """
        try:
            client_folder = self.find_client_folder(client_slug)
        except ClientFolderNotFound:
            logger.warning(
                "detect_sensors: client folder not found for '%s', defaulting to single sensor",
                client_slug,
            )
            return ["DEFAULT"]

        folders = self._list_folders(client_folder["id"])
        ndr_folders = [
            f for f in folders
            if f["name"].upper().startswith("NDR")
        ]

        if not ndr_folders:
            logger.info(
                "detect_sensors: no NDR folders found for '%s', defaulting to single sensor",
                client_slug,
            )
            return ["DEFAULT"]

        # Build search terms for the client slug (resolve display names to slugs first)
        raw = client_slug.strip()
        normalized = raw.lower()
        slug = DISPLAY_NAME_TO_SLUG.get(normalized)
        if slug is None:
            for display_name, s in DISPLAY_NAME_TO_SLUG.items():
                if display_name in normalized or normalized in display_name:
                    slug = s
                    break
        if slug is not None:
            search_terms = CLIENT_FOLDER_MAP.get(slug, [slug])
        else:
            search_terms = CLIENT_FOLDER_MAP.get(normalized, [normalized, raw])

        # Extract sensor IDs from folder names.
        # Pattern: NDR-<client_term>-<SENSOR_ID>-<date_part>
        # The sensor ID is the segment after the client name and before date-like segments.
        sensors: set[str] = set()

        for folder in ndr_folders:
            name = folder["name"]
            # Split by hyphens, stripping outer whitespace from each part
            parts = [p.strip() for p in name.split("-") if p.strip()]

            if len(parts) < 3:
                # Too few parts to contain a sensor ID (e.g. "NDR-ClientName-Date")
                continue

            # parts[0] should be "NDR" (already filtered).
            # Find where the client name ends by matching search terms.
            # Client name may span multiple hyphen-separated parts (e.g. "Monterey-Mech").
            client_end_idx = 1  # default: client is parts[1]

            # Try matching progressively longer sequences starting at index 1
            full_lower = "-".join(parts[1:]).lower()
            for term in search_terms:
                term_lower = term.lower().replace(" ", "")
                # Check if the concatenation of parts starting at 1 begins with the term
                concat = ""
                for i in range(1, len(parts)):
                    concat += parts[i].lower()
                    if concat == term_lower or term_lower.startswith(concat) or concat.startswith(term_lower):
                        client_end_idx = i
                        if len(concat) >= len(term_lower):
                            break

            # The segment immediately after the client name could be the sensor ID
            sensor_candidate_idx = client_end_idx + 1

            if sensor_candidate_idx >= len(parts):
                continue

            candidate = parts[sensor_candidate_idx]

            # Skip if the candidate looks like a date
            if re.match(r"^\d{4}$", candidate):          # year like "2026"
                continue
            if re.match(r"^\d{1,2}$", candidate):        # day or month number
                continue
            if re.match(r"^\d{2,4}-?\d{2}-?\d{2,4}$", candidate):  # full date
                continue
            if re.match(r"^[A-Za-z]{3}\d{2,4}$", candidate):  # "Jan2026"
                continue
            if re.match(r"^\d{1,2}\s*[A-Za-z]+\s*\d{2,4}$", candidate):  # "22 Jan 2026"
                continue

            # Valid sensor IDs are alphanumeric, at least 2 chars, not purely numeric
            if re.match(r"^[A-Za-z0-9]{2,}$", candidate) and not candidate.isdigit():
                sensors.add(candidate.upper())

        # Fallback: group NDR folders by non-date base signature and derive
        # sensor IDs from the distinguishing segments when the main loop missed them.
        def _base_sig(name: str) -> str:
            parts = [p.strip() for p in name.split("-") if p.strip()]
            return "-".join(
                p.lower() for p in parts
                if not re.match(r"^\d{1,4}$", p)
                and not re.match(r"^[A-Za-z]{3}\d{2,4}$", p)
            )

        base_groups: dict[str, list] = {}
        for folder in ndr_folders:
            sig = _base_sig(folder["name"])
            base_groups.setdefault(sig, []).append(folder)

        if len(base_groups) >= 2 and len(sensors) < len(base_groups):
            # Find common segments shared by ALL groups
            all_sigs = list(base_groups.keys())
            common_parts = set(all_sigs[0].split("-"))
            for sig in all_sigs[1:]:
                common_parts &= set(sig.split("-"))

            for sig in all_sigs:
                unique = [
                    p for p in sig.split("-")
                    if p not in common_parts
                    and re.match(r"^[A-Za-z0-9]{2,}$", p)
                    and not p.isdigit()
                ]
                if unique:
                    candidate = unique[0].upper()
                    if candidate not in sensors:
                        sensors.add(candidate)

            if len(sensors) >= 2:
                logger.info(
                    "detect_sensors fallback: derived %d sensor(s) from folder grouping: %s",
                    len(sensors), sorted(sensors),
                )

        if not sensors:
            logger.info(
                "detect_sensors: no sensor IDs found in %d NDR folders for '%s', "
                "defaulting to single sensor",
                len(ndr_folders), client_slug,
            )
            return ["DEFAULT"]

        result = sorted(sensors)
        logger.info(
            "detect_sensors: found %d sensor(s) for '%s': %s",
            len(result), client_slug, result,
        )
        return result

    @staticmethod
    def _extract_sensor_from_filename(filename: str, client_slug: str) -> Optional[str]:
        """Extract sensor identifier from filename.
        
        Examples:
            elephant-vaprd_beaconing-score.csv -> 'VAPRD'
            elephant-vahq_domain.csv -> 'VAHQ'
            heatmap_2026-01-19.png -> None (no sensor)
            monterey_beaconing-score.csv -> None (no sensor pattern)
        
        Returns None if no sensor identifier found.
        """
        stem = Path(filename).stem.lower()
        # Pattern: <client>-<sensor>_<datatype>
        pattern = rf"{re.escape(client_slug.lower())}-(\w+)_"
        match = re.search(pattern, stem)
        return match.group(1).upper() if match else None

    @staticmethod
    def _group_files_by_sensor(
        files: list[dict],
        client_slug: str
    ) -> dict[str, list[dict]]:
        """Group files by sensor identifier.
        
        Returns a dictionary mapping sensor_id -> list of file dicts.
        Files without sensor identifiers are grouped under "DEFAULT".
        """
        grouped: dict[str, list[dict]] = {}
        
        for f in files:
            sensor_id = DriveAgent._extract_sensor_from_filename(f["name"], client_slug)
            if sensor_id is None:
                sensor_id = "DEFAULT"
            
            if sensor_id not in grouped:
                grouped[sensor_id] = []
            grouped[sensor_id].append(f)
        
        return grouped

    # ------------------------------------------------------------------ #
    # File-to-slide mapping
    # ------------------------------------------------------------------ #
    @staticmethod
    def _map_file_to_slide(filename: str) -> Optional[tuple[int, str]]:
        """Map an image filename to ``(slide_index, shape_name)``."""
        stem = Path(filename).stem.lower().replace("-", "_").replace(" ", "_")
        for pattern, mapping in NDR_FILE_MAP.items():
            if pattern in stem:
                return mapping
        return None

    @staticmethod
    def _normalise_title(text: str) -> str:
        """Normalise a slide title / CSV stem for fuzzy matching.

        Lower-case, replace spaces and common punctuation with ``_``,
        collapse multiple underscores.
        """
        out = text.lower().strip()
        out = re.sub(r"[^a-z0-9]+", "_", out)
        return out.strip("_")

    @staticmethod
    def _map_csv_to_slide(filename: str) -> Optional[tuple[int, str]]:
        """Map a CSV filename to ``(slide_index, shape_name)`` using
        ``TABLE_SLIDE_MAP``.

        Matching is by normalised stem: e.g. ``Top_Software_Patches.csv``
        or ``top-software-patches.csv`` → key ``top_software_patches``.
        Fuzzy fallback: stem containing 'software' + 'patch' → software table;
        'microsoft' + 'patch' → Microsoft patches table.
        """
        stem = DriveAgent._normalise_title(Path(filename).stem)
        # Exact match first
        if stem in TABLE_SLIDE_MAP:
            return TABLE_SLIDE_MAP[stem]
        # Substring / contains match
        for pattern, mapping in TABLE_SLIDE_MAP.items():
            if pattern in stem or stem in pattern:
                return mapping
        # Fuzzy: "software" + "patch" -> Top Software Patches (slide 14)
        if "software" in stem and "patch" in stem:
            return (14, "software_patches_table")
        if "microsoft" in stem and "patch" in stem:
            return (15, "ms_patches_table")
        return None

    # ------------------------------------------------------------------ #
    # Nuclei file fetching
    # ------------------------------------------------------------------ #
    def fetch_nuclei_file(
        self,
        client_slug: str,
        dest_dir: str | Path = "assets/ndr_cache",
    ) -> Path | None:
        """Download the nuclei.txt file for a single-sensor client.

        Looks in the latest NDR folder for a ``.txt`` file whose name
        contains ``nuclei``.

        Returns the local path or ``None`` if not found.
        """
        dest = Path(dest_dir) / client_slug
        dest.mkdir(parents=True, exist_ok=True)

        client_folder = self.find_client_folder(client_slug)
        ndr_folder = self.find_latest_ndr_folder(client_folder["id"])

        txt_files = self._list_files_by_ext(ndr_folder["id"], TXT_EXTENSIONS)
        nuclei_files = [
            f for f in txt_files
            if "nuclei" in f["name"].lower()
        ]

        if not nuclei_files:
            logger.info(
                "No nuclei.txt found in NDR folder '%s'", ndr_folder["name"],
            )
            return None

        # Pick the first match (typically there's only one)
        target = nuclei_files[0]
        local_path = dest / target["name"]
        self._download_file(target["id"], local_path)
        logger.info(
            "Downloaded nuclei file '%s' -> %s", target["name"], local_path,
        )
        return local_path

    def fetch_nuclei_files_multi_sensor(
        self,
        client_slug: str,
        sensor_ids: list[str],
        dest_dir: str | Path = "assets/ndr_cache",
    ) -> dict[str, Path]:
        """Download nuclei.txt files from each sensor's NDR folder.

        Parameters
        ----------
        client_slug : str
            Client identifier.
        sensor_ids : list[str]
            Sensor identifiers to fetch.
        dest_dir : str | Path
            Local cache directory.

        Returns
        -------
        dict[str, Path]
            Mapping of sensor_id -> local file path.
            Only sensors with a nuclei file are included.
        """
        dest = Path(dest_dir) / client_slug
        dest.mkdir(parents=True, exist_ok=True)

        client_folder = self.find_client_folder(client_slug)
        sensor_folders = self.find_ndr_folders_by_sensor(
            client_folder["id"], sensor_ids,
        )

        result: dict[str, Path] = {}

        for sensor_id in sensor_ids:
            folder = sensor_folders.get(sensor_id)
            if not folder:
                logger.warning(
                    "Sensor '%s': no NDR folder found, skipping nuclei",
                    sensor_id,
                )
                continue

            txt_files = self._list_files_by_ext(folder["id"], TXT_EXTENSIONS)
            nuclei_files = [
                f for f in txt_files
                if "nuclei" in f["name"].lower()
            ]

            if not nuclei_files:
                logger.info(
                    "Sensor '%s': no nuclei.txt in '%s'",
                    sensor_id, folder["name"],
                )
                continue

            target = nuclei_files[0]
            sensor_dest = dest / sensor_id.lower()
            sensor_dest.mkdir(parents=True, exist_ok=True)
            local_path = sensor_dest / target["name"]
            self._download_file(target["id"], local_path)

            result[sensor_id] = local_path
            logger.info(
                "Sensor '%s': downloaded nuclei file '%s' -> %s",
                sensor_id, target["name"], local_path,
            )

        logger.info(
            "Nuclei files fetched for %d/%d sensors for '%s'",
            len(result), len(sensor_ids), client_slug,
        )
        return result

    # ------------------------------------------------------------------ #
    # Remediation plan (VN folder)
    # ------------------------------------------------------------------ #
    def fetch_remediation_csv(
        self,
        client_slug: str,
        dest_dir: str | Path = "assets/ndr_cache",
    ) -> Path | None:
        """Download the remediation plan CSV from the latest VN-* folder.

        Looks for a CSV whose filename contains both "remediation" and
        "plan" (case-insensitive).

        Returns the local path or ``None`` if no matching CSV is found.
        """
        dest = Path(dest_dir) / client_slug
        dest.mkdir(parents=True, exist_ok=True)

        client_folder = self.find_client_folder(client_slug)
        vn_folder = self.find_latest_vn_folder(client_folder["id"])

        csv_files = self._list_files_by_ext(vn_folder["id"], CSV_EXTENSIONS)
        remediation_files = [
            f for f in csv_files
            if "remediation" in f["name"].lower()
            and "plan" in f["name"].lower()
        ]

        if not remediation_files:
            logger.info(
                "No remediation plan CSV found in VN folder '%s'",
                vn_folder["name"],
            )
            return None

        target = remediation_files[0]
        local_path = dest / target["name"]
        self._download_file(target["id"], local_path)
        logger.info(
            "Downloaded remediation CSV '%s' -> %s",
            target["name"], local_path,
        )
        return local_path

    def fetch_software_csv(
        self,
        client_slug: str,
        dest_dir: str | Path = "assets/ndr_cache",
    ) -> Path | None:
        """Download the non-Windows software updates CSV from the latest VN-* folder.

        Looks for a CSV whose filename starts with
        ``vulnaggregateother_`` or ``non_windows_updates_``
        (case-insensitive).

        Returns the local path or ``None`` if no matching CSV is found.
        """
        dest = Path(dest_dir) / client_slug
        dest.mkdir(parents=True, exist_ok=True)

        client_folder = self.find_client_folder(client_slug)
        vn_folder = self.find_latest_vn_folder(client_folder["id"])

        csv_files = self._list_files_by_ext(vn_folder["id"], CSV_EXTENSIONS)
        software_files = [
            f for f in csv_files
            if f["name"].lower().startswith("vulnaggregateother_")
            or f["name"].lower().startswith("non_windows_updates_")
        ]

        if not software_files:
            logger.info(
                "No software updates CSV found in VN folder '%s'",
                vn_folder["name"],
            )
            return None

        target = software_files[0]
        local_path = dest / target["name"]
        self._download_file(target["id"], local_path)
        logger.info(
            "Downloaded software CSV '%s' -> %s",
            target["name"], local_path,
        )
        return local_path

    def fetch_inventory_csv(
        self,
        client_slug: str,
        dest_dir: str | Path = "assets/ndr_cache",
    ) -> Path | None:
        """Download the inventory report CSV from the latest VN-* folder.

        Looks for a CSV whose filename contains "inventoryreport"
        (case-insensitive).

        Returns the local path or ``None`` if no matching CSV is found.
        """
        dest = Path(dest_dir) / client_slug
        dest.mkdir(parents=True, exist_ok=True)

        client_folder = self.find_client_folder(client_slug)
        vn_folder = self.find_latest_vn_folder(client_folder["id"])

        csv_files = self._list_files_by_ext(vn_folder["id"], CSV_EXTENSIONS)
        inventory_files = [
            f for f in csv_files
            if "inventoryreport" in f["name"].lower().replace(" ", "").replace("_", "")
        ]

        if not inventory_files:
            logger.info(
                "No inventory report CSV found in VN folder '%s'",
                vn_folder["name"],
            )
            return None

        target = inventory_files[0]
        local_path = dest / target["name"]
        self._download_file(target["id"], local_path)
        logger.info(
            "Downloaded inventory CSV '%s' -> %s",
            target["name"], local_path,
        )
        return local_path

    # ------------------------------------------------------------------ #
    # Full pipeline
    # ------------------------------------------------------------------ #
    def fetch_csv_assets(
        self,
        client_slug: str,
        dest_dir: str | Path = "assets/ndr_cache",
    ) -> list[CSVTableAsset]:
        """Find client folder → latest NDR folder → download CSV files.

        Returns a list of ``CSVTableAsset`` objects, each carrying the
        local file path and the target table slide/shape assignment.
        """
        dest = Path(dest_dir) / client_slug
        dest.mkdir(parents=True, exist_ok=True)

        # 1. Client folder
        client_folder = self.find_client_folder(client_slug)

        # 2. Latest NDR subfolder
        ndr_folder = self.find_latest_ndr_folder(client_folder["id"])

        # 3. CSV files (including optional tables/ subfolder)
        files = self._list_csv_files(ndr_folder["id"])
        if not files:
            logger.info(
                "No CSV files found in NDR folder '%s'", ndr_folder["name"]
            )
            return []

        logger.info(
            "Found %d CSV files in '%s'", len(files), ndr_folder["name"]
        )

        # 4. Download & map
        assets: list[CSVTableAsset] = []
        used_slides: set[int] = set()

        for f in files:
            mapping = self._map_csv_to_slide(f["name"])
            if not mapping or mapping[0] in used_slides:
                logger.info(
                    "CSV '%s' does not match any table slide, skipping",
                    f["name"],
                )
                continue

            local_path = dest / f["name"]
            self._download_file(f["id"], local_path)

            slide_idx, shape_name = mapping
            used_slides.add(slide_idx)
            assets.append(CSVTableAsset(
                local_path=local_path,
                slide_index=slide_idx,
                shape_name=shape_name,
                original_filename=f["name"],
            ))
            logger.info(
                "Mapped CSV '%s' -> slide %d (%s)",
                f["name"], slide_idx, shape_name,
            )

        assets.sort(key=lambda a: a.slide_index)
        logger.info(
            "CSV table assets ready: %d files mapped for '%s'",
            len(assets), client_slug,
        )
        return assets

    def fetch_ndr_assets(
        self,
        client_slug: str,
        dest_dir: str | Path = "assets/ndr_cache",
    ) -> list[NDRAsset]:
        """Find client folder -> latest NDR folder -> download images.

        Returns a list of ``NDRAsset`` objects, each carrying the local
        file path and the target slide/shape assignment.
        """
        dest = Path(dest_dir) / client_slug
        dest.mkdir(parents=True, exist_ok=True)

        # 1. Client folder
        client_folder = self.find_client_folder(client_slug)

        # 2. Latest NDR subfolder
        ndr_folder = self.find_latest_ndr_folder(client_folder["id"])

        # 3. Image files
        files = self._list_image_files(ndr_folder["id"])
        if not files:
            logger.warning(
                "No image files found in NDR folder '%s'", ndr_folder["name"]
            )
            return []

        logger.info(
            "Found %d image files in '%s'", len(files), ndr_folder["name"]
        )

        # 4. Download all files first
        for f in files:
            local_path = dest / f["name"]
            self._download_file(f["id"], local_path)

        # 5. Split into outbound (heatmap / outbound) vs other; assign outbound sequentially to slides 5, 6
        assets: list[NDRAsset] = []
        used_slides: set[int] = set()
        outbound_files: list[tuple[Path, str]] = []
        other_files: list[tuple[Path, str]] = []

        for f in files:
            local_path = dest / f["name"]
            stem = Path(f["name"]).stem.lower().replace("-", "_").replace(" ", "_")
            if any(kw in stem for kw in OUTBOUND_STEM_KEYWORDS):
                outbound_files.append((local_path, f["name"]))
            else:
                other_files.append((local_path, f["name"]))

        # If we have 1 or 2 images and none matched outbound, treat all as outbound so they get placed
        if len(files) <= 2 and not outbound_files and other_files:
            outbound_files = list(other_files)
            other_files = []

        outbound_files.sort(key=lambda x: x[1])
        for (local_path, orig_name), (slide_idx, shape_name) in zip(
            outbound_files, OUTBOUND_SLIDES
        ):
            used_slides.add(slide_idx)
            assets.append(NDRAsset(
                local_path=local_path,
                slide_index=slide_idx,
                shape_name=shape_name,
                original_filename=orig_name,
            ))
            logger.info(
                "Outbound (sequential): '%s' -> slide %d (%s)",
                orig_name, slide_idx, shape_name,
            )

        # 6. Map remaining images by keyword or positional fallback
        unmapped: list[tuple[Path, str]] = []
        for local_path, orig_name in other_files:
            mapping = self._map_file_to_slide(orig_name)
            if mapping and mapping[0] not in used_slides:
                slide_idx, shape_name = mapping
                used_slides.add(slide_idx)
                assets.append(NDRAsset(
                    local_path=local_path,
                    slide_index=slide_idx,
                    shape_name=shape_name,
                    original_filename=orig_name,
                ))
                logger.info(
                    "Mapped '%s' -> slide %d (%s)",
                    orig_name, slide_idx, shape_name,
                )
            else:
                unmapped.append((local_path, orig_name))

        available_slots = [
            (idx, name) for idx, name in NDR_SLIDES_ORDERED
            if idx not in used_slides
        ]
        for (local_path, orig_name), (slide_idx, shape_name) in zip(
            unmapped, available_slots
        ):
            used_slides.add(slide_idx)
            assets.append(NDRAsset(
                local_path=local_path,
                slide_index=slide_idx,
                shape_name=shape_name,
                original_filename=orig_name,
            ))
            logger.info(
                "Positional fallback: '%s' -> slide %d (%s)",
                orig_name, slide_idx, shape_name,
            )

        assets.sort(key=lambda a: a.slide_index)
        logger.info(
            "NDR assets ready: %d images mapped for '%s'",
            len(assets), client_slug,
        )
        return assets

    def fetch_ndr_assets_multi_sensor(
        self,
        client_slug: str,
        sensor_ids: list[str],
        dest_dir: str | Path = "assets/ndr_cache",
    ) -> dict[str, list[NDRAsset]]:
        """Fetch NDR images from each sensor's own folder.

        Iterates over per-sensor NDR folders (e.g.
        ``NDR-elephant-vaprd-2026-02-03``) and downloads all images from
        each.  All files from a sensor folder are assigned that sensor ID
        regardless of filename.

        Parameters
        ----------
        client_slug : str
            Client identifier.
        sensor_ids : list[str]
            Sensor identifiers to fetch (e.g. ``["GAPRD", "VAHQ", "VAPRD"]``).
        dest_dir : str | Path
            Local cache directory.

        Returns
        -------
        dict[str, list[NDRAsset]]
            Mapping of sensor_id -> list of NDRAsset objects.
        """
        dest = Path(dest_dir) / client_slug
        dest.mkdir(parents=True, exist_ok=True)

        client_folder = self.find_client_folder(client_slug)
        sensor_folders = self.find_ndr_folders_by_sensor(
            client_folder["id"], sensor_ids,
        )

        sensor_assets: dict[str, list[NDRAsset]] = {}

        for sensor_id in sensor_ids:
            folder = sensor_folders.get(sensor_id)
            if not folder:
                logger.warning(
                    "Sensor '%s': no NDR folder found, skipping images",
                    sensor_id,
                )
                sensor_assets[sensor_id] = []
                continue

            files = self._list_image_files(folder["id"])
            if not files:
                logger.info(
                    "Sensor '%s': no image files in '%s'",
                    sensor_id, folder["name"],
                )
                sensor_assets[sensor_id] = []
                continue

            logger.info(
                "Sensor '%s': found %d images in '%s'",
                sensor_id, len(files), folder["name"],
            )

            # Download all images into sensor-specific subdirectory
            sensor_dest = dest / sensor_id.lower()
            sensor_dest.mkdir(parents=True, exist_ok=True)
            for f in files:
                local_path = sensor_dest / f["name"]
                self._download_file(f["id"], local_path)

            # Classify images: outbound/heatmap vs others
            assets: list[NDRAsset] = []
            used_slides: set[int] = set()
            outbound_files: list[tuple[Path, str]] = []
            other_files: list[tuple[Path, str]] = []

            for f in files:
                local_path = sensor_dest / f["name"]
                stem = Path(f["name"]).stem.lower().replace("-", "_").replace(" ", "_")
                if any(kw in stem for kw in OUTBOUND_STEM_KEYWORDS):
                    outbound_files.append((local_path, f["name"]))
                else:
                    other_files.append((local_path, f["name"]))

            if len(files) <= 2 and not outbound_files and other_files:
                outbound_files = list(other_files)
                other_files = []

            outbound_files.sort(key=lambda x: x[1])
            for (local_path, orig_name), (slide_idx, shape_name) in zip(
                outbound_files, OUTBOUND_SLIDES
            ):
                used_slides.add(slide_idx)
                assets.append(NDRAsset(
                    local_path=local_path,
                    slide_index=slide_idx,
                    shape_name=shape_name,
                    original_filename=orig_name,
                    sensor_id=sensor_id,
                ))

            unmapped: list[tuple[Path, str]] = []
            for local_path, orig_name in other_files:
                mapping = self._map_file_to_slide(orig_name)
                if mapping and mapping[0] not in used_slides:
                    slide_idx, shape_name = mapping
                    used_slides.add(slide_idx)
                    assets.append(NDRAsset(
                        local_path=local_path,
                        slide_index=slide_idx,
                        shape_name=shape_name,
                        original_filename=orig_name,
                        sensor_id=sensor_id,
                    ))
                else:
                    unmapped.append((local_path, orig_name))

            available_slots = [
                (idx, name) for idx, name in NDR_SLIDES_ORDERED
                if idx not in used_slides
            ]
            for (local_path, orig_name), (slide_idx, shape_name) in zip(
                unmapped, available_slots
            ):
                used_slides.add(slide_idx)
                assets.append(NDRAsset(
                    local_path=local_path,
                    slide_index=slide_idx,
                    shape_name=shape_name,
                    original_filename=orig_name,
                    sensor_id=sensor_id,
                ))

            assets.sort(key=lambda a: a.slide_index)
            sensor_assets[sensor_id] = assets
            logger.info(
                "Sensor '%s': %d images mapped from folder '%s'",
                sensor_id, len(assets), folder["name"],
            )

        return sensor_assets

    def fetch_csv_assets_multi_sensor(
        self,
        client_slug: str,
        sensor_ids: list[str],
        dest_dir: str | Path = "assets/ndr_cache",
    ) -> dict[str, list[CSVTableAsset]]:
        """Fetch CSV files from each sensor's own folder.

        Iterates over per-sensor NDR folders and downloads all CSVs from
        each.  Non-sensor CSVs (patches) found in any folder are grouped
        under ``"DEFAULT"``.

        Parameters
        ----------
        client_slug : str
            Client identifier.
        sensor_ids : list[str]
            Sensor identifiers to fetch.
        dest_dir : str | Path
            Local cache directory.

        Returns
        -------
        dict[str, list[CSVTableAsset]]
            Mapping of sensor_id -> list of CSVTableAsset objects.
            Non-sensor CSVs (patches) are under ``"DEFAULT"`` key.
        """
        dest = Path(dest_dir) / client_slug
        dest.mkdir(parents=True, exist_ok=True)

        client_folder = self.find_client_folder(client_slug)
        sensor_folders = self.find_ndr_folders_by_sensor(
            client_folder["id"], sensor_ids,
        )

        sensor_assets: dict[str, list[CSVTableAsset]] = {}

        for sensor_id in sensor_ids:
            folder = sensor_folders.get(sensor_id)
            if not folder:
                logger.warning(
                    "Sensor '%s': no NDR folder found, skipping CSVs",
                    sensor_id,
                )
                sensor_assets[sensor_id] = []
                continue

            files = self._list_csv_files(folder["id"])
            if not files:
                logger.info(
                    "Sensor '%s': no CSV files in '%s'",
                    sensor_id, folder["name"],
                )
                sensor_assets[sensor_id] = []
                continue

            logger.info(
                "Sensor '%s': found %d CSVs in '%s'",
                sensor_id, len(files), folder["name"],
            )

            assets: list[CSVTableAsset] = []
            used_slides: set[int] = set()

            for f in files:
                mapping = self._map_csv_to_slide(f["name"])
                if not mapping or mapping[0] in used_slides:
                    logger.info(
                        "Sensor '%s' - CSV '%s' no slide match, skipping",
                        sensor_id, f["name"],
                    )
                    continue

                sensor_dest = dest / sensor_id.lower()
                sensor_dest.mkdir(parents=True, exist_ok=True)
                local_path = sensor_dest / f["name"]
                self._download_file(f["id"], local_path)

                slide_idx, shape_name = mapping
                used_slides.add(slide_idx)

                # Patches CSVs (slides 14, 15) are not sensor-specific —
                # group them under "DEFAULT" so api.py handles them separately.
                if slide_idx in (14, 15):
                    default_assets = sensor_assets.setdefault("DEFAULT", [])
                    default_assets.append(CSVTableAsset(
                        local_path=local_path,
                        slide_index=slide_idx,
                        shape_name=shape_name,
                        original_filename=f["name"],
                        sensor_id="DEFAULT",
                    ))
                    logger.info(
                        "Sensor '%s' - Patches CSV '%s' -> DEFAULT slide %d",
                        sensor_id, f["name"], slide_idx,
                    )
                else:
                    assets.append(CSVTableAsset(
                        local_path=local_path,
                        slide_index=slide_idx,
                        shape_name=shape_name,
                        original_filename=f["name"],
                        sensor_id=sensor_id,
                    ))
                    logger.info(
                        "Sensor '%s' - Mapped CSV '%s' -> slide %d (%s)",
                        sensor_id, f["name"], slide_idx, shape_name,
                    )

            assets.sort(key=lambda a: a.slide_index)
            sensor_assets[sensor_id] = assets
            logger.info(
                "Sensor '%s': %d CSV files mapped",
                sensor_id, len(assets),
            )

        return sensor_assets


# ---------------------------------------------------------------------------
# Post-processing: place downloaded images onto the PPTX
# ---------------------------------------------------------------------------
def place_ndr_images(pptx_path: Path, assets: list[NDRAsset]) -> int:
    """Open a generated PPTX and insert NDR chart images onto the
    correct slides, replacing the dashed-border placeholder shapes.

    Uses the same contain-fit algorithm as the main report generator.
    When a matching placeholder shape is found, its actual position and
    size are used as the bounding box; otherwise falls back to
    ``NDR_BOX_INCHES``.

    Returns the number of images successfully placed.
    """
    if not assets:
        return 0

    from pptx.util import Inches

    prs = PptxPresentation(str(pptx_path))
    default_box = Rect(
        x=int(Inches(NDR_BOX_INCHES["left"])),
        y=int(Inches(NDR_BOX_INCHES["top"])),
        width=int(Inches(NDR_BOX_INCHES["width"])),
        height=int(Inches(NDR_BOX_INCHES["height"])),
    )

    placed = 0

    for asset in assets:
        if asset.slide_index >= len(prs.slides):
            logger.warning(
                "Slide %d out of range, skipping '%s'",
                asset.slide_index, asset.original_filename,
            )
            continue

        slide = prs.slides[asset.slide_index]

        # Find the placeholder shape and read its position before removing it
        box = default_box
        for shape in list(slide.shapes):
            if shape.name == asset.shape_name:
                box = Rect(
                    x=int(shape.left),
                    y=int(shape.top),
                    width=int(shape.width),
                    height=int(shape.height),
                )
                sp = shape._element
                sp.getparent().remove(sp)
                break

        # Read image dimensions
        with PILImage.open(asset.local_path) as img:
            w_px, h_px = img.size
            dpi_x, dpi_y = img.info.get("dpi", (96, 96))
            dpi_x = dpi_x if dpi_x and dpi_x > 0 else 96
            dpi_y = dpi_y if dpi_y and dpi_y > 0 else 96

        w_emu = int(round(w_px / dpi_x * 914400))
        h_emu = int(round(h_px / dpi_y * 914400))
        source = Rect(x=0, y=0, width=w_emu, height=h_emu)

        # Contain-fit transform
        transform = compute_contain_transform(source, box)

        slide.shapes.add_picture(
            str(asset.local_path),
            Emu(transform.dest_x),
            Emu(transform.dest_y),
            Emu(transform.dest_width),
            Emu(transform.dest_height),
        )
        placed += 1
        logger.info(
            "Placed '%s' on slide %d [%dx%d -> %dx%d]",
            asset.original_filename, asset.slide_index,
            w_emu, h_emu, transform.dest_width, transform.dest_height,
        )

    prs.save(str(pptx_path))
    logger.info("Saved PPTX with %d NDR images -> %s", placed, pptx_path)
    return placed


# ---------------------------------------------------------------------------
# Source-column detection
# ---------------------------------------------------------------------------
def _is_source_column(col_name: str) -> bool:
    """Return ``True`` if *col_name* looks like a "source" column to exclude.

    Keeps "Unique sources" (count/metric column). Excludes columns that list
    sources, e.g. "Source", "Sources", "Data source", "References".
    """
    normalised = col_name.strip().lower().replace(" ", "_")
    # Always keep "Unique sources" — it's a metric, not a source list
    if normalised in ("unique_sources", "uniquesources"):
        return False
    if "unique" in normalised and "source" in normalised:
        return False
    for pattern in SOURCE_COLUMN_PATTERNS:
        if pattern in normalised or normalised in pattern:
            return True
    return False


# ---------------------------------------------------------------------------
# KB count helpers
# ---------------------------------------------------------------------------
def count_unique_kb_patches(csv_path: Path) -> int | None:
    """Count unique non-empty values in the ``Required KB`` column.

    Opens the CSV with ``utf-8-sig`` encoding and performs a
    case-insensitive header match for a column named ``required_kb``
    (after normalising spaces/underscores).

    Returns the count, or ``None`` if the column is not found.
    """
    try:
        with open(csv_path, newline="", encoding="utf-8-sig") as fh:
            reader = csv.DictReader(fh)
            if reader.fieldnames is None:
                return None

            # Case-insensitive column match: normalise to "required_kb"
            target_col: str | None = None
            for col in reader.fieldnames:
                normalised = col.strip().lower().replace(" ", "_")
                if normalised == "required_kb":
                    target_col = col
                    break

            if target_col is None:
                logger.info(
                    "Column 'Required KB' not found in %s (columns: %s)",
                    csv_path, reader.fieldnames,
                )
                return None

            unique_kbs: set[str] = set()
            for row in reader:
                value = (row.get(target_col) or "").strip()
                if value:
                    unique_kbs.add(value)

            logger.info(
                "Counted %d unique KB patches from %s", len(unique_kbs), csv_path,
            )
            return len(unique_kbs)
    except Exception as exc:
        logger.error("Failed to read remediation CSV %s: %s", csv_path, exc)
        return None


def build_ms_patches_csv(
    csv_path: Path,
    dest_dir: Path,
) -> tuple[Path | None, int]:
    """Build a Microsoft patches CSV from remediation plan data.
    
    Extracts unique KB patches and their associated Windows product names from the
    remediation plan CSV. The output CSV has two columns:
    - Package Name: The Windows OS version (e.g., "Windows 11 23H2", "Windows Server 2019")
    - KB Patch: The KB number (e.g., KB5068865)
    
    Parameters
    ----------
    csv_path : Path
        Path to the remediation plan CSV
    dest_dir : Path
        Directory to write the cleaned CSV into
    
    Returns
    -------
    (clean_csv_path, unique_count)
        Path to the MS patches CSV and the number of unique KB patches.
        ``(None, 0)`` if the required columns are not found.
    """
    try:
        with open(csv_path, newline="", encoding="utf-8-sig") as fh:
            reader = csv.DictReader(fh)
            if reader.fieldnames is None:
                return None, 0
            
            # Find required columns (case-insensitive)
            kb_col: str | None = None
            package_col: str | None = None
            
            # Log available columns for debugging
            logger.info("Remediation CSV columns: %s", reader.fieldnames)
            
            for col in reader.fieldnames:
                normalised = col.strip().lower().replace(" ", "_")
                
                # KB column
                if normalised in ("required_kb", "kb_patch", "kb", "required_patch"):
                    kb_col = col
                
                # OS/Package name column - prioritize "os" column which has Windows version
                if normalised == "os":
                    package_col = col
                elif package_col is None and normalised in ("package_name", "product", "product_name", 
                                                            "windows_version", "os_version", "operating_system"):
                    package_col = col
            
            if kb_col is None:
                logger.warning(
                    "Required KB column not found in %s (columns: %s)",
                    csv_path, reader.fieldnames,
                )
                return None, 0
            
            if package_col is None:
                # Fallback: look for any column with "package" in name
                for col in reader.fieldnames:
                    if "package" in col.lower() and "name" in col.lower():
                        package_col = col
                        logger.info("Using column '%s' for package names", col)
                        break
            
            if package_col is None:
                logger.warning(
                    "Package name column not found in %s. Using KB column as fallback.",
                    csv_path,
                )
                package_col = kb_col
            
            # Collect KB patches grouped by package name
            # Format: package_name -> [list of KB numbers]
            package_kbs: dict[str, list[str]] = {}
            
            for row in reader:
                kb_value = (row.get(kb_col) or "").strip()
                package_value = (row.get(package_col) or "").strip()
                
                if kb_value and package_value:
                    if package_value not in package_kbs:
                        package_kbs[package_value] = []
                    
                    # Add KB if not already in list for this package
                    if kb_value not in package_kbs[package_value]:
                        package_kbs[package_value].append(kb_value)
            
            if not package_kbs:
                logger.info("No KB patches found in %s", csv_path)
                return None, 0
            
            # Write output CSV
            dest_dir.mkdir(parents=True, exist_ok=True)
            out_csv = dest_dir / "top_ms_patches.csv"
            
            with open(out_csv, "w", newline="", encoding="utf-8") as fout:
                writer = csv.writer(fout)
                writer.writerow(["Package Name", "KB Patch"])
                
                # Sort by package name for logical grouping (Windows versions together)
                for package_name in sorted(package_kbs.keys()):
                    kb_list = package_kbs[package_name]
                    # Sort KBs within each package and join with commas
                    kb_list.sort()
                    kb_string = ", ".join(kb_list)
                    writer.writerow([package_name, kb_string])
            
            total_kb_count = sum(len(kbs) for kbs in package_kbs.values())
            logger.info(
                "Built MS patches CSV: %d packages with %d total KBs from columns '%s' + '%s' -> %s",
                len(package_kbs), total_kb_count, package_col, kb_col, out_csv,
            )
            return out_csv, total_kb_count
            
    except Exception as exc:
        logger.error("Failed to build MS patches CSV from %s: %s", csv_path, exc)
        return None, 0


def _strip_version(name: str) -> str:
    """Remove embedded version numbers and architecture tags from a package name.

    Strips version patterns (``6.0.25``, ``3.1.10``), architecture
    suffixes (``(x64)``, ``(x86)``), and normalises whitespace / dashes
    so that names differing only by version, arch, or punctuation
    collapse together.

    Examples
    --------
    >>> _strip_version("Microsoft ASP.NET Core 6.0.25 - Shared Framework (x64)")
    'Microsoft ASP.NET Core Shared Framework'
    >>> _strip_version("Wazuh Agent")
    'Wazuh Agent'
    """
    # Remove version numbers: digit groups separated by dots (e.g. 6.0.25)
    stripped = re.sub(r"\d+(?:\.\d+)+", "", name)
    # Remove architecture tags: (x64), (x86), (arm64), etc.
    stripped = re.sub(r"\(\s*x(?:64|86)\s*\)", "", stripped, flags=re.IGNORECASE)
    stripped = re.sub(r"\(\s*arm\d*\s*\)", "", stripped, flags=re.IGNORECASE)
    # Normalise separators: " - " → " " so "Foo - Bar" == "Foo Bar"
    stripped = re.sub(r"\s*-\s*", " ", stripped)
    # Collapse multiple spaces
    stripped = re.sub(r"\s{2,}", " ", stripped)
    return stripped.strip()


def count_unique_software_packages(csv_path: Path) -> int | None:
    """Count unique software packages in the ``Package Name`` column.

    Version numbers are stripped before comparing, so
    ``Foo 6.0.25`` and ``Foo 8.0.0`` count as one package.

    Returns the count, or ``None`` if the column is not found.
    """
    try:
        with open(csv_path, newline="", encoding="utf-8-sig") as fh:
            reader = csv.DictReader(fh)
            if reader.fieldnames is None:
                return None

            target_col: str | None = None
            for col in reader.fieldnames:
                normalised = col.strip().lower().replace(" ", "_")
                if normalised == "package_name":
                    target_col = col
                    break

            if target_col is None:
                logger.info(
                    "Column 'Package Name' not found in %s (columns: %s)",
                    csv_path, reader.fieldnames,
                )
                return None

            unique_packages: set[str] = set()
            for row in reader:
                value = (row.get(target_col) or "").strip()
                if value:
                    unique_packages.add(_strip_version(value))

            logger.info(
                "Counted %d unique software packages from %s",
                len(unique_packages), csv_path,
            )
            return len(unique_packages)
    except Exception as exc:
        logger.error("Failed to read software CSV %s: %s", csv_path, exc)
        return None


def count_inventory_systems(csv_path: Path) -> int | None:
    """Count total systems in the inventory report CSV.
    
    Counts the number of data rows (excluding the header row).
    
    Parameters
    ----------
    csv_path : Path
        Path to the inventory report CSV file
    
    Returns
    -------
    int | None
        Number of systems (rows) in the inventory, or None if reading fails
    """
    try:
        with open(csv_path, newline="", encoding="utf-8-sig") as fh:
            reader = csv.DictReader(fh)
            if reader.fieldnames is None:
                return None
            
            # Count all data rows
            count = sum(1 for _ in reader)
            
            logger.info(
                "Counted %d systems in inventory from %s",
                count, csv_path,
            )
            return count
    except Exception as exc:
        logger.error("Failed to read inventory CSV %s: %s", csv_path, exc)
        return None


def count_siem_installations(csv_path: Path) -> int | None:
    """Count systems with SIEM installed from the inventory report CSV.
    
    Looks for a column named "SIEM" (case-insensitive) and counts all rows
    where the value is "installed" (case-insensitive).
    
    Parameters
    ----------
    csv_path : Path
        Path to the inventory report CSV file
    
    Returns
    -------
    int | None
        Number of systems with SIEM installed, or None if reading fails
        or the SIEM column is not found
    """
    try:
        with open(csv_path, newline="", encoding="utf-8-sig") as fh:
            reader = csv.DictReader(fh)
            if reader.fieldnames is None:
                return None
            
            # Find the SIEM column (case-insensitive)
            siem_col: str | None = None
            for col in reader.fieldnames:
                if col.strip().lower() == "siem":
                    siem_col = col
                    break
            
            if siem_col is None:
                logger.info(
                    "Column 'SIEM' not found in %s (columns: %s)",
                    csv_path, reader.fieldnames,
                )
                return None
            
            # Count rows where SIEM value is "installed"
            count = 0
            for row in reader:
                value = (row.get(siem_col) or "").strip().lower()
                if value == "installed":
                    count += 1
            
            logger.info(
                "Counted %d SIEM installations from %s",
                count, csv_path,
            )
            return count
    except Exception as exc:
        logger.error("Failed to read inventory CSV for SIEM count %s: %s", csv_path, exc)
        return None


def count_dfir_installations(csv_path: Path) -> int | None:
    """Count systems with DFIR/IR Agent installed from the inventory report CSV.
    
    Looks for a column named "IR_Agent" (case-insensitive) and counts all rows
    where the value is "installed" (case-insensitive).
    
    Parameters
    ----------
    csv_path : Path
        Path to the inventory report CSV file
    
    Returns
    -------
    int | None
        Number of systems with DFIR agent installed, or None if reading fails
        or the IR_Agent column is not found
    """
    try:
        with open(csv_path, newline="", encoding="utf-8-sig") as fh:
            reader = csv.DictReader(fh)
            if reader.fieldnames is None:
                return None
            
            # Find the IR_Agent column (case-insensitive, with various separators)
            ir_col: str | None = None
            for col in reader.fieldnames:
                normalized = col.strip().lower().replace(" ", "").replace("_", "").replace("-", "")
                if normalized == "iragent":
                    ir_col = col
                    break
            
            if ir_col is None:
                logger.info(
                    "Column 'IR_Agent' not found in %s (columns: %s)",
                    csv_path, reader.fieldnames,
                )
                return None
            
            # Count rows where IR_Agent value is "installed"
            count = 0
            for row in reader:
                value = (row.get(ir_col) or "").strip().lower()
                if value == "installed":
                    count += 1
            
            logger.info(
                "Counted %d DFIR agent installations from %s",
                count, csv_path,
            )
            return count
    except Exception as exc:
        logger.error("Failed to read inventory CSV for DFIR count %s: %s", csv_path, exc)
        return None


def count_xdr_installations(csv_path: Path) -> int | None:
    """Count systems with XDR (SentinelOne) installed from the inventory report CSV.
    
    Looks for a column named "XDR" (case-insensitive) and counts all rows
    where the value is "installed" (case-insensitive).
    
    Parameters
    ----------
    csv_path : Path
        Path to the inventory report CSV file
    
    Returns
    -------
    int | None
        Number of systems with XDR installed, or None if reading fails
        or the XDR column is not found
    """
    try:
        with open(csv_path, newline="", encoding="utf-8-sig") as fh:
            reader = csv.DictReader(fh)
            if reader.fieldnames is None:
                return None
            
            # Find the XDR column (case-insensitive)
            xdr_col: str | None = None
            for col in reader.fieldnames:
                if col.strip().lower() == "xdr":
                    xdr_col = col
                    break
            
            if xdr_col is None:
                logger.info(
                    "Column 'XDR' not found in %s (columns: %s)",
                    csv_path, reader.fieldnames,
                )
                return None
            
            # Count rows where XDR value is "installed"
            count = 0
            for row in reader:
                value = (row.get(xdr_col) or "").strip().lower()
                if value == "installed":
                    count += 1
            
            logger.info(
                "Counted %d XDR installations from %s",
                count, csv_path,
            )
            return count
    except Exception as exc:
        logger.error("Failed to read inventory CSV for XDR count %s: %s", csv_path, exc)
        return None


def _extract_update_target(action: str) -> str:
    """Extract the 'Update to higher than X.Y.Z' suffix from an action string.

    Returns the target version string (e.g. ``"6.0.27"``) or ``""`` if
    the pattern is not found.
    """
    m = re.search(r"[Uu]pdate to (?:higher than |)(\d+(?:\.\d+)+)", action)
    return m.group(1) if m else ""


def _version_tuple(version: str) -> tuple[int, ...]:
    """Parse a dotted version string into a tuple for comparison."""
    try:
        return tuple(int(x) for x in version.split("."))
    except ValueError:
        return (0,)


def build_software_patches_csv(
    csv_path: Path,
    dest_dir: Path,
) -> tuple[Path | None, int]:
    """Deduplicate a software-updates CSV into a two-column table CSV.

    Reads the raw CSV (``vulnaggregateother_*`` / ``non_windows_updates_*``),
    finds the ``Package Name`` and ``Recommended Action`` columns, strips
    version numbers from package names so that entries differing only by
    installed version are collapsed into a single row, and writes a clean
    CSV with the two target columns.

    When multiple rows collapse, the recommended action with the highest
    target version (``Update to higher than X.Y.Z``) is kept.

    Parameters
    ----------
    csv_path : Path
        Path to the downloaded software CSV.
    dest_dir : Path
        Directory to write the cleaned CSV into.

    Returns
    -------
    (clean_csv_path, unique_count)
        Path to the deduplicated CSV and the number of unique packages.
        ``(None, 0)`` if the required columns are not found.
    """
    try:
        with open(csv_path, newline="", encoding="utf-8-sig") as fh:
            reader = csv.DictReader(fh)
            if reader.fieldnames is None:
                return None, 0

            # Case-insensitive column lookup
            pkg_col: str | None = None
            action_col: str | None = None
            for col in reader.fieldnames:
                normalised = col.strip().lower().replace(" ", "_")
                if normalised == "package_name" and pkg_col is None:
                    pkg_col = col
                elif normalised == "recommended_action" and action_col is None:
                    action_col = col

            if pkg_col is None:
                logger.info(
                    "Column 'Package Name' not found in %s (columns: %s)",
                    csv_path, reader.fieldnames,
                )
                return None, 0

            # Deduplicate by version-stripped package name.
            # Keep the recommended action with the highest target version.
            seen: dict[str, str] = {}       # base_name -> best action
            best_ver: dict[str, tuple] = {}  # base_name -> highest target version tuple
            for row in reader:
                raw_name = (row.get(pkg_col) or "").strip()
                if not raw_name:
                    continue
                base_name = _strip_version(raw_name)
                action = (row.get(action_col) or "").strip() if action_col else ""

                target_ver = _version_tuple(_extract_update_target(action))
                prev_ver = best_ver.get(base_name, (0,))
                if base_name not in seen or target_ver > prev_ver:
                    # Store a clean action: "<base_name> - Update to higher than X.Y.Z"
                    seen[base_name] = action
                    best_ver[base_name] = target_ver

            if not seen:
                return None, 0

            # Build display actions: extract the "Update to ..." portion
            # so the action reads cleanly against the version-free name.
            display_rows: list[tuple[str, str]] = []
            for base_name, action in seen.items():
                # Try to extract "Update to higher than X.Y.Z"
                m = re.search(
                    r"([Uu]pdate to (?:higher than )?\d+(?:\.\d+)*)",
                    action,
                )
                clean_action = m.group(1) if m else action
                display_rows.append((base_name, clean_action))

            # Write clean two-column CSV
            dest_dir.mkdir(parents=True, exist_ok=True)
            clean_path = dest_dir / "top_software_patches.csv"
            with open(clean_path, "w", newline="", encoding="utf-8-sig") as fh:
                writer = csv.writer(fh)
                writer.writerow(["Package Name", "Recommended Action"])
                for name, action in display_rows:
                    writer.writerow([name, action])

            logger.info(
                "Built software patches CSV: %d unique packages -> %s",
                len(display_rows), clean_path,
            )
            return clean_path, len(display_rows)

    except Exception as exc:
        logger.error(
            "Failed to build software patches CSV from %s: %s", csv_path, exc,
        )
        return None, 0


def replace_kb_count_in_slides(
    pptx_path: Path,
    kb_count: int,
    software_count: int | None = None,
) -> int:
    """Replace ``[Count]`` placeholders on KB- and software-related slides.

    Targets shapes by name:
        ``cyber_insight_body``     (slide 3)
        ``required_patches_body``  (slide 14)

    Uses contextual replacement so each ``[Count]`` is replaced with the
    correct number:
      - Paragraphs containing "microsoft kb" or "windows systems"
        → *kb_count*
      - Paragraphs containing "software packages" or "third-party software"
        → *software_count* (if provided)

    Returns the total number of replacements made.
    """
    prs = PptxPresentation(str(pptx_path))
    target_shapes = {"cyber_insight_body", "required_patches_body"}
    kb_keywords = ("microsoft kb", "windows systems")
    sw_keywords = ("software packages", "third-party software")
    replacements = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name not in target_shapes:
                continue
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in para.runs).lower()

                # Determine which count applies to this paragraph
                count_value: int | None = None
                if any(kw in para_text for kw in kb_keywords):
                    count_value = kb_count
                elif software_count is not None and any(
                    kw in para_text for kw in sw_keywords
                ):
                    count_value = software_count

                if count_value is None:
                    continue

                for run in para.runs:
                    if "[Count]" in run.text:
                        run.text = run.text.replace(
                            "[Count]", str(count_value),
                        )
                        replacements += 1
                        logger.info(
                            "Replaced [Count] -> %d in shape '%s'",
                            count_value, shape.name,
                        )

    if replacements:
        prs.save(str(pptx_path))
        logger.info(
            "Patch counts: %d replacement(s) saved to %s",
            replacements, pptx_path,
        )
    return replacements


def replace_ai_insight_in_slide(
    pptx_path: Path,
    insight_text: str,
) -> bool:
    """Replace the AI insight placeholder on slide 3 (Cyber Insight Summary).

    Searches for the paragraph containing "[Additional observations and 
    recommendations to be populated]" and replaces it with the provided
    AI-generated insight text.

    Parameters
    ----------
    pptx_path : Path
        Path to the PowerPoint file
    insight_text : str
        The AI-generated observation/recommendation text

    Returns
    -------
    bool
        True if the replacement was made, False otherwise
    """
    prs = PptxPresentation(str(pptx_path))
    target_shape = "cyber_insight_body"
    placeholder = "[Additional observations and recommendations to be populated]"
    replaced = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name != target_shape:
                continue
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in para.runs)
                
                if placeholder in para_text:
                    # Replace the entire paragraph text
                    # Keep the bullet formatting but ensure text is not bolded
                    for run in para.runs:
                        if placeholder in run.text:
                            run.text = run.text.replace(
                                placeholder, 
                                f"  {insight_text}"
                            )
                            # Ensure text is not bolded
                            run.font.bold = False
                            logger.info(
                                "Replaced AI insight placeholder in shape '%s'",
                                shape.name,
                            )
                            replaced = True
                            break
                
                if replaced:
                    break
            
            if replaced:
                break
        
        if replaced:
            break

    if replaced:
        prs.save(str(pptx_path))
        logger.info("AI insight saved to %s", pptx_path)
    
    return replaced


def replace_endpoint_count_in_slide(
    pptx_path: Path,
    endpoint_count: int,
    siem_count: int | None = None,
    dfir_count: int | None = None,
    xdr_count: int | None = None,
) -> int:
    """Replace endpoint and agent count placeholders on slide 4 (Service Coverage & Deployment).

    Searches for "[Count]" in the service_coverage_body shape and replaces:
    - "[Count] systems detected" with the endpoint count
    - "EPIC SIEM: [Count] agents" with the SIEM count (if provided)
    - "SentinelOne XDR: [Count] agents" with XDR count or "Client Managed" if count is 0
    - "DFIR Agent: [Count] agents" with the DFIR count (if provided)
    - "Host-Based (OS / 3rd Party SW): [Count] systems reporting" with SIEM count (if provided)
      (Note: Host-Based count equals SIEM agent count)

    Parameters
    ----------
    pptx_path : Path
        Path to the PowerPoint file
    endpoint_count : int
        Number of systems detected from inventory report
    siem_count : int | None
        Number of systems with SIEM installed (optional)
        Also used for Host-Based vulnerability management count
    dfir_count : int | None
        Number of systems with DFIR/IR Agent installed (optional)
    xdr_count : int | None
        Number of systems with XDR installed (optional)
        If 0, displays "Client Managed" instead

    Returns
    -------
    int
        Number of replacements made (0-5)
    """
    prs = PptxPresentation(str(pptx_path))
    target_shape = "service_coverage_body"
    replacements = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name != target_shape:
                continue
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in para.runs)
                
                # Look for the endpoint deployment line
                if "systems detected" in para_text.lower() and "[Count]" in para_text:
                    for run in para.runs:
                        if "[Count]" in run.text:
                            run.text = run.text.replace(
                                "[Count]", 
                                str(endpoint_count)
                            )
                            logger.info(
                                "Replaced endpoint count [Count] -> %d in shape '%s'",
                                endpoint_count, shape.name,
                            )
                            replacements += 1
                            break
                
                # Look for the EPIC SIEM line
                elif siem_count is not None and "epic siem" in para_text.lower() and "[Count]" in para_text:
                    for run in para.runs:
                        if "[Count]" in run.text:
                            # If SIEM count is 0, show "Client Managed" instead of the number
                            if siem_count == 0:
                                # Replace "[Count] agents" with "Client Managed"
                                run.text = run.text.replace("[Count] agents", "Client Managed")
                                logger.info(
                                    "Replaced SIEM [Count] agents -> Client Managed in shape '%s'",
                                    shape.name,
                                )
                            else:
                                run.text = run.text.replace(
                                    "[Count]", 
                                    str(siem_count)
                                )
                                logger.info(
                                    "Replaced SIEM count [Count] -> %d in shape '%s'",
                                    siem_count, shape.name,
                                )
                            replacements += 1
                            break
                
                # Look for the SentinelOne XDR line
                elif xdr_count is not None and "sentinelone xdr" in para_text.lower() and "[Count]" in para_text:
                    for run in para.runs:
                        if "[Count]" in run.text:
                            # If XDR count is 0, show "Client Managed" instead of the number
                            if xdr_count == 0:
                                # Replace "[Count] agents" with "Client Managed"
                                run.text = run.text.replace("[Count] agents", "Client Managed")
                                logger.info(
                                    "Replaced XDR [Count] agents -> Client Managed in shape '%s'",
                                    shape.name,
                                )
                            else:
                                run.text = run.text.replace(
                                    "[Count]", 
                                    str(xdr_count)
                                )
                                logger.info(
                                    "Replaced XDR count [Count] -> %d in shape '%s'",
                                    xdr_count, shape.name,
                                )
                            replacements += 1
                            break
                
                # Look for the DFIR Agent line
                elif dfir_count is not None and "dfir agent" in para_text.lower() and "[Count]" in para_text:
                    for run in para.runs:
                        if "[Count]" in run.text:
                            # If DFIR count is 0, show "Client Managed" instead of the number
                            if dfir_count == 0:
                                # Replace "[Count] agents" with "Client Managed"
                                run.text = run.text.replace("[Count] agents", "Client Managed")
                                logger.info(
                                    "Replaced DFIR [Count] agents -> Client Managed in shape '%s'",
                                    shape.name,
                                )
                            else:
                                run.text = run.text.replace(
                                    "[Count]", 
                                    str(dfir_count)
                                )
                                logger.info(
                                    "Replaced DFIR count [Count] -> %d in shape '%s'",
                                    dfir_count, shape.name,
                                )
                            replacements += 1
                            break
                
                # Look for the Host-Based vulnerability management line (uses SIEM count)
                elif siem_count is not None and "host-based" in para_text.lower() and "[Count]" in para_text:
                    for run in para.runs:
                        if "[Count]" in run.text:
                            run.text = run.text.replace(
                                "[Count]", 
                                str(siem_count)
                            )
                            logger.info(
                                "Replaced Host-Based count [Count] -> %d (SIEM count) in shape '%s'",
                                siem_count, shape.name,
                            )
                            replacements += 1
                            break

    if replacements > 0:
        prs.save(str(pptx_path))
        logger.info("Service coverage counts saved to %s (%d replacements)", pptx_path, replacements)
    
    return replacements


# ---------------------------------------------------------------------------
# Copy service coverage slide content from a reference PPTX
# ---------------------------------------------------------------------------
_COUNT_PATTERNS = [
    # "123 systems detected" / "123 system detected"
    (re.compile(r"(\d+)(\s+systems?\s+detected)", re.IGNORECASE), r"[Count]\2"),
    # "123 agents" / "123 agent"
    (re.compile(r"(\d+)(\s+agents?)", re.IGNORECASE), r"[Count]\2"),
    # "123 systems reporting" / "123 system reporting"
    (re.compile(r"(\d+)(\s+systems?\s+reporting)", re.IGNORECASE), r"[Count]\2"),
    # "accounts: 123" / "account: 123"
    (re.compile(r"(accounts?:\s*)\d+", re.IGNORECASE), r"\1[Count]"),
]


def _replace_counts_with_placeholder(text: str) -> str:
    """Replace numeric count values in *text* with ``[Count]`` placeholders."""
    for pattern, replacement in _COUNT_PATTERNS:
        text = pattern.sub(replacement, text)
    return text


def copy_service_coverage_from_reference(
    target_pptx_path: Path,
    reference_pptx_path: Path,
) -> bool:
    """Copy the service coverage text frame from *reference_pptx_path* into
    the target PPTX, preserving all paragraph formatting.

    Numeric count values in the copied text are replaced with ``[Count]``
    placeholders so that ``replace_endpoint_count_in_slide`` can fill them
    with fresh inventory data afterwards.

    Parameters
    ----------
    target_pptx_path : Path
        Path to the report PPTX being generated (modified in-place).
    reference_pptx_path : Path
        Path to the most recent client PPTX downloaded from Drive.

    Returns
    -------
    bool
        ``True`` if the copy succeeded, ``False`` otherwise.
    """
    SHAPE_NAME = "service_coverage_body"
    FALLBACK_KEYWORD = "Endpoint Deployment"

    # --- Open reference and locate the source text frame -----------------
    ref_prs = PptxPresentation(str(reference_pptx_path))
    ref_txBody = None

    # First try: shape named service_coverage_body
    for slide in ref_prs.slides:
        for shape in slide.shapes:
            if shape.name == SHAPE_NAME and shape.has_text_frame:
                ref_txBody = shape.text_frame._txBody
                break
        if ref_txBody is not None:
            break

    # Fallback: any text frame containing the keyword
    if ref_txBody is None:
        for slide in ref_prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                frame_text = shape.text_frame.text
                if FALLBACK_KEYWORD in frame_text:
                    ref_txBody = shape.text_frame._txBody
                    logger.info(
                        "Reference: using fallback shape '%s' (contains '%s')",
                        shape.name, FALLBACK_KEYWORD,
                    )
                    break
            if ref_txBody is not None:
                break

    if ref_txBody is None:
        logger.warning(
            "No service coverage text frame found in reference PPTX: %s",
            reference_pptx_path,
        )
        return False

    # Deep-copy all <a:p> elements from the reference
    ap_tag = qn("a:p")
    ref_paragraphs = [
        copy.deepcopy(p) for p in ref_txBody.findall(ap_tag)
    ]
    if not ref_paragraphs:
        logger.warning("Reference text frame has no paragraphs")
        return False

    # Replace numeric counts with [Count] placeholders in the copies
    ar_tag = qn("a:r")
    at_tag = qn("a:t")
    for p_elem in ref_paragraphs:
        for run_elem in p_elem.findall(f".//{ar_tag}"):
            t_elem = run_elem.find(at_tag)
            if t_elem is not None and t_elem.text:
                t_elem.text = _replace_counts_with_placeholder(t_elem.text)

    # --- Open target and locate its text frame ---------------------------
    tgt_prs = PptxPresentation(str(target_pptx_path))
    tgt_txBody = None

    for slide in tgt_prs.slides:
        for shape in slide.shapes:
            if shape.name == SHAPE_NAME and shape.has_text_frame:
                tgt_txBody = shape.text_frame._txBody
                break
        if tgt_txBody is not None:
            break

    if tgt_txBody is None:
        logger.warning(
            "No '%s' shape found in target PPTX: %s",
            SHAPE_NAME, target_pptx_path,
        )
        return False

    # Remove existing <a:p> paragraphs from the target
    for old_p in tgt_txBody.findall(ap_tag):
        tgt_txBody.remove(old_p)

    # Append the reference paragraphs (bodyPr/lstStyle stay untouched)
    for new_p in ref_paragraphs:
        tgt_txBody.append(new_p)

    tgt_prs.save(str(target_pptx_path))
    logger.info(
        "Copied %d service coverage paragraphs from reference '%s' into '%s'",
        len(ref_paragraphs), reference_pptx_path, target_pptx_path,
    )
    return True


# ---------------------------------------------------------------------------
# Module-level cell styling (shared by old and paginated code paths)
# ---------------------------------------------------------------------------
def _lock_cell(cell, font_pt, *, bold=False, bg=CLR_TBL_DARK):
    """Style a cell with word-wrap OFF and noAutofit to prevent row expansion."""
    cell.text_frame.word_wrap = False
    cell.text_frame.margin_top = Pt(1)
    cell.text_frame.margin_bottom = Pt(1)
    cell.text_frame.margin_left = Pt(2)
    cell.text_frame.margin_right = Pt(2)
    bodyPr = cell.text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        for child in list(bodyPr):
            if child.tag.endswith("}normAutofit") or child.tag.endswith("}spAutoFit"):
                bodyPr.remove(child)
        if bodyPr.find(qn("a:noAutofit")) is None:
            bodyPr.append(bodyPr.makeelement(qn("a:noAutofit"), {}))
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_pt)
        p.font.bold = bold
        p.font.color.rgb = CLR_WHITE
        p.font.name = "Calibri"
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg


# ---------------------------------------------------------------------------
# Pagination helpers
# ---------------------------------------------------------------------------
def _get_slide_title(slide) -> str:
    """Read text from the shape named ``slide_title``."""
    for shape in slide.shapes:
        if shape.name == "slide_title" and shape.has_text_frame:
            return shape.text_frame.text
    return ""


def _update_slide_title(slide, new_title: str) -> None:
    """Update runs in the ``slide_title`` shape."""
    for shape in slide.shapes:
        if shape.name == "slide_title" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = new_title
            return


def _remove_shape_by_name(slide, name: str) -> bool:
    """Remove the first shape matching *name*, falling back to first table.

    Returns ``True`` if a shape was removed.
    """
    # Try by name first
    for shape in list(slide.shapes):
        if shape.name == name:
            shape._element.getparent().remove(shape._element)
            return True
    # Fallback: first table on the slide
    for shape in list(slide.shapes):
        if shape.has_table:
            shape._element.getparent().remove(shape._element)
            return True
    return False


def _compute_col_widths_for_pagination(
    headers: list[str],
    rows: list[list[str]],
    total_emu: int,
) -> list[int]:
    """Proportional column widths using same algorithm as table_optimizer.

    Computes across ALL rows for cross-page consistency.
    """
    n = len(headers)
    if n == 0:
        return []
    if n == 1:
        return [total_emu]

    col_max = [len(h) for h in headers]
    col_sum = [float(len(h)) for h in headers]
    for row in rows:
        for ci in range(n):
            ln = len(row[ci]) if ci < len(row) else 0
            if ln > col_max[ci]:
                col_max[ci] = ln
            col_sum[ci] += ln

    denom = len(rows) + 1
    col_avg = [s / denom for s in col_sum]

    scores = [
        max(3.0, mx * 0.7 + av * 0.3)
        for mx, av in zip(col_max, col_avg)
    ]
    score_sum = sum(scores)
    min_col_w = int(0.55 * EMU_PER_INCH)
    widths = [max(min_col_w, int(total_emu * s / score_sum)) for s in scores]

    current = sum(widths)
    if current != total_emu and current > 0:
        factor = total_emu / current
        widths = [int(w * factor) for w in widths]
        widths[-1] += total_emu - sum(widths)

    return widths


def _build_paginated_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    col_widths: list[int],
    shape_name: str,
) -> None:
    """Create one table page at TABLE_BOX_INCHES, constrained to 3.0" height."""
    n_cols = len(headers)
    n_data = len(rows)
    n_total = 1 + max(n_data, 1)

    tbl_left_emu = int(TABLE_BOX_INCHES["left"] * EMU_PER_INCH)
    tbl_top_emu = int(TABLE_BOX_INCHES["top"] * EMU_PER_INCH)
    tbl_width_emu = int(TABLE_BOX_INCHES["width"] * EMU_PER_INCH)
    tbl_height_emu = int(TABLE_BOX_INCHES["height"] * EMU_PER_INCH)

    row_height_emu = max(1, tbl_height_emu // n_total)

    # Font sizing — reduced for better fit in 3.0" height
    if n_total > 8:
        header_pt, data_pt = 7, 6
    elif n_total > 5:
        header_pt, data_pt = 8, 7
    else:
        header_pt, data_pt = 9, 8

    shape = slide.shapes.add_table(
        n_total, n_cols,
        Emu(tbl_left_emu), Emu(tbl_top_emu),
        Emu(tbl_width_emu), Emu(tbl_height_emu),
    )
    shape.name = shape_name
    tbl = shape.table

    # Column widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    # Row heights
    for r in range(n_total):
        tbl.rows[r].height = row_height_emu

    # Header
    for ci, col_name in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = col_name
        _lock_cell(cell, header_pt, bold=True, bg=CLR_BLACK)

    # Data
    if n_data == 0:
        for ci in range(n_cols):
            cell = tbl.cell(1, ci)
            cell.text = "No data" if ci == 0 else ""
            _lock_cell(cell, data_pt, bg=CLR_TBL_DARK)
    else:
        for ri, row_data in enumerate(rows):
            row_idx = ri + 1
            row_clr = CLR_TBL_DARK if row_idx % 2 == 1 else CLR_TBL_LIGHT
            for ci in range(n_cols):
                cell = tbl.cell(row_idx, ci)
                cell.text = row_data[ci] if ci < len(row_data) else ""
                _lock_cell(cell, data_pt, bg=row_clr)

    # Re-enforce dimensions to prevent PowerPoint expansion
    for r in range(n_total):
        tbl.rows[r].height = row_height_emu
    shape.height = Emu(tbl_height_emu)
    shape.width = Emu(tbl_width_emu)


def _add_tip_text(slide, text: str) -> None:
    """Add small italic gray text below the table on page 1."""
    tip_top = TABLE_BOX_INCHES["top"] + TABLE_BOX_INCHES["height"] + 0.05
    box = slide.shapes.add_textbox(
        Inches(TABLE_BOX_INCHES["left"]),
        Inches(tip_top),
        Inches(TABLE_BOX_INCHES["width"]),
        Inches(0.25),
    )
    box.name = "pagination_tip"
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7)
    p.font.italic = True
    p.font.color.rgb = CLR_MUTED
    p.font.name = "Calibri"


# ---------------------------------------------------------------------------
# Table placement (single slide per table, no pagination)
# ---------------------------------------------------------------------------
def place_csv_tables_paginated(
    pptx_path: Path,
    assets: list[CSVTableAsset],
    *,
    max_rows_per_page: int = MAX_ROWS_PER_PAGE,
) -> tuple[int, set[int]]:
    """Place CSV data onto slides — one table per slide (no pagination).

    Reads each CSV, filters source columns, removes the old placeholder
    shape, and builds a single table with ALL data rows on the existing
    slide.  The table optimizer will later convert these to images.

    Returns ``(tables_placed, empty_set)`` for API compatibility.
    """
    if not assets:
        return 0, set()

    prs = PptxPresentation(str(pptx_path))
    count = 0

    for asset in assets:
        if asset.slide_index >= len(prs.slides):
            logger.warning(
                "Slide %d out of range, skipping CSV '%s'",
                asset.slide_index, asset.original_filename,
            )
            continue

        # --- Read CSV ---------------------------------------------------------
        try:
            with open(asset.local_path, newline="", encoding="utf-8-sig") as fh:
                reader = csv.reader(fh)
                all_rows = list(reader)
        except Exception as exc:
            logger.error(
                "Failed to read CSV '%s': %s", asset.local_path, exc,
            )
            continue

        if not all_rows:
            logger.warning("CSV '%s' is empty, skipping", asset.original_filename)
            continue

        # Skip metadata first row
        if len(all_rows) >= 2 and all_rows[0][0].strip() == "":
            all_rows = all_rows[1:]

        if not all_rows:
            continue

        # Filter source columns
        header = all_rows[0]
        keep_indices = [
            i for i, col in enumerate(header) if not _is_source_column(col)
        ]
        if not keep_indices:
            logger.warning(
                "CSV '%s' has no columns after removing source columns",
                asset.original_filename,
            )
            continue

        filtered_header = [header[i] for i in keep_indices]
        filtered_rows = [
            [row[i] if i < len(row) else "" for i in keep_indices]
            for row in all_rows[1:]
        ]

        n_data_rows = len(filtered_rows)

        slide = prs.slides[asset.slide_index]

        # Compute col widths across ALL rows
        tbl_width_emu = int(TABLE_BOX_INCHES["width"] * EMU_PER_INCH)
        col_widths = _compute_col_widths_for_pagination(
            filtered_header, filtered_rows, tbl_width_emu,
        )

        # Remove old shape from existing slide
        _remove_shape_by_name(slide, asset.shape_name)

        # Build a single table with ALL rows (no pagination)
        _build_paginated_table(
            slide, filtered_header, filtered_rows, col_widths, asset.shape_name,
        )

        count += 1
        logger.info(
            "Placed '%s' on slide %d: %d data rows (single table)",
            asset.original_filename, asset.slide_index, n_data_rows,
        )

    prs.save(str(pptx_path))
    logger.info("Saved PPTX with %d CSV tables -> %s", count, pptx_path)
    return count, set()


# ---------------------------------------------------------------------------
# Post-processing: populate table slides from CSV data (DEPRECATED)
# ---------------------------------------------------------------------------
def place_csv_tables(pptx_path: Path, assets: list[CSVTableAsset]) -> int:
    """Open a generated PPTX and populate table shapes from CSV data.

    .. deprecated::
        Use :func:`place_csv_tables_paginated` instead, which splits
        large tables across multiple slides.

    Returns the number of tables successfully populated.
    """
    if not assets:
        return 0

    prs = PptxPresentation(str(pptx_path))
    populated = 0

    for asset in assets:
        if asset.slide_index >= len(prs.slides):
            logger.warning(
                "Slide %d out of range, skipping CSV '%s'",
                asset.slide_index, asset.original_filename,
            )
            continue

        # --- Read CSV ---------------------------------------------------------
        try:
            with open(asset.local_path, newline="", encoding="utf-8-sig") as fh:
                reader = csv.reader(fh)
                all_rows = list(reader)
        except Exception as exc:
            logger.error(
                "Failed to read CSV '%s': %s", asset.local_path, exc
            )
            continue

        if not all_rows:
            logger.warning("CSV '%s' is empty, skipping", asset.original_filename)
            continue

        # --- Skip metadata first row if present --------------------------------
        # NDR CSVs have a metadata row like: ,wiegmann1,2026-01-04,Sensitive Data
        # Detected by an empty first cell in row 0.
        if len(all_rows) >= 2 and all_rows[0][0].strip() == "":
            logger.info(
                "CSV '%s' has metadata first row, skipping it",
                asset.original_filename,
            )
            all_rows = all_rows[1:]

        # --- Filter out source columns ----------------------------------------
        header = all_rows[0]
        keep_indices = [
            i for i, col in enumerate(header) if not _is_source_column(col)
        ]
        if not keep_indices:
            logger.warning(
                "CSV '%s' has no columns left after removing source columns",
                asset.original_filename,
            )
            continue

        filtered_header = [header[i] for i in keep_indices]
        filtered_rows = [
            [row[i] if i < len(row) else "" for i in keep_indices]
            for row in all_rows[1:]
        ]

        n_cols = len(filtered_header)
        n_data_rows = len(filtered_rows)

        # --- Find and replace the existing shape --------------------------------
        slide = prs.slides[asset.slide_index]
        old_shape = None
        is_table = False

        # Try table shape first (existing behaviour for slides 14-15)
        for shape in slide.shapes:
            if shape.name == asset.shape_name and shape.has_table:
                old_shape = shape
                is_table = True
                break

        # Fallback: any shape with matching name (NDR placeholder shapes)
        if old_shape is None:
            for shape in slide.shapes:
                if shape.name == asset.shape_name:
                    old_shape = shape
                    break

        # Fallback for table slides 14/15: use first table on slide if name match failed
        if old_shape is None and asset.slide_index in (14, 15):
            for shape in slide.shapes:
                if shape.has_table:
                    old_shape = shape
                    is_table = True
                    logger.info(
                        "Using first table on slide %d for CSV '%s' (name '%s' not found)",
                        asset.slide_index, asset.original_filename, asset.shape_name,
                    )
                    break

        if old_shape is None:
            logger.warning(
                "Shape '%s' not found on slide %d, skipping CSV '%s'",
                asset.shape_name, asset.slide_index, asset.original_filename,
            )
            continue

        # Use conservative table box to ensure tables never extend beyond visible slide area
        tbl_left_emu = int(TABLE_BOX_INCHES["left"] * EMU_PER_INCH)
        tbl_top_emu = int(TABLE_BOX_INCHES["top"] * EMU_PER_INCH)
        tbl_width_emu = int(TABLE_BOX_INCHES["width"] * EMU_PER_INCH)
        max_table_height_emu = int(TABLE_BOX_INCHES["height"] * EMU_PER_INCH)
        max_table_height_emu = max(1, max_table_height_emu)
        # Cap number of rows so table fits in the fixed box. Use larger min height to prevent too many rows.
        MIN_ROW_HEIGHT_INCHES = 0.25
        min_row_emu = int(MIN_ROW_HEIGHT_INCHES * EMU_PER_INCH)
        max_data_rows = max(0, (max_table_height_emu // min_row_emu) - 1)
        original_count = len(filtered_rows)
        if original_count > max_data_rows:
            filtered_rows = filtered_rows[:max_data_rows]
            overflow = original_count - max_data_rows
            filtered_rows.append([f"... and {overflow} more row(s)"] + [""] * (n_cols - 1))
        n_data_rows = len(filtered_rows)
        n_total_rows = 1 + max(n_data_rows, 1)
        # Fixed table size = heatmap box; row height = box height / row count
        row_height_emu = max_table_height_emu // n_total_rows
        row_height_emu = max(1, row_height_emu)
        tbl_height_emu = max_table_height_emu  # same height as heatmap for every table

        # Remove old shape (table or placeholder)
        sp = old_shape._element
        sp.getparent().remove(sp)

        # --- Create new table (same position and size as heatmap image) -------
        new_shape = slide.shapes.add_table(
            n_total_rows,
            n_cols,
            Emu(tbl_left_emu),
            Emu(tbl_top_emu),
            Emu(tbl_width_emu),
            Emu(tbl_height_emu),
        )
        new_shape.name = asset.shape_name
        tbl = new_shape.table

        # Lock the shape height to prevent PowerPoint from expanding beyond fixed size
        new_shape.height = Emu(tbl_height_emu)
        new_shape.width = Emu(tbl_width_emu)

        # Set every row to the same fixed height (EMU); must be set before and after content
        for r in range(n_total_rows):
            tbl.rows[r].height = row_height_emu

        # Scale font size aggressively based on row count so text fits without expanding rows
        if n_total_rows > 20:
            header_pt, data_pt = 6, 5
        elif n_total_rows > 15:
            header_pt, data_pt = 7, 6
        elif n_total_rows > 10:
            header_pt, data_pt = 8, 7
        else:
            header_pt, data_pt = 9, 8

        # --- Populate header --------------------------------------------------
        for ci, col_name in enumerate(filtered_header):
            cell = tbl.cell(0, ci)
            cell.text = col_name
            _lock_cell(cell, header_pt, bold=True, bg=CLR_BLACK)

        # --- Populate data rows -----------------------------------------------
        if n_data_rows == 0:
            for ci in range(n_cols):
                cell = tbl.cell(1, ci)
                cell.text = "No data" if ci == 0 else ""
                _lock_cell(cell, data_pt, bg=CLR_TBL_DARK)
        else:
            for ri, row_data in enumerate(filtered_rows):
                row_idx = ri + 1
                row_clr = CLR_TBL_DARK if row_idx % 2 == 1 else CLR_TBL_LIGHT
                for ci in range(n_cols):
                    cell = tbl.cell(row_idx, ci)
                    cell.text = row_data[ci] if ci < len(row_data) else ""
                    _lock_cell(cell, data_pt, bg=row_clr)

        # Re-apply row heights so table fills the fixed box (same size as heatmap)
        for r in range(n_total_rows):
            tbl.rows[r].height = row_height_emu
        new_shape.height = Emu(tbl_height_emu)

        populated += 1
        logger.info(
            "Populated table '%s' on slide %d: %d cols x %d rows from '%s'",
            asset.shape_name, asset.slide_index,
            n_cols, n_data_rows, asset.original_filename,
        )

    prs.save(str(pptx_path))
    logger.info(
        "Saved PPTX with %d CSV tables populated -> %s", populated, pptx_path
    )
    return populated


# ---------------------------------------------------------------------------
# Multi-sensor placement functions
# ---------------------------------------------------------------------------
def place_ndr_images_multi_sensor(
    pptx_path: Path,
    sensor_assets: dict[str, list[NDRAsset]],
    sensor_slide_map: dict[str, list[int]]
) -> dict[str, int]:
    """Place NDR images for multiple sensors.
    
    Args:
        pptx_path: Path to PPTX file
        sensor_assets: Mapping of sensor_id -> [NDRAssets]
        sensor_slide_map: Mapping of sensor_id -> [slide_indices]
    
    Returns:
        Mapping of sensor_id -> count_of_images_placed
    """
    if not sensor_assets:
        return {}

    from pptx.util import Inches

    prs = PptxPresentation(str(pptx_path))
    box = Rect(
        x=int(Inches(NDR_BOX_INCHES["left"])),
        y=int(Inches(NDR_BOX_INCHES["top"])),
        width=int(Inches(NDR_BOX_INCHES["width"])),
        height=int(Inches(NDR_BOX_INCHES["height"])),
    )

    placed_counts: dict[str, int] = {}

    for sensor_id, assets in sensor_assets.items():
        placed = 0
        slide_indices = sensor_slide_map.get(sensor_id, [])
        
        if not slide_indices:
            logger.warning("No slide mapping found for sensor '%s'", sensor_id)
            continue

        for asset in assets:
            # Map asset's base slide index to sensor-specific slide index
            try:
                base_idx = asset.slide_index
                # Find position in base slides [5,6,7,8,9,10,11,12]
                base_slides = [5, 6, 7, 8, 9, 10, 11, 12]
                if base_idx in base_slides:
                    offset = base_slides.index(base_idx)
                    if offset < len(slide_indices):
                        actual_slide_idx = slide_indices[offset]
                    else:
                        logger.warning(
                            "Sensor '%s': offset %d out of range for slide mapping",
                            sensor_id, offset
                        )
                        continue
                else:
                    actual_slide_idx = base_idx
                    
            except (ValueError, IndexError) as e:
                logger.warning(
                    "Sensor '%s': failed to map slide %d: %s",
                    sensor_id, asset.slide_index, e
                )
                continue

            if actual_slide_idx >= len(prs.slides):
                logger.warning(
                    "Sensor '%s': Slide %d out of range, skipping '%s'",
                    sensor_id, actual_slide_idx, asset.original_filename,
                )
                continue

            slide = prs.slides[actual_slide_idx]

            # Build shape name with sensor suffix
            shape_name_with_sensor = f"{asset.shape_name}_{sensor_id.lower()}"

            # Remove the placeholder shape (try sensor-specific name first, then base name)
            removed = False
            for shape in list(slide.shapes):
                if shape.name in (shape_name_with_sensor, asset.shape_name):
                    sp = shape._element
                    sp.getparent().remove(sp)
                    removed = True
                    break

            if not removed:
                logger.debug(
                    "Sensor '%s': Shape '%s' not found on slide %d",
                    sensor_id, shape_name_with_sensor, actual_slide_idx
                )

            # Read image dimensions
            with PILImage.open(asset.local_path) as img:
                w_px, h_px = img.size
                dpi_x, dpi_y = img.info.get("dpi", (96, 96))
                dpi_x = dpi_x if dpi_x and dpi_x > 0 else 96
                dpi_y = dpi_y if dpi_y and dpi_y > 0 else 96

            w_emu = int(round(w_px / dpi_x * 914400))
            h_emu = int(round(h_px / dpi_y * 914400))
            source = Rect(x=0, y=0, width=w_emu, height=h_emu)

            # Contain-fit transform
            transform = compute_contain_transform(source, box)

            slide.shapes.add_picture(
                str(asset.local_path),
                Emu(transform.dest_x),
                Emu(transform.dest_y),
                Emu(transform.dest_width),
                Emu(transform.dest_height),
            )
            placed += 1
            logger.info(
                "Sensor '%s': Placed '%s' on slide %d [%dx%d -> %dx%d]",
                sensor_id, asset.original_filename, actual_slide_idx,
                w_emu, h_emu, transform.dest_width, transform.dest_height,
            )

        placed_counts[sensor_id] = placed
        logger.info("Sensor '%s': Placed %d NDR images", sensor_id, placed)

    prs.save(str(pptx_path))
    total_placed = sum(placed_counts.values())
    logger.info("Saved PPTX with %d total NDR images (multi-sensor) -> %s", 
                total_placed, pptx_path)
    return placed_counts


def place_csv_tables_multi_sensor(
    pptx_path: Path,
    sensor_assets: dict[str, list[CSVTableAsset]],
    sensor_slide_map: dict[str, list[int]]
) -> dict[str, int]:
    """Place CSV tables for multiple sensors.
    
    Args:
        pptx_path: Path to PPTX file
        sensor_assets: Mapping of sensor_id -> [CSVTableAssets]
        sensor_slide_map: Mapping of sensor_id -> [slide_indices]
    
    Returns:
        Mapping of sensor_id -> count_of_tables_placed
    """
    if not sensor_assets:
        return {}

    prs = PptxPresentation(str(pptx_path))
    placed_counts: dict[str, int] = {}

    for sensor_id, assets in sensor_assets.items():
        count = 0
        slide_indices = sensor_slide_map.get(sensor_id, [])
        
        if not slide_indices:
            logger.warning("No slide mapping found for sensor '%s'", sensor_id)
            continue

        for asset in assets:
            # Map asset's base slide index to sensor-specific slide index
            try:
                base_idx = asset.slide_index
                base_slides = [5, 6, 7, 8, 9, 10, 11, 12]
                if base_idx in base_slides:
                    offset = base_slides.index(base_idx)
                    if offset < len(slide_indices):
                        actual_slide_idx = slide_indices[offset]
                    else:
                        actual_slide_idx = base_idx
                else:
                    actual_slide_idx = base_idx
                    
            except (ValueError, IndexError):
                actual_slide_idx = asset.slide_index

            if actual_slide_idx >= len(prs.slides):
                logger.warning(
                    "Sensor '%s': Slide %d out of range, skipping CSV '%s'",
                    sensor_id, actual_slide_idx, asset.original_filename,
                )
                continue

            # Read CSV
            try:
                with open(asset.local_path, newline="", encoding="utf-8-sig") as fh:
                    reader = csv.reader(fh)
                    all_rows = list(reader)
            except Exception as exc:
                logger.error(
                    "Sensor '%s': Failed to read CSV '%s': %s",
                    sensor_id, asset.local_path, exc,
                )
                continue

            if not all_rows:
                logger.warning(
                    "Sensor '%s': CSV '%s' is empty, skipping",
                    sensor_id, asset.original_filename
                )
                continue

            # Skip metadata first row
            if len(all_rows) >= 2 and all_rows[0][0].strip() == "":
                all_rows = all_rows[1:]

            if not all_rows:
                continue

            # Filter source columns
            header = all_rows[0]
            keep_indices = [
                i for i, col in enumerate(header) if not _is_source_column(col)
            ]
            if not keep_indices:
                logger.warning(
                    "Sensor '%s': CSV '%s' has no columns after removing source columns",
                    sensor_id, asset.original_filename,
                )
                continue

            filtered_header = [header[i] for i in keep_indices]
            filtered_rows = [
                [row[i] if i < len(row) else "" for i in keep_indices]
                for row in all_rows[1:]
            ]

            slide = prs.slides[actual_slide_idx]

            # Compute col widths across ALL rows
            tbl_width_emu = int(TABLE_BOX_INCHES["width"] * EMU_PER_INCH)
            col_widths = _compute_col_widths_for_pagination(
                filtered_header, filtered_rows, tbl_width_emu,
            )

            # Build shape name with sensor suffix
            shape_name_with_sensor = f"{asset.shape_name}_{sensor_id.lower()}"

            # Remove old shape from existing slide (try both names)
            for name_to_try in (shape_name_with_sensor, asset.shape_name):
                if _remove_shape_by_name(slide, name_to_try):
                    break

            # Build table
            _build_paginated_table(
                slide, filtered_header, filtered_rows, col_widths, 
                shape_name_with_sensor,
            )

            count += 1
            logger.info(
                "Sensor '%s': Placed CSV '%s' on slide %d: %d data rows",
                sensor_id, asset.original_filename, actual_slide_idx, len(filtered_rows),
            )

        placed_counts[sensor_id] = count
        logger.info("Sensor '%s': Placed %d CSV tables", sensor_id, count)

    prs.save(str(pptx_path))
    total_placed = sum(placed_counts.values())
    logger.info("Saved PPTX with %d total CSV tables (multi-sensor) -> %s",
                total_placed, pptx_path)
    return placed_counts
