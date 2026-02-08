"""
Nuclei vulnerability report parser and slide population.

Parses ``nuclei.txt`` files produced by the Nuclei scanner, generates
mitigation recommendations, and populates the Internal Network
Vulnerabilities / Mitigation slides in the PPTX report.

Handles pagination by duplicating slides when the number of
vulnerabilities exceeds the per-slide capacity (~3 per slide).
"""

from __future__ import annotations

import logging
import math
import re
from dataclasses import dataclass
from pathlib import Path

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from core.slide_utils import duplicate_slide, update_slide_title

logger = logging.getLogger("shm.nuclei_parser")

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
MAX_PER_SLIDE = 2  # Reduced from 3 to prevent overflow with consolidated multi-target vulnerabilities

# Content area dimensions (from create_template.py)
CONTENT_W = 7.12
CONTENT_L = (10 - CONTENT_W) / 2   # 1.44"
BODY_TOP = 1.10

TEXTBOX_LEFT = CONTENT_L + 0.20     # 1.64"
TEXTBOX_TOP = BODY_TOP + 0.90       # 2.00"
TEXTBOX_WIDTH = CONTENT_W - 0.40    # 6.72"
TEXTBOX_HEIGHT = 2.98                # Extends to footer: BODY_TOP(1.10) + 0.90 + 2.98 = 4.98 (footer at 5.08, 0.10" padding)

CLR_DARK = RGBColor(0x1B, 0x21, 0x2C)
CLR_BLACK = RGBColor(0x00, 0x00, 0x00)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_PANEL = RGBColor(0xD6, 0xDB, 0xDF)

# Severity colours used for the severity tag in vulnerability headers
SEVERITY_COLOURS: dict[str, RGBColor] = {
    "critical": RGBColor(0xCC, 0x00, 0x00),
    "high":     RGBColor(0xE0, 0x4B, 0x00),
    "medium":   RGBColor(0xE6, 0x9B, 0x00),
    "low":      RGBColor(0x2E, 0x7D, 0x32),
    "info":     RGBColor(0x41, 0x69, 0xE1),
}

# Regex for parsing nuclei output lines (without numbering)
NUCLEI_LINE_RE = re.compile(
    r"^\[([^\]]+)\]\s+\[([^\]]+)\]\s+\[([^\]]+)\]\s+(.+)$"
)


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------
@dataclass
class Vulnerability:
    """A single vulnerability entry parsed from nuclei.txt."""
    index: int           # Entry number (1, 2, 3...)
    vuln_id: str         # e.g. "ldap-anonymous-login-detect"
    protocol: str        # e.g. "javascript", "tcp"
    severity: str        # "critical", "high", "medium", "low", "info"
    target: str          # IP:port or URL
    description: str     # "Issue:" text (may be empty)
    raw_line: str        # Original header line for fallback


# ---------------------------------------------------------------------------
# Parser
# ---------------------------------------------------------------------------
def parse_nuclei_file(file_path: Path) -> list[Vulnerability]:
    """Parse a nuclei.txt file into a list of Vulnerability objects.

    Expected format per entry::

        [vuln-id] [protocol] [severity] target
        Issue: description text...
        (optional continuation lines)

    Entries are auto-numbered sequentially (1, 2, 3...).
    Vulnerabilities with the same vuln_id are consolidated with multiple targets.
    Returns an empty list if the file is missing or empty.
    """
    if not file_path or not file_path.exists():
        logger.info("Nuclei file not found: %s", file_path)
        return []

    try:
        text = file_path.read_text(encoding="utf-8", errors="replace").strip()
    except Exception as exc:
        logger.error("Failed to read nuclei file %s: %s", file_path, exc)
        return []

    if not text:
        logger.info("Nuclei file is empty: %s", file_path)
        return []

    # First pass: parse all entries
    raw_vulns: list[Vulnerability] = []
    lines = text.splitlines()
    i = 0
    temp_index = 1

    while i < len(lines):
        line = lines[i].strip()
        m = NUCLEI_LINE_RE.match(line)
        if m:
            vuln_id = m.group(1).strip()
            protocol = m.group(2).strip()
            severity = m.group(3).strip().lower()
            target = m.group(4).strip()
            raw_line = line

            # Collect subsequent "Issue:" lines as description
            description_parts: list[str] = []
            j = i + 1
            while j < len(lines):
                next_line = lines[j].strip()
                # Stop if we hit the next numbered entry or an empty line
                if not next_line or NUCLEI_LINE_RE.match(next_line):
                    break
                if next_line.lower().startswith("issue:"):
                    description_parts.append(next_line)
                else:
                    # Continuation of previous description
                    description_parts.append(next_line)
                j += 1

            description = " ".join(description_parts).strip()
            i = j

            raw_vulns.append(Vulnerability(
                index=temp_index,  # Temporary, will be renumbered after consolidation
                vuln_id=vuln_id,
                protocol=protocol,
                severity=severity,
                target=target,
                description=description,
                raw_line=raw_line,
            ))
            temp_index += 1
        else:
            i += 1

    # Second pass: consolidate duplicates by vuln_id
    # Group by vuln_id, keeping first occurrence's protocol, severity, description
    consolidated: dict[str, Vulnerability] = {}
    
    for vuln in raw_vulns:
        if vuln.vuln_id in consolidated:
            # Add this target to the existing entry
            existing = consolidated[vuln.vuln_id]
            # Append target if not already included
            if vuln.target not in existing.target:
                existing.target = f"{existing.target}, {vuln.target}"
        else:
            # First occurrence - add to consolidated list
            consolidated[vuln.vuln_id] = vuln
    
    # Convert back to list and renumber
    vulns = list(consolidated.values())
    for idx, vuln in enumerate(vulns, start=1):
        vuln.index = idx

    logger.info("Parsed %d vulnerabilities from %s (consolidated from %d raw entries)", 
                len(vulns), file_path, len(raw_vulns))
    return vulns


# ---------------------------------------------------------------------------
# AI-Generated Descriptions
# ---------------------------------------------------------------------------
def generate_vulnerability_description(vuln: Vulnerability) -> str:
    """Generate an AI-style description for a vulnerability based on its CVE/ID and severity.
    
    Returns a concise, technical description explaining the vulnerability's impact.
    """
    vuln_id_lower = vuln.vuln_id.lower()
    severity = vuln.severity.lower()
    
    # Check for specific vulnerability patterns and generate appropriate descriptions
    
    # CVE-2023-34048 (VMware vCenter)
    if "cve-2023-34048" in vuln_id_lower:
        return ("This vulnerability allows unauthenticated remote code execution in VMware vCenter Server. "
                "An attacker with network access can trigger this flaw to execute arbitrary code with elevated privileges, "
                "potentially leading to complete system compromise.")
    
    # LDAP Anonymous Login
    if "ldap-anonymous" in vuln_id_lower:
        return ("LDAP server allows anonymous bind operations, enabling unauthenticated access to directory information. "
                "This can expose sensitive organizational data including user accounts, group memberships, and network topology.")
    
    # SSH/Terrapin
    if "cve-2023-48795" in vuln_id_lower or "terrapin" in vuln_id_lower:
        return ("SSH service is vulnerable to the Terrapin attack (CVE-2023-48795), a prefix truncation attack that can "
                "downgrade connection security by manipulating sequence numbers during the handshake phase.")
    
    # Default credentials
    if "default" in vuln_id_lower and ("credential" in vuln_id_lower or "login" in vuln_id_lower or "password" in vuln_id_lower):
        return ("System is configured with default factory credentials. This allows unauthorized access using "
                "publicly available default usernames and passwords, representing a critical security misconfiguration.")
    
    # SSL/TLS issues
    if any(x in vuln_id_lower for x in ["ssl", "tls", "certificate"]):
        return ("SSL/TLS configuration weakness detected. This may allow man-in-the-middle attacks, protocol downgrade attacks, "
                "or exposure of encrypted communications through cipher suite vulnerabilities or certificate validation issues.")
    
    # HTTP misconfigurations
    if any(x in vuln_id_lower for x in ["http-", "apache", "nginx", "iis"]):
        return ("Web server misconfiguration or known vulnerability detected. This may enable information disclosure, "
                "directory traversal, or unauthorized access through improper security headers or outdated server versions.")
    
    # Database exposures
    if any(x in vuln_id_lower for x in ["mysql", "postgres", "mongodb", "redis", "sql"]):
        return ("Database service exposure or misconfiguration detected. This could allow unauthorized data access, "
                "extraction of sensitive information, or potential data manipulation through weak authentication or exposed interfaces.")
    
    # SMB/Windows file sharing
    if any(x in vuln_id_lower for x in ["smb", "cifs", "netbios"]):
        return ("SMB/CIFS file sharing vulnerability detected. This may enable lateral movement, unauthorized file access, "
                "or remote code execution through exposed shares or protocol vulnerabilities.")
    
    # RDP issues
    if "rdp" in vuln_id_lower:
        return ("Remote Desktop Protocol (RDP) service vulnerability or exposure detected. This could enable brute-force attacks, "
                "credential harvesting, or exploitation of RDP-specific vulnerabilities for unauthorized remote access.")
    
    # DNS issues
    if "dns" in vuln_id_lower:
        return ("DNS service misconfiguration or vulnerability detected. This may allow zone transfer attacks, "
                "DNS poisoning, or information disclosure about internal network structure.")
    
    # SNMP exposures
    if "snmp" in vuln_id_lower:
        return ("SNMP service running with weak community strings or outdated protocol version. This enables "
                "unauthorized device configuration access and network reconnaissance through MIB tree enumeration.")
    
    # Generic by severity
    if severity == "critical":
        return (f"Critical severity vulnerability ({vuln.vuln_id}) detected that requires immediate attention. "
                "This issue could lead to complete system compromise, remote code execution, or unauthorized administrative access.")
    elif severity == "high":
        return (f"High severity vulnerability ({vuln.vuln_id}) identified that significantly impacts system security. "
                "Exploitation could result in unauthorized access, data breach, or service disruption.")
    elif severity == "medium":
        return (f"Medium severity vulnerability ({vuln.vuln_id}) detected. While not immediately critical, "
                "this issue could be chained with other vulnerabilities or exploited under specific conditions.")
    else:
        return (f"Network vulnerability ({vuln.vuln_id}) identified that should be addressed as part of "
                "comprehensive security hardening efforts to reduce the overall attack surface.")


# ---------------------------------------------------------------------------
# Mitigation rules
# ---------------------------------------------------------------------------
MITIGATION_RULES: list[tuple[list[str], str]] = [
    # (keyword patterns matched against vuln_id + description, remediation text)
    (
        ["ldap", "anonymous"],
        "Disable anonymous LDAP binding. Configure LDAP servers to require "
        "authenticated access and enforce strong credential policies.",
    ),
    (
        ["smb", "signing"],
        "Enable SMB signing on all Windows hosts via Group Policy. "
        "Set 'Microsoft network server: Digitally sign communications "
        "(always)' to Enabled.",
    ),
    (
        ["ssl", "expired"],
        "Replace expired SSL/TLS certificates with valid ones from a "
        "trusted Certificate Authority. Implement certificate expiry monitoring.",
    ),
    (
        ["ssl", "self-signed"],
        "Replace self-signed certificates with certificates issued by a "
        "trusted Certificate Authority.",
    ),
    (
        ["ssl", "weak", "cipher"],
        "Disable weak SSL/TLS cipher suites. Configure servers to use "
        "only TLS 1.2+ with strong cipher suites (AES-GCM, ChaCha20).",
    ),
    (
        ["tls", "version", "deprecated"],
        "Disable deprecated TLS versions (TLS 1.0, TLS 1.1). Configure "
        "servers to require TLS 1.2 or higher.",
    ),
    (
        ["ssh", "weak"],
        "Update SSH server configuration to disable weak algorithms. "
        "Use strong key exchange, cipher, and MAC algorithms.",
    ),
    (
        ["ssh", "password"],
        "Disable SSH password authentication where possible. Enforce "
        "key-based authentication and implement fail2ban or similar "
        "brute-force protection.",
    ),
    (
        ["ftp", "anonymous"],
        "Disable anonymous FTP access. If FTP is required, enforce "
        "authenticated access and consider migrating to SFTP.",
    ),
    (
        ["dns", "zone-transfer"],
        "Restrict DNS zone transfers to authorized secondary DNS servers "
        "only. Configure allow-transfer ACLs.",
    ),
    (
        ["snmp", "default", "community"],
        "Change default SNMP community strings. Migrate to SNMPv3 with "
        "authentication and encryption.",
    ),
    (
        ["snmp", "public"],
        "Change the SNMP community string from 'public' to a strong, "
        "unique value. Preferably migrate to SNMPv3.",
    ),
    (
        ["telnet"],
        "Disable Telnet and migrate to SSH for remote management. "
        "Telnet transmits credentials in plaintext.",
    ),
    (
        ["rdp", "encryption"],
        "Configure RDP to use Network Level Authentication (NLA) and "
        "TLS encryption. Disable legacy RDP encryption.",
    ),
    (
        ["default", "credential"],
        "Change all default credentials immediately. Implement a strong "
        "password policy and use a credential management solution.",
    ),
    (
        ["default", "login"],
        "Change default login credentials. Implement account lockout "
        "policies and multi-factor authentication where possible.",
    ),
    (
        ["open", "port"],
        "Review and close unnecessary open ports. Implement firewall "
        "rules to restrict access to required services only.",
    ),
    (
        ["http", "missing", "header"],
        "Configure recommended HTTP security headers (X-Frame-Options, "
        "X-Content-Type-Options, Content-Security-Policy, etc.).",
    ),
    (
        ["cors", "misconfiguration"],
        "Review and restrict CORS policies to allow only trusted origins. "
        "Avoid using wildcard (*) for Access-Control-Allow-Origin.",
    ),
    (
        ["xss"],
        "Implement input validation and output encoding. Deploy "
        "Content-Security-Policy headers to mitigate cross-site scripting.",
    ),
    (
        ["cve-"],
        "Apply the vendor-supplied security patch for this CVE. Monitor "
        "vendor advisories for updates and implement a regular patching cycle.",
    ),
    (
        ["exposed", "panel"],
        "Restrict access to administrative panels using network-level "
        "controls (firewall rules, VPN). Implement strong authentication.",
    ),
    (
        ["unencrypted"],
        "Enable encryption for this service. Migrate to the encrypted "
        "variant of the protocol (e.g. HTTP -> HTTPS, FTP -> SFTP).",
    ),
    (
        ["smb"],
        "Review SMB configuration. Disable SMBv1, enable SMB signing, "
        "and restrict SMB access to authorized networks only.",
    ),
    (
        ["dns"],
        "Review DNS server configuration. Restrict recursive queries "
        "and zone transfers to authorized hosts only.",
    ),
]


def generate_mitigation(vuln: Vulnerability) -> list[str]:
    """Generate AI-style mitigation recommendations for a vulnerability.

    Returns a list of at least 2 specific, actionable mitigation steps.
    """
    vuln_id_lower = vuln.vuln_id.lower()
    severity = vuln.severity.lower()
    
    # CVE-2023-34048 (VMware vCenter)
    if "cve-2023-34048" in vuln_id_lower:
        return [
            "Apply VMware vCenter Server security patch immediately to address the remote code execution vulnerability (refer to VMSA-2023-0023).",
            "Isolate vCenter Server management interfaces behind a firewall or VPN, restricting access to authorized administrators only.",
            "Monitor vCenter logs for suspicious authentication attempts or unusual API calls that could indicate exploitation attempts.",
        ]
    
    # LDAP Anonymous Login
    if "ldap-anonymous" in vuln_id_lower:
        return [
            "Disable anonymous LDAP bind on all domain controllers and directory servers by configuring authentication requirements.",
            "Implement LDAP signing and channel binding to prevent man-in-the-middle attacks and ensure authenticated access only.",
            "Review and restrict LDAP query permissions to prevent enumeration of sensitive directory information.",
        ]
    
    # SSH/Terrapin (CVE-2023-48795)
    if "cve-2023-48795" in vuln_id_lower or "terrapin" in vuln_id_lower:
        return [
            "Update SSH server and client software to versions that include the Terrapin vulnerability patch (OpenSSH 9.6+ or vendor equivalent).",
            "Configure SSH to use secure encryption algorithms and disable vulnerable cipher suites that are susceptible to prefix truncation attacks.",
        ]
    
    # Default credentials
    if "default" in vuln_id_lower and ("credential" in vuln_id_lower or "login" in vuln_id_lower or "password" in vuln_id_lower):
        return [
            "Change all default credentials immediately using strong, unique passwords that meet organizational complexity requirements (minimum 12 characters, mixed case, numbers, special characters).",
            "Implement multi-factor authentication (MFA) where supported to add an additional layer of security beyond passwords.",
            "Deploy a credential management solution to enforce password rotation policies and prevent reuse of default credentials.",
        ]
    
    # SSL/TLS issues
    if any(x in vuln_id_lower for x in ["ssl", "tls", "certificate"]):
        return [
            "Update SSL/TLS configuration to disable weak cipher suites and protocols (disable SSLv3, TLS 1.0, TLS 1.1), enforcing TLS 1.2+ with strong ciphers only.",
            "Implement proper certificate validation and ensure certificates are signed by trusted Certificate Authorities with valid expiration dates.",
            "Enable HTTP Strict Transport Security (HSTS) headers to force encrypted connections and prevent protocol downgrade attacks.",
        ]
    
    # HTTP misconfigurations
    if any(x in vuln_id_lower for x in ["http-", "apache", "nginx", "iis"]):
        return [
            "Update web server software to the latest stable version with all security patches applied.",
            "Configure security headers (X-Frame-Options, X-Content-Type-Options, Content-Security-Policy, X-XSS-Protection) to prevent common web attacks.",
            "Review and harden web server configuration by disabling unnecessary modules, restricting directory listings, and implementing proper access controls.",
        ]
    
    # Database exposures
    if any(x in vuln_id_lower for x in ["mysql", "postgres", "mongodb", "redis", "sql"]):
        return [
            "Configure database authentication to require strong passwords and disable anonymous or guest access entirely.",
            "Implement network-level restrictions using firewall rules to limit database access to authorized application servers only.",
            "Enable database audit logging to monitor for unauthorized access attempts and suspicious query patterns.",
        ]
    
    # SMB/Windows file sharing
    if any(x in vuln_id_lower for x in ["smb", "cifs", "netbios"]):
        return [
            "Disable SMBv1 protocol completely and enforce SMBv2/v3 with encryption to prevent exploitation of legacy protocol vulnerabilities.",
            "Enable SMB signing requirements to prevent man-in-the-middle attacks and session hijacking.",
            "Restrict SMB access using firewall rules to block ports 445, 139 from untrusted networks and implement network segmentation.",
        ]
    
    # RDP issues
    if "rdp" in vuln_id_lower:
        return [
            "Enable Network Level Authentication (NLA) to require authentication before establishing a full RDP session, reducing brute-force attack surface.",
            "Implement account lockout policies and deploy multi-factor authentication for all RDP access.",
            "Restrict RDP access using firewall rules or VPN, limiting connections to authorized source IP addresses only, and consider using an RDP gateway.",
        ]
    
    # DNS issues
    if "dns" in vuln_id_lower:
        return [
            "Configure DNS server access control lists (ACLs) to restrict zone transfers to authorized secondary DNS servers only.",
            "Disable DNS recursion for external queries on authoritative DNS servers to prevent DNS amplification attacks.",
            "Implement DNSSEC where possible to add cryptographic validation and prevent DNS spoofing attacks.",
        ]
    
    # SNMP exposures
    if "snmp" in vuln_id_lower:
        return [
            "Change default SNMP community strings from 'public'/'private' to strong, unique values, or preferably migrate to SNMPv3 with encryption and authentication.",
            "Restrict SNMP access using firewall rules and access control lists to allow only authorized network management systems.",
            "Configure SNMP to read-only mode unless write access is explicitly required, and regularly audit SNMP configuration for security compliance.",
        ]
    
    # Generic mitigation based on severity
    if severity == "critical":
        return [
            "Apply vendor-supplied security patches immediately as this critical vulnerability could lead to complete system compromise.",
            "Implement network segmentation to isolate affected systems and restrict access using firewall rules until remediation is complete.",
            "Monitor affected systems for indicators of compromise and review logs for any suspicious activity that could indicate exploitation.",
        ]
    elif severity == "high":
        return [
            "Apply security updates and patches according to organizational change management procedures, prioritizing remediation within 30 days.",
            "Implement compensating controls such as network access restrictions or enhanced monitoring until patching is complete.",
        ]
    elif severity == "medium":
        return [
            "Schedule remediation activities within the next 60-90 days according to normal patch management cycles.",
            "Review security configurations and implement hardening measures to reduce overall attack surface while planning for patching.",
        ]
    else:
        return [
            "Address this vulnerability as part of regular system hardening and maintenance activities.",
            "Document the finding and include in the next scheduled security review for long-term remediation planning.",
        ]


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def _find_slide_by_shape_name(prs, shape_name: str) -> int | None:
    """Find the 0-based slide index containing a shape with *shape_name*.

    Returns ``None`` if no slide has a matching shape.
    """
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.name == shape_name:
                return idx
    return None


def _remove_shape_by_name(slide, name: str) -> bool:
    """Remove the first shape matching *name*. Returns True if removed."""
    for shape in list(slide.shapes):
        if shape.name == name:
            shape._element.getparent().remove(shape._element)
            return True
    return False


def _add_bullet_para(
    tf,
    text: str,
    *,
    sz: int = 10,
    bold: bool = False,
    clr: RGBColor = CLR_DARK,
    level: int = 0,
    bullet_char: str | None = None,
    spc_before: int = 2,
    spc_after: int = 1,
):
    """Append a formatted paragraph with optional bullet to a text frame.

    Mirrors the bullet formatting from ``create_template._add_para()``.
    """
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = clr
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(spc_before)
    p.space_after = Pt(spc_after)

    if bullet_char:
        pPr = p._p.get_or_add_pPr()
        indent_emu = 228600  # 0.25" hanging indent
        margin_emu = indent_emu + (level * 457200)  # +0.5" per level
        pPr.set("marL", str(margin_emu))
        pPr.set("indent", str(-indent_emu))
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "+mj-lt")
        bc = etree.SubElement(pPr, qn("a:buChar"))
        bc.set("char", bullet_char)
    elif level > 0:
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(level * 457200))

    return p


def _build_vuln_textbox(slide, vulns_chunk: list[Vulnerability], shape_name: str):
    """Create a formatted textbox with vulnerability entries on *slide*.

    Each vulnerability gets:
    - Level-0 numbered header: ``1. [vuln-id] [protocol] [severity]``
    - Level-1 sub-bullet:  ``Targets: host1, host2, ...``
    - Level-1 sub-bullet:  ``Issue: AI-generated description``
    """
    box = slide.shapes.add_textbox(
        Inches(TEXTBOX_LEFT), Inches(TEXTBOX_TOP),
        Inches(TEXTBOX_WIDTH), Inches(TEXTBOX_HEIGHT),
    )
    box.name = shape_name
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Set top anchor
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "t")

    first = True
    for vuln in vulns_chunk:
        # Header with number: "1. [CVE-ID] [protocol] [severity]"
        severity_tag = f"[{vuln.severity.upper()}]"
        header = f"{vuln.index}. [{vuln.vuln_id}] [{vuln.protocol}] {severity_tag}"

        if first:
            # Use the existing first paragraph
            p = tf.paragraphs[0]
            p.text = header
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = CLR_DARK
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0)
            p.space_after = Pt(2)
            # No bullet for numbered entries - just bold text
            first = False
        else:
            # Add subsequent headers
            p_header = tf.add_paragraph()
            p_header.text = header
            p_header.font.size = Pt(10)
            p_header.font.bold = True
            p_header.font.color.rgb = CLR_DARK
            p_header.font.name = "Calibri"
            p_header.alignment = PP_ALIGN.LEFT
            p_header.space_before = Pt(6)
            p_header.space_after = Pt(2)

        # Targets sub-bullet
        # Ensure long target lists wrap properly
        targets_text = f"Targets: {vuln.target}"
        
        p_targets = tf.add_paragraph()
        p_targets.font.size = Pt(9)
        p_targets.font.color.rgb = CLR_DARK
        p_targets.font.name = "Calibri"
        p_targets.alignment = PP_ALIGN.LEFT
        p_targets.space_before = Pt(1)
        p_targets.space_after = Pt(1)
        p_targets.level = 1
        
        # Set bullet formatting
        pPr = p_targets._p.get_or_add_pPr()
        indent_emu = 228600
        margin_emu = indent_emu + (1 * 457200)  # level 1
        pPr.set("marL", str(margin_emu))
        pPr.set("indent", str(-indent_emu))
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "+mj-lt")
        bc = etree.SubElement(pPr, qn("a:buChar"))
        bc.set("char", "\u2013")
        
        # Add "Targets:" in bold
        run_targets_label = p_targets.add_run()
        run_targets_label.text = "Targets: "
        run_targets_label.font.bold = True
        run_targets_label.font.size = Pt(9)
        run_targets_label.font.color.rgb = CLR_DARK
        run_targets_label.font.name = "Calibri"
        
        # Add target list in regular
        run_targets_list = p_targets.add_run()
        run_targets_list.text = vuln.target
        run_targets_list.font.bold = False
        run_targets_list.font.size = Pt(9)
        run_targets_list.font.color.rgb = CLR_DARK
        run_targets_list.font.name = "Calibri"

        # Description sub-bullet (AI-generated)
        # Need to add "Issue:" in bold, then rest in regular
        desc = generate_vulnerability_description(vuln)
        p_desc = tf.add_paragraph()
        p_desc.font.size = Pt(9)
        p_desc.font.color.rgb = CLR_DARK
        p_desc.font.name = "Calibri"
        p_desc.alignment = PP_ALIGN.LEFT
        p_desc.space_before = Pt(1)
        p_desc.space_after = Pt(3)
        p_desc.level = 1
        
        # Set bullet formatting
        pPr = p_desc._p.get_or_add_pPr()
        indent_emu = 228600
        margin_emu = indent_emu + (1 * 457200)  # level 1
        pPr.set("marL", str(margin_emu))
        pPr.set("indent", str(-indent_emu))
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "+mj-lt")
        bc = etree.SubElement(pPr, qn("a:buChar"))
        bc.set("char", "\u2013")
        
        # Add "Issue:" in bold
        run_issue = p_desc.add_run()
        run_issue.text = "Issue: "
        run_issue.font.bold = True
        run_issue.font.size = Pt(9)
        run_issue.font.color.rgb = CLR_DARK
        run_issue.font.name = "Calibri"
        
        # Add description in regular
        run_desc = p_desc.add_run()
        run_desc.text = desc
        run_desc.font.bold = False
        run_desc.font.size = Pt(9)
        run_desc.font.color.rgb = CLR_DARK
        run_desc.font.name = "Calibri"


def _build_mitigation_textbox(slide, vulns_chunk: list[Vulnerability], shape_name: str):
    """Create a formatted textbox with mitigation entries on *slide*.

    Each vulnerability gets:
    - Level-0 numbered header: ``1. [severity]`` (matching vuln number, no CVE ID)
    - Level-1 sub-bullets: multiple mitigation steps (at least 2 per vulnerability)
    """
    box = slide.shapes.add_textbox(
        Inches(TEXTBOX_LEFT), Inches(TEXTBOX_TOP),
        Inches(TEXTBOX_WIDTH), Inches(TEXTBOX_HEIGHT),
    )
    box.name = shape_name
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "t")

    first = True
    for vuln in vulns_chunk:
        # Header with number and severity only (no CVE ID)
        severity_tag = f"[{vuln.severity.upper()}]"
        header = f"{vuln.index}. {severity_tag}"
        mitigation_steps = generate_mitigation(vuln)  # Returns list of steps

        if first:
            p = tf.paragraphs[0]
            p.text = header
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = CLR_DARK
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0)
            p.space_after = Pt(2)
            first = False
        else:
            # Add subsequent headers
            p_header = tf.add_paragraph()
            p_header.text = header
            p_header.font.size = Pt(10)
            p_header.font.bold = True
            p_header.font.color.rgb = CLR_DARK
            p_header.font.name = "Calibri"
            p_header.alignment = PP_ALIGN.LEFT
            p_header.space_before = Pt(6)
            p_header.space_after = Pt(2)

        # Add each mitigation step with hardcoded bullet
        for step in mitigation_steps:
            # Add bullet directly to the text
            bullet_text = f"\u2022  {step}"
            p_step = tf.add_paragraph()
            p_step.text = bullet_text
            p_step.font.size = Pt(9)
            p_step.font.bold = False
            p_step.font.color.rgb = CLR_DARK
            p_step.font.name = "Calibri"
            p_step.alignment = PP_ALIGN.LEFT
            p_step.space_before = Pt(1)
            p_step.space_after = Pt(2)


def _build_no_vulns_textbox(slide, shape_name: str, message: str = "No vulnerabilities detected."):
    """Place a simple 'no vulnerabilities' message on the slide."""
    box = slide.shapes.add_textbox(
        Inches(TEXTBOX_LEFT), Inches(TEXTBOX_TOP),
        Inches(TEXTBOX_WIDTH), Inches(TEXTBOX_HEIGHT),
    )
    box.name = shape_name
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "t")

    p = tf.paragraphs[0]
    p.text = message
    p.font.size = Pt(11)
    p.font.bold = False
    p.font.color.rgb = CLR_DARK
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT


# ---------------------------------------------------------------------------
# Pagination & population
# ---------------------------------------------------------------------------
def populate_vulnerability_slides(
    prs,
    vulns: list[Vulnerability],
    vuln_shape_name: str,
    mit_shape_name: str,
    vuln_intro_shape_name: str | None = None,
    mit_intro_shape_name: str | None = None,
    sensor_suffix: str = "",
) -> int:
    """Populate vulnerability and mitigation slides with parsed data.

    Handles pagination by duplicating slides when content overflows.

    Parameters
    ----------
    prs : Presentation
        The open PPTX presentation object.
    vulns : list[Vulnerability]
        Parsed vulnerability entries (may be empty).
    vuln_shape_name : str
        Name of the placeholder shape on the vulnerability slide
        (e.g. ``"internal_vulns_content"`` or ``"internal_vulns_content_vaprd"``).
    mit_shape_name : str
        Name of the placeholder shape on the mitigation slide.
    vuln_intro_shape_name : str | None
        Name of the intro text shape on the vulnerability slide (removed
        on continuation slides).
    mit_intro_shape_name : str | None
        Name of the intro text shape on the mitigation slide.
    sensor_suffix : str
        Sensor tag for logging (e.g. ``"VAPRD"``).

    Returns
    -------
    int
        Total number of extra slides inserted (continuation slides).
    """
    tag = f" [{sensor_suffix}]" if sensor_suffix else ""

    # Locate slides by shape name
    vuln_slide_idx = _find_slide_by_shape_name(prs, vuln_shape_name)
    mit_slide_idx = _find_slide_by_shape_name(prs, mit_shape_name)

    if vuln_slide_idx is None:
        logger.warning("Vuln slide not found for shape '%s'%s", vuln_shape_name, tag)
        return 0
    if mit_slide_idx is None:
        logger.warning("Mitigation slide not found for shape '%s'%s", mit_shape_name, tag)
        return 0

    # Handle empty / no vulnerabilities
    if not vulns:
        slide = prs.slides[vuln_slide_idx]
        _remove_shape_by_name(slide, vuln_shape_name)
        _build_no_vulns_textbox(slide, vuln_shape_name)

        mit_slide = prs.slides[mit_slide_idx]
        _remove_shape_by_name(mit_slide, mit_shape_name)
        _build_no_vulns_textbox(
            mit_slide, mit_shape_name,
            "No mitigation actions required — no vulnerabilities were detected.",
        )
        logger.info("No vulnerabilities to populate%s", tag)
        return 0

    # Calculate pagination
    n_pages = max(1, math.ceil(len(vulns) / MAX_PER_SLIDE))
    vuln_slides_inserted = 0
    mit_slides_inserted = 0

    # --- Phase 1: Duplicate vulnerability slides ---
    # Duplicate BEFORE populating so each copy has the placeholder shape
    for page in range(1, n_pages):
        target_idx = vuln_slide_idx + page
        duplicate_slide(prs, vuln_slide_idx, target_idx)
        vuln_slides_inserted += 1

    # Recalculate mitigation slide position (shifted by inserted vuln slides)
    mit_slide_idx = _find_slide_by_shape_name(prs, mit_shape_name)
    if mit_slide_idx is None:
        logger.error("Mitigation slide lost after vuln duplication%s", tag)
        return vuln_slides_inserted

    # --- Phase 2: Populate each vulnerability slide ---
    for page in range(n_pages):
        slide_idx = vuln_slide_idx + page
        slide = prs.slides[slide_idx]
        chunk = vulns[page * MAX_PER_SLIDE : (page + 1) * MAX_PER_SLIDE]

        # Remove placeholder
        _remove_shape_by_name(slide, vuln_shape_name)

        # Build content
        _build_vuln_textbox(slide, chunk, vuln_shape_name)

        # Continuation slides: update title (keep intro text)
        if page > 0:
            # For multi-sensor, preserve sensor name in continuation title
            if sensor_suffix and sensor_suffix != "aggregate":
                cont_title = f"Internal Network Vulnerabilities - {sensor_suffix} (cont.)"
            else:
                cont_title = "Internal Network Vulnerabilities (cont.)"
            update_slide_title(slide, cont_title)

    logger.info(
        "Populated %d vuln slide(s) with %d vulnerabilities%s",
        n_pages, len(vulns), tag,
    )

    # --- Phase 3: Duplicate mitigation slides ---
    for page in range(1, n_pages):
        target_idx = mit_slide_idx + page
        duplicate_slide(prs, mit_slide_idx, target_idx)
        mit_slides_inserted += 1

    # --- Phase 4: Populate each mitigation slide ---
    for page in range(n_pages):
        slide_idx = mit_slide_idx + page
        slide = prs.slides[slide_idx]
        chunk = vulns[page * MAX_PER_SLIDE : (page + 1) * MAX_PER_SLIDE]

        _remove_shape_by_name(slide, mit_shape_name)
        _build_mitigation_textbox(slide, chunk, mit_shape_name)

        if page > 0:
            # For multi-sensor, preserve sensor name in continuation title
            if sensor_suffix and sensor_suffix != "aggregate":
                cont_title = f"Internal Network Mitigation - {sensor_suffix} (cont.)"
            else:
                cont_title = "Internal Network Mitigation (cont.)"
            update_slide_title(slide, cont_title)

    logger.info(
        "Populated %d mitigation slide(s)%s",
        n_pages, tag,
    )

    total_inserted = vuln_slides_inserted + mit_slides_inserted
    return total_inserted


def update_all_page_numbers(prs) -> None:
    """Update every slide's ``page_number`` shape to match its 1-based position."""
    for idx, slide in enumerate(prs.slides):
        page_str = str(idx + 1)
        for shape in slide.shapes:
            if shape.name == "page_number" and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = page_str
                break
