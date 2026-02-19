"""
Patch History Management

Tracks historical patch counts over time and generates trend charts
that automatically update with the current reporting period's data.
"""

from __future__ import annotations

import csv
import json
import logging
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

logger = logging.getLogger("shm.patch_history")


class PatchHistoryManager:
    """Manages historical patch count data for trend chart generation."""
    
    def __init__(self, storage_path: str | Path = "assets/patch_history"):
        """Initialize the patch history manager.
        
        Parameters
        ----------
        storage_path : str | Path
            Directory where patch history JSON files are stored
        """
        self.storage_path = Path(storage_path)
        self.storage_path.mkdir(parents=True, exist_ok=True)
    
    def _get_history_file(self, client_slug: str) -> Path:
        """Get the path to a client's history file."""
        return self.storage_path / f"{client_slug}_patch_history.json"
    
    def load_history(self, client_slug: str) -> list[dict]:
        """Load patch history for a client.
        
        Returns
        -------
        list[dict]
            List of history entries, each with: date, microsoft_count, software_count
        """
        history_file = self._get_history_file(client_slug)
        
        if not history_file.exists():
            logger.info("No history file found for '%s', starting fresh", client_slug)
            return []
        
        try:
            with open(history_file, 'r', encoding='utf-8') as f:
                history = json.load(f)
            logger.info("Loaded %d history entries for '%s'", len(history), client_slug)
            return history
        except Exception as exc:
            logger.error("Failed to load history for '%s': %s", client_slug, exc)
            return []
    
    def save_history(self, client_slug: str, history: list[dict]) -> None:
        """Save patch history for a client.
        
        Parameters
        ----------
        client_slug : str
            Client identifier
        history : list[dict]
            List of history entries
        """
        history_file = self._get_history_file(client_slug)
        
        try:
            with open(history_file, 'w', encoding='utf-8') as f:
                json.dump(history, f, indent=2)
            logger.info("Saved %d history entries for '%s'", len(history), client_slug)
        except Exception as exc:
            logger.error("Failed to save history for '%s': %s", client_slug, exc)
    
    def add_entry(
        self,
        client_slug: str,
        end_date: str,
        microsoft_count: int,
        software_count: int,
        max_history: int = 6
    ) -> None:
        """Add or update a history entry for the current report period.
        
        Parameters
        ----------
        client_slug : str
            Client identifier
        end_date : str
            Report period end date (e.g., "2026-01-05")
        microsoft_count : int
            Count of missing Microsoft KB patches
        software_count : int
            Count of missing software packages
        max_history : int
            Maximum number of history entries to keep (default: 6)
        """
        history = self.load_history(client_slug)
        
        # Check if entry for this date already exists
        existing_idx = None
        for i, entry in enumerate(history):
            if entry.get('date') == end_date:
                existing_idx = i
                break
        
        new_entry = {
            'date': end_date,
            'microsoft_count': microsoft_count,
            'software_count': software_count,
            'updated_at': datetime.now().isoformat(),
        }
        
        if existing_idx is not None:
            # Update existing entry
            history[existing_idx] = new_entry
            logger.info("Updated history entry for '%s' on %s", client_slug, end_date)
        else:
            # Add new entry
            history.append(new_entry)
            logger.info("Added history entry for '%s' on %s", client_slug, end_date)
        
        # Sort by date
        history.sort(key=lambda x: x['date'])
        
        # Keep only the most recent entries
        if len(history) > max_history:
            history = history[-max_history:]
            logger.info("Trimmed history to %d most recent entries", max_history)
        
        self.save_history(client_slug, history)
    
    def get_trend_data(
        self,
        client_slug: str,
        end_date: str,
        microsoft_count: int,
        software_count: int,
        num_points: int = 3
    ) -> tuple[list[str], list[int], list[int]]:
        """Get trend data for chart generation.

        Returns the most recent *num_points* history entries ending with the
        current report month (``end_date``).  Entries are taken from whatever
        months exist in history (no forced consecutive-month gaps).

        Parameters
        ----------
        client_slug : str
            Client identifier
        end_date : str
            Current report period end date (e.g., "2026-01-05")
        microsoft_count : int
            Current count of missing Microsoft KB patches
        software_count : int
            Current count of missing software packages
        num_points : int
            Number of data points to include in the trend (default: 3)

        Returns
        -------
        tuple[list[str], list[int], list[int]]
            (dates, microsoft_counts, software_counts) for chart generation
            Dates are formatted as month names (e.g., "January")
        """
        # Add/update current entry
        self.add_entry(client_slug, end_date, microsoft_count, software_count)

        # Load updated history
        history = self.load_history(client_slug)

        if not history:
            logger.warning("No history available for '%s'", client_slug)
            return [], [], []

        current_dt = datetime.strptime(end_date, '%Y-%m-%d')
        current_ym = (current_dt.year, current_dt.month)

        # Freshness cutoff: drop entries older than num_points+3 months
        cutoff_m = current_dt.month - (num_points + 3)
        cutoff_y = current_dt.year
        while cutoff_m <= 0:
            cutoff_m += 12
            cutoff_y -= 1
        cutoff_ym = (cutoff_y, cutoff_m)

        # Deduplicate by (year, month) keeping the most recent entry per
        # month, and discard entries that are too old or in the future.
        month_best: dict[tuple[int, int], dict] = {}
        for entry in history:
            try:
                dt = datetime.strptime(entry['date'], '%Y-%m-%d')
                ym = (dt.year, dt.month)
            except Exception:
                continue
            if ym < cutoff_ym or ym > current_ym:
                continue
            if ym not in month_best or entry['date'] > month_best[ym]['date']:
                month_best[ym] = entry

        # Sort by month and take the most recent num_points entries
        deduped = sorted(month_best.values(), key=lambda e: e['date'])
        recent = deduped[-num_points:]

        # Persist the cleaned-up history so stale entries don't return
        self.save_history(client_slug, deduped)

        # --- Guarantee exactly num_points months ----------------------
        # If we have fewer entries than required, pad the beginning with
        # the consecutive months before the earliest entry so the chart
        # always has the right number of data points.  Padded months
        # carry the same counts as the earliest real entry so the line
        # stays flat rather than dropping to zero.
        if len(recent) < num_points:
            earliest = recent[0] if recent else {
                'date': end_date,
                'microsoft_count': microsoft_count,
                'software_count': software_count,
            }
            earliest_dt = datetime.strptime(earliest['date'], '%Y-%m-%d')
            pad_entries: list[dict] = []
            for i in range(num_points - len(recent), 0, -1):
                pm = earliest_dt.month - i
                py = earliest_dt.year
                while pm <= 0:
                    pm += 12
                    py -= 1
                pad_entries.append({
                    'date': f"{py:04d}-{pm:02d}-01",
                    'microsoft_count': earliest['microsoft_count'],
                    'software_count': earliest['software_count'],
                })
            recent = pad_entries + recent

        # Format dates for display (month name only)
        dates = []
        for entry in recent:
            try:
                date_obj = datetime.strptime(entry['date'], '%Y-%m-%d')
                dates.append(date_obj.strftime('%B'))
            except:
                dates.append(entry['date'])

        microsoft_counts = [entry['microsoft_count'] for entry in recent]
        software_counts = [entry['software_count'] for entry in recent]

        logger.info(
            "Generated trend data for '%s': %d points, latest = MS:%d SW:%d",
            client_slug, len(dates), microsoft_counts[-1], software_counts[-1]
        )

        return dates, microsoft_counts, software_counts


def extract_patch_counts_from_pptx(pptx_path: str | Path) -> tuple[int, int] | None:
    """Extract Microsoft KB and software patch counts from a PPTX report.

    First tries the named shape ``required_patches_body``.  If that is not
    found (e.g. the file was round-tripped through Google Slides, which
    renames shapes), falls back to scanning *every* text frame for the
    characteristic patch-count sentences.

    Returns
    -------
    tuple[int, int] | None
        ``(ms_count, sw_count)`` or ``None`` if the counts cannot be parsed.
    """
    from pptx import Presentation as PptxPresentation

    _MS_RE = re.compile(r"currently\s+(\d+)\s+Microsoft\s+KB", re.IGNORECASE)
    _SW_RE = re.compile(r"(\d+)\s+software\s+packages", re.IGNORECASE)

    try:
        prs = PptxPresentation(str(pptx_path))
    except Exception as exc:
        logger.warning("Could not open PPTX '%s': %s", pptx_path, exc)
        return None

    def _try_extract(text: str) -> tuple[int, int] | None:
        ms_match = _MS_RE.search(text)
        sw_match = _SW_RE.search(text)
        ms = int(ms_match.group(1)) if ms_match else None
        sw = int(sw_match.group(1)) if sw_match else None
        if ms is not None or sw is not None:
            return (ms or 0, sw or 0)
        return None

    # Pass 1: look for the named shape (our own template)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == "required_patches_body" and shape.has_text_frame:
                result = _try_extract(shape.text_frame.text)
                if result:
                    logger.info("Extracted patch counts (named shape) from '%s': MS=%s, SW=%s",
                                pptx_path, *result)
                    return result

    # Pass 2: scan all text frames (handles Google Slides renamed shapes)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            result = _try_extract(shape.text_frame.text)
            if result:
                logger.info("Extracted patch counts (text scan) from '%s': MS=%s, SW=%s",
                            pptx_path, *result)
                return result

    logger.warning("No patch count text found in '%s'", pptx_path)
    return None


def estimate_report_date_from_pptx(
    pptx_path: str | Path, drive_modified_time: str | None = None
) -> str | None:
    """Estimate the report date from a PPTX file.

    Looks for text matching ``"… to <day> <Month> <year>"`` — first in a
    shape named ``reporting_period`` on the cover slide, then by scanning
    all text frames on the first few slides (handles Google Slides renamed
    shapes).

    Falls back to the Drive ``modifiedTime`` if parsing fails.

    Returns
    -------
    str | None
        Date string in ``YYYY-MM-DD`` format, or ``None``.
    """
    from pptx import Presentation as PptxPresentation

    _DATE_RE = re.compile(r"to\s+(\d{1,2}\s+\w+\s+\d{4})", re.IGNORECASE)

    def _try_parse(text: str) -> str | None:
        match = _DATE_RE.search(text)
        if match:
            try:
                return datetime.strptime(match.group(1), "%d %B %Y").strftime("%Y-%m-%d")
            except ValueError:
                pass
        return None

    try:
        prs = PptxPresentation(str(pptx_path))
        if prs.slides:
            cover = prs.slides[0]

            # Pass 1: named shape
            for shape in cover.shapes:
                if shape.name == "reporting_period" and shape.has_text_frame:
                    result = _try_parse(shape.text_frame.text)
                    if result:
                        return result

            # Pass 2: any text frame on the cover slide
            for shape in cover.shapes:
                if shape.has_text_frame:
                    result = _try_parse(shape.text_frame.text)
                    if result:
                        logger.info("Parsed report date (text scan) from cover slide: %s", result)
                        return result
    except Exception as exc:
        logger.warning("Could not parse report date from '%s': %s", pptx_path, exc)

    # Fallback to Drive modifiedTime
    if drive_modified_time:
        try:
            dt = datetime.fromisoformat(drive_modified_time.replace("Z", "+00:00"))
            return dt.strftime("%Y-%m-%d")
        except Exception as exc:
            logger.warning("Could not parse Drive modifiedTime '%s': %s", drive_modified_time, exc)

    return None


def seed_history_from_drive(
    client_slug: str,
    drive_agent,
    manager: PatchHistoryManager | None = None,
    storage_path: str | Path = "assets/patch_history",
    count: int = 4,
    current_end_date: str | None = None,
) -> int:
    """Seed local patch history by extracting data from previous PPTX reports
    found in the client's top-level Drive folder (Archive is excluded).

    PPTX files whose report date falls in the same month as
    *current_end_date* are skipped because the current report already
    provides that month's data.

    Parameters
    ----------
    client_slug : str
        Client identifier.
    drive_agent
        A ``DriveAgent`` instance.
    manager : PatchHistoryManager | None
        Existing manager instance (created if ``None``).
    storage_path : str | Path
        Directory for patch history storage.
    count : int
        Number of PPTX files to download (default: 4, extra buffer so
        that after skipping the current month we still have enough).
    current_end_date : str | None
        The current report's end date (``YYYY-MM-DD``).  PPTX reports
        from the same month are skipped.

    Returns
    -------
    int
        Number of history entries added.
    """
    if manager is None:
        manager = PatchHistoryManager(storage_path)

    # Determine current month so we can skip PPTX files from it
    current_ym: tuple[int, int] | None = None
    if current_end_date:
        try:
            _dt = datetime.strptime(current_end_date, "%Y-%m-%d")
            current_ym = (_dt.year, _dt.month)
        except Exception:
            pass

    client_folder = drive_agent.find_client_folder(client_slug)
    client_folder_id = client_folder["id"]

    pptx_files = drive_agent.find_client_pptx_files(
        client_folder_id=client_folder_id,
        client_slug=client_slug,
        count=count,
    )
    if not pptx_files:
        logger.info("No previous PPTX files found on Drive for '%s'", client_slug)
        return 0

    added = 0
    for entry in pptx_files:
        local_path = entry["local_path"]
        try:
            counts = extract_patch_counts_from_pptx(local_path)
            if counts is None:
                continue
            ms_count, sw_count = counts

            report_date = estimate_report_date_from_pptx(
                local_path, drive_modified_time=entry.get("modified_time"),
            )
            if report_date is None:
                logger.warning("Could not determine date for '%s', skipping", entry["name"])
                continue

            # Skip if this PPTX is from the same month as the current report
            if current_ym:
                try:
                    rd = datetime.strptime(report_date, "%Y-%m-%d")
                    if (rd.year, rd.month) == current_ym:
                        logger.info(
                            "Skipping '%s' (same month as current report)", entry["name"],
                        )
                        continue
                except Exception:
                    pass

            manager.add_entry(client_slug, report_date, ms_count, sw_count)
            added += 1
            logger.info(
                "Seeded history for '%s': date=%s MS=%d SW=%d (from '%s')",
                client_slug, report_date, ms_count, sw_count, entry["name"],
            )
        except Exception as exc:
            logger.warning("Failed to extract data from '%s': %s", entry["name"], exc)
        finally:
            # Clean up temp file
            try:
                Path(local_path).unlink(missing_ok=True)
            except Exception:
                pass

    logger.info("Seeded %d history entries for '%s' from Drive", added, client_slug)
    return added


def generate_patch_trend_chart_for_report(
    client_slug: str,
    end_date: str,
    microsoft_count: int,
    software_count: int,
    output_path: str | Path,
    *,
    num_points: int = 3,
    storage_path: str | Path = "assets/patch_history",
    drive_agent=None,
) -> Path:
    """Generate a patch trend chart for a report, automatically managing history.
    
    This is the main function to call from the API. It:
    1. Loads historical patch data for the client
    2. Adds/updates the current report period's counts
    3. Generates a trend chart with the most recent data points
    4. Saves the updated history for future reports
    
    Parameters
    ----------
    client_slug : str
        Client identifier (e.g., "elephant", "wiegmann")
    end_date : str
        Report period end date (format: "YYYY-MM-DD", e.g., "2026-01-05")
    microsoft_count : int
        Current count of missing Microsoft KB patches
    software_count : int
        Current count of missing software packages
    output_path : str | Path
        Where to save the generated chart image
    num_points : int
        Number of historical data points to show (default: 3)
    storage_path : str | Path
        Directory for patch history storage (default: "assets/patch_history")
    
    Returns
    -------
    Path
        Path to the generated chart image
    
    Examples
    --------
    >>> chart_path = generate_patch_trend_chart_for_report(
    ...     client_slug="elephant",
    ...     end_date="2026-01-05",
    ...     microsoft_count=4,
    ...     software_count=2,
    ...     output_path="assets/ndr_cache/elephant/patch_trend.png"
    ... )
    """
    from core.patch_chart import generate_patch_trend_chart

    manager = PatchHistoryManager(storage_path)

    # Seed history from previous PPTX reports in the client folder
    if drive_agent is not None:
        try:
            seed_history_from_drive(
                client_slug=client_slug,
                drive_agent=drive_agent,
                manager=manager,
                storage_path=storage_path,
                current_end_date=end_date,
            )
        except Exception as exc:
            logger.warning("Drive history seeding failed (non-fatal): %s", exc)

    # Get trend data (automatically updates history)
    dates, ms_counts, sw_counts = manager.get_trend_data(
        client_slug=client_slug,
        end_date=end_date,
        microsoft_count=microsoft_count,
        software_count=software_count,
        num_points=num_points,
    )
    
    if not dates:
        logger.warning(
            "No trend data available for '%s', cannot generate chart",
            client_slug
        )
        return None
    
    # Generate the chart
    return generate_patch_trend_chart(
        dates=dates,
        microsoft_counts=ms_counts,
        software_counts=sw_counts,
        output_path=output_path,
    )
