"""
Slide manipulation utilities for paginated tables.

Provides helpers to insert new slides at arbitrary positions and decorate
them with the standard Mantix4 header/footer chrome so continuation pages
look identical to the originals.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

logger = logging.getLogger("shm.slide_utils")

# ---------------------------------------------------------------------------
# Layout constants (duplicated from create_template.py)
# ---------------------------------------------------------------------------
TITLE_BAR_Y = 0.13
TITLE_BAR_H = 0.58
LINE_Y = 0.94
FOOTER_Y = 5.08
FOOTER_H = 0.55
MARGIN_L = 0.38
MARGIN_R = 0.49

CLR_BLACK = RGBColor(0x00, 0x00, 0x00)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_GREEN = RGBColor(0x38, 0x76, 0x1D)
CLR_DARK = RGBColor(0x1B, 0x21, 0x2C)

LOGO_PATH = Path("assets/mantix4_logo.png")


# ---------------------------------------------------------------------------
# Slide reordering
# ---------------------------------------------------------------------------
def delete_slide(prs, slide_index: int) -> None:
    """Delete a slide at *slide_index* (0-based) from *prs*.

    Removes the relationship and the ``sldId`` entry so the slide
    is fully excised from the presentation.
    """
    sldIdLst = prs.slides._sldIdLst
    sldId_items = list(sldIdLst)
    if slide_index < 0 or slide_index >= len(sldId_items):
        raise IndexError(f"slide_index {slide_index} out of range (0..{len(sldId_items)-1})")
    rId = sldId_items[slide_index].get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId_items[slide_index])


def move_slide(prs, from_index: int, to_index: int) -> None:
    """Move a slide from *from_index* to *to_index* within *prs*.

    Manipulates the ``sldIdLst`` XML element directly.
    """
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    if from_index < 0 or from_index >= len(slides):
        raise IndexError(f"from_index {from_index} out of range (0..{len(slides)-1})")
    if to_index < 0 or to_index > len(slides):
        raise IndexError(f"to_index {to_index} out of range (0..{len(slides)})")

    elem = slides[from_index]
    sldIdLst.remove(elem)
    # After removal, re-read the list so the insert target is correct.
    slides = list(sldIdLst)
    if to_index >= len(slides):
        sldIdLst.append(elem)
    else:
        sldIdLst.insert(to_index, elem)


def insert_slide_at(prs, position: int):
    """Add a blank slide and move it to *position*.

    Returns the new ``Slide`` object.
    """
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)
    # The new slide is appended at the end; move it to the desired position.
    from_index = len(prs.slides) - 1
    if from_index != position:
        move_slide(prs, from_index, position)
    return prs.slides[position]


# ---------------------------------------------------------------------------
# Decoration helpers (mirror create_template._add_header / _add_footer)
# ---------------------------------------------------------------------------
def _add_rect(slide, l, t, w, h, *, fill=None, name=None):
    """Add a filled rectangle shape."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    if name:
        s.name = name
    return s


def _add_textbox(slide, l, t, w, h, text, *, sz=14, bold=False,
                 clr=CLR_BLACK, align=PP_ALIGN.LEFT, name=None,
                 anchor=MSO_ANCHOR.TOP):
    """Add a text box."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        anchor_map = {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }
        bodyPr.set("anchor", anchor_map.get(anchor, "t"))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = clr
    p.font.name = "Calibri"
    p.alignment = align
    return box


def decorate_slide(slide, title: str, *, wide: bool = False,
                   page_num: str = "") -> None:
    """Apply full Mantix4 chrome (background, title bar, accent line,
    footer, logo, page number) to a blank *slide*.

    Replicates ``_add_header()`` + ``_add_footer()`` from
    ``create_template.py``.
    """
    # 1. White background rect
    _add_rect(slide, 0, 0, 10, 5.625, fill=CLR_WHITE, name="bg_rect")

    # 2. Title bar
    bw = 7.5 if wide else 5.5
    bl = (10 - bw) / 2
    _add_rect(slide, bl, TITLE_BAR_Y, bw, TITLE_BAR_H,
              fill=CLR_BLACK, name="title_bar")

    # 3. Title text
    _add_textbox(slide, bl, TITLE_BAR_Y + 0.05, bw, TITLE_BAR_H - 0.05,
                 title, sz=20, bold=True, clr=CLR_WHITE,
                 align=PP_ALIGN.CENTER, name="slide_title",
                 anchor=MSO_ANCHOR.MIDDLE)

    # 4. Green accent line
    thickness = Pt(3)
    line_left = MARGIN_L
    line_right = 10 - MARGIN_R
    r = _add_rect(slide, line_left, LINE_Y, line_right - line_left, 0.01,
                  fill=CLR_GREEN, name="accent_line")
    r.height = thickness

    # 5. Dark footer bar
    _add_rect(slide, 0, FOOTER_Y, 10, FOOTER_H,
              fill=CLR_DARK, name="footer_bar")

    # 6. Logo
    if LOGO_PATH.exists():
        logo_w = Inches(1.3)
        logo_h = Inches(0.34)
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(0.38), Inches(FOOTER_Y + 0.10),
            logo_w, logo_h,
        )

    # 7. Page number
    _add_textbox(slide, 9.0, FOOTER_Y + 0.05, 0.7, 0.4,
                 page_num, sz=10, clr=CLR_WHITE,
                 align=PP_ALIGN.RIGHT, name="page_number")


# ---------------------------------------------------------------------------
# Multi-sensor slide duplication
# ---------------------------------------------------------------------------
def duplicate_slide(prs, source_index: int, target_index: int = None):
    """Duplicate a slide at source_index and insert at target_index.

    If target_index is None, appends to the end.
    Returns the new slide object.
    """
    import copy

    source_slide = prs.slides[source_index]
    blank_layout = prs.slide_layouts[6]

    # Create a new blank slide
    new_slide = prs.slides.add_slide(blank_layout)

    # Deep copy all shapes from source to new slide,
    # re-registering image relationships so pictures (logo etc.) work.
    for shape in source_slide.shapes:
        el_copy = copy.deepcopy(shape.element)

        # Fix embedded image references (a:blip r:embed="rIdX")
        for blip in el_copy.findall('.//' + qn('a:blip')):
            embed_rId = blip.get(qn('r:embed'))
            if embed_rId:
                try:
                    rel = source_slide.part.rels[embed_rId]
                    new_rId = new_slide.part.relate_to(
                        rel.target_part, rel.reltype,
                    )
                    blip.set(qn('r:embed'), new_rId)
                except KeyError:
                    pass

        new_slide.shapes._spTree.insert_element_before(el_copy, 'p:extLst')

    # Move to target position if specified
    if target_index is not None:
        from_index = len(prs.slides) - 1
        if from_index != target_index:
            move_slide(prs, from_index, target_index)
        return prs.slides[target_index]

    return new_slide


def update_slide_title(slide, new_title: str) -> bool:
    """Update the text in the slide_title shape.
    
    Returns True if successful, False if shape not found.
    """
    for shape in slide.shapes:
        if shape.name == "slide_title" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = new_title
            return True
    return False


def rename_shape_with_suffix(slide, old_name: str, suffix: str) -> bool:
    """Rename a shape by appending a suffix.
    
    Example: rename_shape_with_suffix(slide, "ndr_outbound_data_1", "_vaprd")
    -> shape renamed to "ndr_outbound_data_1_vaprd"
    
    Returns True if shape was found and renamed.
    """
    for shape in slide.shapes:
        if shape.name == old_name:
            shape.name = f"{old_name}{suffix}"
            return True
    return False


def duplicate_ndr_slides_for_sensors(
    prs,
    sensors: list[str],
    base_slide_indices: list[int] = None,
) -> dict[str, list[int]]:
    """Duplicate NDR slide templates for each sensor.
    
    Args:
        prs: PowerPoint presentation object
        sensors: List of sensor IDs (e.g., ['VAPRD', 'VAHQ', 'GAPRD'])
        base_slide_indices: Template slides to duplicate (default: [5,6,7,8,9,10,11,12])
    
    Returns:
        Mapping of sensor_id -> [slide_indices]
        
    Example:
        Input: sensors=['VAPRD', 'VAHQ'], base=[5,6,7,8,9,10,11,12]
        Output: {
            'VAPRD': [5,6,7,8,9,10,11,12],      # Original slides
            'VAHQ': [13,14,15,16,17,18,19,20]   # Duplicated slides
        }
    """
    if base_slide_indices is None:
        base_slide_indices = [5, 6, 7, 8, 9, 10, 11, 12]
    
    if not sensors:
        return {}
    
    # Single sensor or DEFAULT: use original slides
    if len(sensors) == 1 or (len(sensors) == 1 and sensors[0] == "DEFAULT"):
        sensor_id = sensors[0]
        logger.info("Single sensor '%s': using original slides %s", sensor_id, base_slide_indices)
        return {sensor_id: base_slide_indices}
    
    sensor_slide_map: dict[str, list[int]] = {}
    
    # First sensor uses original slides
    first_sensor = sensors[0]
    sensor_slide_map[first_sensor] = list(base_slide_indices)
    
    # Update shape names and titles for first sensor
    for slide_idx in base_slide_indices:
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        
        # Update slide title to include sensor name
        for shape in slide.shapes:
            if shape.name == "slide_title" and shape.has_text_frame:
                current_title = shape.text_frame.text
                if first_sensor not in current_title:
                    new_title = f"{current_title} ({first_sensor})"
                    update_slide_title(slide, new_title)
                break
        
        # Rename placeholder shapes with sensor suffix
        placeholder_names = [
            "ndr_outbound_data_1", "ndr_outbound_data_2", "ndr_top_ip",
            "ndr_top_urls", "ndr_ext_dest", "ndr_country",
            "ndr_beaconing", "ndr_sensitive_data"
        ]
        for placeholder_name in placeholder_names:
            rename_shape_with_suffix(slide, placeholder_name, f"_{first_sensor.lower()}")
    
    logger.info("Sensor '%s': assigned to original slides %s", first_sensor, base_slide_indices)
    
    # Duplicate slides for remaining sensors
    for sensor_idx, sensor_id in enumerate(sensors[1:], start=1):
        new_slide_indices = []
        
        # Calculate insertion point (after previous sensor's slides)
        insert_position = base_slide_indices[-1] + 1 + (sensor_idx - 1) * len(base_slide_indices)
        
        for offset, source_idx in enumerate(base_slide_indices):
            if source_idx >= len(prs.slides):
                logger.warning("Source slide %d out of range, skipping", source_idx)
                continue
            
            target_idx = insert_position + offset
            new_slide = duplicate_slide(prs, source_idx, target_idx)
            new_slide_indices.append(target_idx)
            
            # Update slide title to include sensor name
            for shape in new_slide.shapes:
                if shape.name == "slide_title" and shape.has_text_frame:
                    current_title = shape.text_frame.text
                    # Remove previous sensor name if present
                    for prev_sensor in sensors[:sensor_idx]:
                        current_title = current_title.replace(f" ({prev_sensor})", "")
                    # Add new sensor name
                    if sensor_id not in current_title:
                        new_title = f"{current_title} ({sensor_id})"
                        update_slide_title(new_slide, new_title)
                    break
            
            # Rename placeholder shapes with sensor suffix
            placeholder_names = [
                "ndr_outbound_data_1", "ndr_outbound_data_2", "ndr_top_ip",
                "ndr_top_urls", "ndr_ext_dest", "ndr_country",
                "ndr_beaconing", "ndr_sensitive_data"
            ]
            for placeholder_name in placeholder_names:
                # Try both original name and first sensor's name
                if not rename_shape_with_suffix(new_slide, placeholder_name, f"_{sensor_id.lower()}"):
                    rename_shape_with_suffix(new_slide, f"{placeholder_name}_{first_sensor.lower()}", f"_{sensor_id.lower()}")
        
        sensor_slide_map[sensor_id] = new_slide_indices
        logger.info("Sensor '%s': duplicated to slides %s", sensor_id, new_slide_indices)
    
    return sensor_slide_map
