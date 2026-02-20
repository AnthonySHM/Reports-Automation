"""
Table Optimization Agent for PowerPoint Reports.

Post-processes tables created by ``place_csv_tables()`` by rendering each
table as a styled PNG image that fits exactly within the grey content
panel on each slide.

Usage::

    from core.table_optimizer import optimize_tables
    n = optimize_tables(pptx_path, csv_assets)
"""

from __future__ import annotations

import csv
import logging
import os
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation as PptxPresentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from core.drive_agent import CSVTableAsset, _is_source_column
from core.slide_utils import duplicate_slide, update_slide_title

logger = logging.getLogger("shm.table_optimizer")

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
EMU_PER_INCH = 914_400

# Image rendering resolution
IMG_DPI = 200

# Maximum data rows per table image before paginating onto continuation slides
MAX_IMAGE_ROWS = 40

# Colours (RGBA tuples for PIL) — match Mantix4 template
CLR_TRANSPARENT = (0, 0, 0, 0)
CLR_BLACK = (0x00, 0x00, 0x00, 0xFF)
CLR_WHITE = (0xFF, 0xFF, 0xFF, 0xFF)
CLR_TBL_DARK = (0x0B, 0x53, 0x94, 0xFF)
CLR_TBL_LIGHT = (0x3D, 0x85, 0xC6, 0xFF)
CLR_TITLE_BG = (0x07, 0x3B, 0x6B, 0xFF)   # slightly darker blue for title row

# ---------------------------------------------------------------------------
# Per-slide usable area (inches) — matches grey panel exactly
# ---------------------------------------------------------------------------
NDR_SLIDE_INDICES = set(range(5, 13))     # NDR chart slides (0-based)
TABLE_SLIDE_INDICES = {14, 15}            # patch-table slides

# Grey panel dimensions (from create_template.py)
#   NDR slides:     panel at (1.44", 1.00", 7.12" x 4.00")
#   Table slides:   panel at (1.44", 1.10", 7.12" x height up to note area)
#   Regular slides: panel at (1.44", 1.10", 7.12" x 3.70")
# Image fills the grey panel exactly — no extra padding.
_CONTENT_L = 1.44   # grey panel left
_CONTENT_W = 7.12   # grey panel width
_FOOTER_Y = 5.08    # footer top position
_NOTE_HEIGHT = 0.50 # note panel height
_NOTE_GAP = 0.08    # gap between note and footer
_NOTE_PADDING = 0.15 # padding between table and note

# For table slides: calculate max height to note area
# note_top = FOOTER_Y - NOTE_HEIGHT - NOTE_GAP = 5.08 - 0.50 - 0.08 = 4.50
# table_max_height = note_top - 1.10 (table_top) - NOTE_PADDING = 4.50 - 1.10 - 0.15 = 3.25
_NOTE_TOP = _FOOTER_Y - _NOTE_HEIGHT - _NOTE_GAP  # 4.50"
_TABLE_MAX_HEIGHT = _NOTE_TOP - 1.10 - _NOTE_PADDING  # 3.25"


@dataclass(frozen=True)
class _Area:
    """Usable rectangle for a table on a given slide type (inches)."""
    left: float
    top: float
    width: float
    max_height: float


# NDR chart slides: panel 1.00"–5.00"
_AREA_NDR = _Area(left=_CONTENT_L, top=1.00, width=_CONTENT_W, max_height=4.00)
# Table slides: constrained to not overlap note area (max 3.25" height)
_AREA_TBL = _Area(left=_CONTENT_L, top=1.10, width=_CONTENT_W, max_height=_TABLE_MAX_HEIGHT)
# Regular slides: default panel area
_AREA_DEF = _Area(left=_CONTENT_L, top=1.10, width=_CONTENT_W, max_height=3.70)


def _area_for(slide_index: int, ndr_indices: set[int] | None = None) -> _Area:
    ndr = ndr_indices if ndr_indices is not None else NDR_SLIDE_INDICES
    if slide_index in ndr:
        return _AREA_NDR
    if slide_index in TABLE_SLIDE_INDICES:
        return _AREA_TBL
    return _AREA_DEF


# ---------------------------------------------------------------------------
# CSV reader
# ---------------------------------------------------------------------------
@dataclass
class _CSVData:
    """Filtered CSV content with per-column metrics."""
    headers: list[str]
    rows: list[list[str]]
    col_max_len: list[int]
    col_avg_len: list[float]
    title: str | None = None        # metadata title (e.g. "Beaconing Score")
    date: str | None = None          # date stamp (e.g. "2026-01-04")
    sensor_name: str | None = None   # sensor identifier (e.g. "VAPRD")


def _read_csv(path: Path) -> _CSVData | None:
    """Read a CSV, extract metadata title, drop source columns."""
    try:
        with open(path, newline="", encoding="utf-8-sig") as fh:
            raw = list(csv.reader(fh))
    except Exception as exc:
        logger.error("Cannot read CSV '%s': %s", path, exc)
        return None

    if not raw:
        return None

    # Extract title, date, and sensor from metadata first row (empty first cell)
    # Metadata row format: [empty, sensor_id, date, title]
    title: str | None = None
    date: str | None = None
    sensor_name: str | None = None
    if len(raw) >= 2 and raw[0][0].strip() == "":
        meta = raw[0]
        # Title is the last non-empty cell in the metadata row
        title = next((c.strip() for c in reversed(meta) if c.strip()), None)
        # Sensor name is cell index 1, date is cell index 2
        if len(meta) > 1 and meta[1].strip():
            sensor_name = meta[1].strip()
        if len(meta) > 2 and meta[2].strip():
            date = meta[2].strip()
        raw = raw[1:]
    if not raw:
        return None

    header = raw[0]
    keep = [i for i, c in enumerate(header) if not _is_source_column(c)]
    if not keep:
        return None

    headers = [header[i] for i in keep]
    rows = [[r[i] if i < len(r) else "" for i in keep] for r in raw[1:]]

    nc = len(headers)
    col_max = [len(h) for h in headers]
    col_sum = [float(len(h)) for h in headers]
    for row in rows:
        for ci in range(nc):
            ln = len(row[ci]) if ci < len(row) else 0
            if ln > col_max[ci]:
                col_max[ci] = ln
            col_sum[ci] += ln

    denom = len(rows) + 1
    return _CSVData(
        headers=headers,
        rows=rows,
        col_max_len=col_max,
        col_avg_len=[s / denom for s in col_sum],
        title=title,
        date=date,
        sensor_name=sensor_name,
    )


# ---------------------------------------------------------------------------
# Font loading
# ---------------------------------------------------------------------------
_font_cache: dict[tuple[int, bool], ImageFont.FreeTypeFont] = {}

_FONT_NAMES_BOLD = ["calibrib.ttf", "Calibri Bold.ttf", "arialbd.ttf", "Arial Bold.ttf"]
_FONT_NAMES_REGULAR = ["calibri.ttf", "Calibri.ttf", "arial.ttf", "Arial.ttf"]


def _load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    """Load a TrueType font, with caching. Falls back to PIL default."""
    key = (size, bold)
    if key in _font_cache:
        return _font_cache[key]

    names = _FONT_NAMES_BOLD if bold else _FONT_NAMES_REGULAR
    for name in names:
        try:
            font = ImageFont.truetype(name, size)
            _font_cache[key] = font
            return font
        except OSError:
            full = os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts", name)
            try:
                font = ImageFont.truetype(full, size)
                _font_cache[key] = font
                return font
            except OSError:
                continue

    font = ImageFont.load_default()
    _font_cache[key] = font
    return font


# ---------------------------------------------------------------------------
# Column width measurement (content-fitted)
# ---------------------------------------------------------------------------
def _text_width(text: str, font: ImageFont.FreeTypeFont) -> int:
    """Return pixel width of *text* rendered with *font*."""
    try:
        return int(font.getlength(text))
    except AttributeError:
        return font.getsize(text)[0]


def _measure_natural_col_widths(
    data: _CSVData,
    font: ImageFont.FreeTypeFont,
    font_bold: ImageFont.FreeTypeFont,
    pad_x: int,
) -> list[int]:
    """Measure pixel width needed for each column based on actual content.

    Returns natural widths (no capping) — the caller decides whether to
    shrink the font if the total exceeds the available width.
    """
    n = len(data.headers)
    if n == 0:
        return []

    widths: list[int] = []
    for ci in range(n):
        max_w = _text_width(data.headers[ci], font_bold)
        for row in data.rows:
            cell = row[ci] if ci < len(row) else ""
            tw = _text_width(cell, font)
            if tw > max_w:
                max_w = tw
        widths.append(max_w + 2 * pad_x)

    return widths


# ---------------------------------------------------------------------------
# Image rendering
# ---------------------------------------------------------------------------
def _render_table_image(data: _CSVData, width_px: int, height_px: int) -> Image.Image:
    """Render CSV data as a transparent-background table image.

    - If a metadata title exists, it is rendered as a full-width row
      at the top of the table.
    - Columns are sized to fit their longest content (no truncation).
    - Font is auto-sized so all columns fit horizontally.
    - Table is centred horizontally; transparent background lets the
      grey slide panel show through.
    """
    n_data = len(data.rows)

    # Build combined title row: "Title  |  Sensor  |  Date"
    title_parts = []
    if data.title:
        title_parts.append(data.title)
    if data.sensor_name:
        title_parts.append(data.sensor_name)
    if data.date:
        title_parts.append(data.date)
    combined_title = "  |  ".join(title_parts) if title_parts else None
    has_title = bool(combined_title)

    # Total visible rows: optional title + header + data (min 1 for "No data")
    n_total = (1 if has_title else 0) + 1 + max(n_data, 1)

    # --- Pick the largest font where all columns fit horizontally ----------
    max_font_from_height = max(8, min(28, int((height_px / n_total) * 0.60)))

    font_size = max_font_from_height
    while font_size > 8:
        test_font = _load_font(font_size, bold=False)
        test_bold = _load_font(font_size, bold=True)
        test_pad = max(3, int(font_size * 0.4))
        natural = _measure_natural_col_widths(data, test_font, test_bold, test_pad)
        if sum(natural) <= width_px:
            break
        font_size -= 1

    font = _load_font(font_size, bold=False)
    font_bold = _load_font(font_size, bold=True)
    pad_x = max(3, int(font_size * 0.4))

    # Column widths at the chosen font size
    col_widths = _measure_natural_col_widths(data, font, font_bold, pad_x)
    table_w = sum(col_widths)

    # If still too wide at min font, scale down proportionally (last resort)
    if table_w > width_px and table_w > 0:
        factor = width_px / table_w
        col_widths = [max(1, int(w * factor)) for w in col_widths]
        col_widths[-1] += width_px - sum(col_widths)
        table_w = width_px

    # Centre the table horizontally
    x_offset = max(0, (width_px - table_w) // 2)

    # Compact height: size to fit rows instead of stretching to full panel
    preferred_rh = max(int(font_size * 2.5), font_size + 10)
    compact_h = min(preferred_rh * n_total, height_px)

    # Row heights — distribute evenly, remainder to first rows
    base_rh = compact_h // n_total
    remainder = compact_h - base_rh * n_total
    row_heights = [base_rh + (1 if i < remainder else 0) for i in range(n_total)]

    row_h = base_rh
    pad_y = max(2, int((row_h - font_size) / 2))

    img = Image.new("RGBA", (width_px, compact_h), CLR_TRANSPARENT)
    draw = ImageDraw.Draw(img)

    def _draw_row(y: int, rh: int, texts: list[str], bg: tuple, use_bold: bool):
        f = font_bold if use_bold else font
        x = x_offset
        for ci, cw in enumerate(col_widths):
            draw.rectangle([x, y, x + cw - 1, y + rh - 1], fill=bg)
            cell_text = texts[ci] if ci < len(texts) else ""
            draw.text((x + pad_x, y + pad_y), cell_text, fill=CLR_WHITE, font=f)
            x += cw

    def _draw_title_row(y: int, rh: int, title: str):
        """Full-width title bar spanning all columns."""
        draw.rectangle(
            [x_offset, y, x_offset + table_w - 1, y + rh - 1],
            fill=CLR_TITLE_BG,
        )
        draw.text(
            (x_offset + pad_x, y + pad_y), title,
            fill=CLR_WHITE, font=font_bold,
        )

    y = 0
    ri_offset = 0  # index into row_heights

    # Title row (title | sensor | date)
    if has_title:
        _draw_title_row(y, row_heights[ri_offset], combined_title)
        y += row_heights[ri_offset]
        ri_offset += 1

    # Header row
    _draw_row(y, row_heights[ri_offset], data.headers, CLR_BLACK, use_bold=True)
    y += row_heights[ri_offset]
    ri_offset += 1

    # Data rows
    if n_data == 0:
        _draw_row(y, row_heights[ri_offset], ["No data"], CLR_TBL_DARK, use_bold=False)
    else:
        for ri, row_data in enumerate(data.rows):
            rh = row_heights[ri_offset + ri]
            bg = CLR_TBL_DARK if ri % 2 == 0 else CLR_TBL_LIGHT
            _draw_row(y, rh, row_data, bg, use_bold=False)
            y += rh

    return img


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def _get_slide_title(slide) -> str:
    """Return the text of the ``slide_title`` shape, or ``""``."""
    for shape in slide.shapes:
        if shape.name == "slide_title" and shape.has_text_frame:
            return shape.text_frame.text
    return ""


def _render_and_place(slide, data: _CSVData, area: _Area, shape_name: str) -> bool:
    """Render *data* as a PNG table image and place it on *slide*.

    Removes the first shape whose name matches *shape_name* (or the first
    table shape) before inserting the image.  Returns ``True`` on success.
    """
    img_w = int(area.width * IMG_DPI)
    img_h = int(area.max_height * IMG_DPI)
    img = _render_table_image(data, img_w, img_h)
    actual_h = img.size[1]

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".png")
    os.close(tmp_fd)
    try:
        img.save(tmp_path, "PNG")

        # Remove old table / image shape
        removed = False
        for shape in list(slide.shapes):
            if shape.name == shape_name:
                shape._element.getparent().remove(shape._element)
                removed = True
                break
        if not removed:
            for shape in list(slide.shapes):
                if shape.has_table:
                    shape._element.getparent().remove(shape._element)
                    removed = True
                    break
        if not removed:
            return False

        pic = slide.shapes.add_picture(
            tmp_path,
            Emu(int(area.left * EMU_PER_INCH)),
            Emu(int(area.top * EMU_PER_INCH)),
            Emu(int(area.width * EMU_PER_INCH)),
            Emu(int(actual_h / IMG_DPI * EMU_PER_INCH)),
        )
        pic.name = shape_name
        return True
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------
def optimize_tables(
    pptx_path: Path,
    assets: list[CSVTableAsset],
    *,
    skip_indices: set[int] | None = None,
    ndr_slide_indices: set[int] | None = None,
    report_date: str | None = None,
) -> int:
    """Convert tables to images in a PPTX at grey-panel dimensions.

    Re-reads the original CSV files, renders each as a styled PNG image
    that fills the grey content panel exactly, removes the old table
    shape, and inserts the image on the slide.

    Args:
        pptx_path: Path to PPTX file with tables to convert.
        assets: CSVTableAsset objects produced by ``DriveAgent``.
            Each asset must have its ``slide_index`` and ``shape_name``
            set to the *actual* values on the PPTX (i.e. already
            remapped for multi-sensor templates).
        skip_indices: Ignored (kept for API compatibility).
        ndr_slide_indices: Set of 0-based slide indices that are NDR
            chart slides (use the wider 4.00" panel).  Defaults to
            ``{5..12}`` for the standard 18-slide template.

    Returns:
        Number of tables converted to images.
    """
    if not assets:
        return 0

    prs = PptxPresentation(str(pptx_path))
    count = 0

    # Sort by slide_index ascending so that slide insertions at higher
    # indices don't shift earlier assets.
    sorted_assets = sorted(assets, key=lambda a: a.slide_index)
    slide_offset = 0  # running count of continuation slides inserted

    for asset in sorted_assets:
        actual_idx = asset.slide_index + slide_offset

        if actual_idx >= len(prs.slides):
            logger.warning("Slide %d out of range, skipping", actual_idx)
            continue

        # 1. Read full CSV data (including metadata title)
        data = _read_csv(asset.local_path)
        if data is None:
            logger.warning(
                "Could not read CSV for '%s', skipping",
                asset.original_filename,
            )
            continue

        # Set sensor name from asset if not already extracted from CSV metadata
        if asset.sensor_id and asset.sensor_id != "DEFAULT":
            if not data.sensor_name:
                data.sensor_name = asset.sensor_id
        # Set date from report_date if not already extracted from CSV metadata
        if not data.date and report_date:
            data.date = report_date

        # 2. Slide-specific constraints (computed from *original* index so
        #    continuation slides inherit the correct area)
        area = _area_for(asset.slide_index, ndr_slide_indices)
        slide = prs.slides[actual_idx]

        # 2b. Empty table → replace with "No Data" text box
        if not data.rows:
            removed = False
            for shape in list(slide.shapes):
                if shape.name == asset.shape_name or shape.has_table:
                    shape._element.getparent().remove(shape._element)
                    removed = True
                    break
            _add_no_data_textbox(slide, area)
            logger.info(
                "Empty CSV '%s' on slide %d → replaced with 'No Data'",
                asset.original_filename, actual_idx,
            )
            count += 1
            continue

        # 3. Pagination: split into pages of MAX_IMAGE_ROWS
        n_rows = len(data.rows)

        if n_rows > MAX_IMAGE_ROWS:
            pages = [
                data.rows[i : i + MAX_IMAGE_ROWS]
                for i in range(0, n_rows, MAX_IMAGE_ROWS)
            ]
            n_cont = len(pages) - 1  # number of continuation slides

            # Read the slide title before any modifications
            orig_title = _get_slide_title(slide)

            # Duplicate the original slide for each continuation page.
            # Done BEFORE modifying the original so duplicates inherit all
            # shapes (header, footer, note panel, table placeholder).
            for ci in range(n_cont):
                duplicate_slide(prs, actual_idx, actual_idx + 1 + ci)

            # --- Page 1: render on the original slide ---
            page1_data = _CSVData(
                headers=data.headers,
                rows=pages[0],
                col_max_len=data.col_max_len,
                col_avg_len=data.col_avg_len,
                title=data.title,
                date=data.date,
                sensor_name=data.sensor_name,
            )
            if _render_and_place(slide, page1_data, area, asset.shape_name):
                count += 1
                logger.info(
                    "Converted '%s' slide %d to image (page 1/%d): "
                    "%d cols, %d rows",
                    asset.original_filename, actual_idx,
                    len(pages), len(data.headers), len(pages[0]),
                )
            else:
                logger.warning(
                    "No table found on slide %d for '%s', skipping",
                    actual_idx, asset.original_filename,
                )

            # --- Continuation pages ---
            for ci in range(n_cont):
                cont_slide = prs.slides[actual_idx + 1 + ci]

                # Update title to include "(cont.)"
                cont_title = f"{orig_title} (cont.)" if orig_title else "(cont.)"
                update_slide_title(cont_slide, cont_title)

                page_data = _CSVData(
                    headers=data.headers,
                    rows=pages[1 + ci],
                    col_max_len=data.col_max_len,
                    col_avg_len=data.col_avg_len,
                    title=data.title,
                    date=data.date,
                    sensor_name=data.sensor_name,
                )
                if _render_and_place(cont_slide, page_data, area, asset.shape_name):
                    count += 1
                    logger.info(
                        "Continuation slide %d (page %d/%d): %d rows",
                        actual_idx + 1 + ci, ci + 2, len(pages),
                        len(pages[1 + ci]),
                    )

            slide_offset += n_cont

        else:
            # No pagination needed — single page
            if _render_and_place(slide, data, area, asset.shape_name):
                count += 1
                logger.info(
                    "Converted '%s' slide %d to image: %d cols, %d rows%s, "
                    "%.2f\"x%.2f\"",
                    asset.original_filename, actual_idx,
                    len(data.headers), n_rows,
                    f" (title: {data.title})" if data.title else "",
                    area.width, area.max_height,
                )
            else:
                logger.warning(
                    "No table found on slide %d for '%s', skipping",
                    actual_idx, asset.original_filename,
                )

    prs.save(str(pptx_path))
    logger.info("Saved PPTX with %d table images -> %s", count, pptx_path)
    return count


# ---------------------------------------------------------------------------
# "No Data" text box helper
# ---------------------------------------------------------------------------
def _add_no_data_textbox(slide, area: _Area) -> None:
    """Add a centred 'No Data' text box in the given area."""
    box = slide.shapes.add_textbox(
        Emu(int(area.left * EMU_PER_INCH)),
        Emu(int(area.top * EMU_PER_INCH)),
        Emu(int(area.width * EMU_PER_INCH)),
        Inches(0.5),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "ctr")
    p = tf.paragraphs[0]
    p.text = "No Data"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# Cleanup leftover template tables
# ---------------------------------------------------------------------------
def cleanup_placeholder_tables(pptx_path: Path, table_slide_indices: set[int]) -> int:
    """Replace leftover template TABLE shapes with 'No Data' text boxes.

    Opens the PPTX, checks each slide in *table_slide_indices* for shapes
    with ``.has_table``.  Any found are removed and replaced with a centred
    'No Data' text box.

    Returns:
        Number of placeholder tables cleaned up.
    """
    prs = PptxPresentation(str(pptx_path))
    cleaned = 0

    for idx in sorted(table_slide_indices):
        if idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        found_table = False
        for shape in list(slide.shapes):
            if shape.has_table:
                shape._element.getparent().remove(shape._element)
                found_table = True
        if found_table:
            area = _area_for(idx)
            _add_no_data_textbox(slide, area)
            cleaned += 1
            logger.info("Cleaned placeholder table on slide %d", idx)

    if cleaned:
        prs.save(str(pptx_path))
        logger.info("Saved PPTX after cleaning %d placeholder tables", cleaned)

    return cleaned
