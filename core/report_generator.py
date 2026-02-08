"""
SHM Report Generator — Orchestrator

Reads a layout manifest (CSV or JSON), resolves source images from the
RPA Legacy Bridge, computes contain-fit transforms, and produces the
final PowerPoint presentation with a full audit trail.

Designed to be called from:
  - The Flask API (Google Apps Script → HTTP → this module)
  - Direct CLI invocation for debugging / manual runs
"""

from __future__ import annotations

import json
import logging
import os
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from core.transform import Rect, TransformResult, compute_contain_transform

logger = logging.getLogger("shm.report_generator")


# ---------------------------------------------------------------------------
# Error types for the Legacy Bridge integration
# ---------------------------------------------------------------------------
class SlideLayoutError(Exception):
    """Raised when a referenced slide or placeholder is missing."""


class ImageSourceError(Exception):
    """Raised when an image file cannot be found or read."""


# ---------------------------------------------------------------------------
# Layout Manifest Schema
# ---------------------------------------------------------------------------
MANIFEST_COLUMNS = [
    "slide_index",       # 0-based slide number
    "placeholder_id",    # placeholder index on that slide  (-1 = free-place)
    "image_path",        # path to source image (absolute or relative to assets_dir)
    "box_x_emu",         # destination box origin X  (EMU, used when placeholder_id == -1)
    "box_y_emu",         # destination box origin Y  (EMU)
    "box_w_emu",         # destination box width     (EMU)
    "box_h_emu",         # destination box height    (EMU)
    "label",             # human-readable label for audit log
]


# ---------------------------------------------------------------------------
# Placement Record (one per image)
# ---------------------------------------------------------------------------
@dataclass
class PlacementRecord:
    """Single entry in the audit trail."""
    label: str
    slide_index: int
    image_path: str
    transform: TransformResult
    status: str = "ok"          # "ok" | "skipped" | "error"
    error_message: str = ""
    timestamp: float = 0.0

    def to_dict(self) -> dict:
        return {
            "label": self.label,
            "slide_index": self.slide_index,
            "image_path": self.image_path,
            "status": self.status,
            "error_message": self.error_message,
            "timestamp": self.timestamp,
            "transform": self.transform.to_audit_dict() if self.transform else None,
        }


# ---------------------------------------------------------------------------
# Core Generator
# ---------------------------------------------------------------------------
class ReportGenerator:
    """Reads a manifest, applies Algorithm 1 transforms, and writes a PPTX.

    Parameters
    ----------
    template_path : str | Path
        Path to a .pptx template. Use ``None`` for a blank presentation.
    output_path : str | Path
        Destination path for the generated report.
    assets_dir : str | Path
        Base directory for resolving relative image paths.
    """

    def __init__(
        self,
        template_path: Optional[str | Path] = None,
        output_path: str | Path = "output/report.pptx",
        assets_dir: str | Path = "assets",
    ):
        self.template_path = Path(template_path) if template_path else None
        self.output_path = Path(output_path)
        self.assets_dir = Path(assets_dir)
        self.audit_trail: list[PlacementRecord] = []

    # ------------------------------------------------------------------
    # Manifest loading
    # ------------------------------------------------------------------
    def load_manifest(self, manifest_path: str | Path) -> pd.DataFrame:
        """Load a layout manifest from CSV or JSON.

        Returns a DataFrame normalised to ``MANIFEST_COLUMNS``.
        """
        path = Path(manifest_path)
        if not path.exists():
            raise FileNotFoundError(f"Manifest not found: {path}")

        if path.suffix.lower() == ".csv":
            df = pd.read_csv(path)
        elif path.suffix.lower() == ".json":
            df = pd.DataFrame(json.loads(path.read_text(encoding="utf-8")))
        else:
            raise ValueError(f"Unsupported manifest format: {path.suffix}")

        missing = set(MANIFEST_COLUMNS) - set(df.columns)
        if missing:
            raise ValueError(f"Manifest is missing columns: {missing}")

        logger.info("Loaded manifest with %d entries from %s", len(df), path)
        return df

    # ------------------------------------------------------------------
    # Image dimension reader
    # ------------------------------------------------------------------
    def _resolve_image(self, image_path: str) -> Path:
        """Resolve an image path against the assets directory."""
        p = Path(image_path)
        if p.is_absolute() and p.exists():
            return p
        resolved = self.assets_dir / p
        if resolved.exists():
            return resolved
        raise ImageSourceError(
            f"Image not found at '{image_path}' "
            f"(also checked '{resolved}')"
        )

    @staticmethod
    def _get_image_dimensions_emu(image_path: Path) -> tuple[int, int]:
        """Read image pixel dimensions and convert to EMU.

        PowerPoint uses EMU (English Metric Units):
            1 inch = 914400 EMU
        We assume 96 DPI unless the image provides its own DPI metadata.
        """
        with Image.open(image_path) as img:
            w_px, h_px = img.size
            dpi_x, dpi_y = img.info.get("dpi", (96, 96))
            # Ensure DPI values are valid
            dpi_x = dpi_x if dpi_x and dpi_x > 0 else 96
            dpi_y = dpi_y if dpi_y and dpi_y > 0 else 96

        w_emu = int(round(w_px / dpi_x * 914400))
        h_emu = int(round(h_px / dpi_y * 914400))
        return w_emu, h_emu

    # ------------------------------------------------------------------
    # Placeholder resolver
    # ------------------------------------------------------------------
    @staticmethod
    def _get_placeholder_rect(slide, placeholder_id: int) -> Rect:
        """Extract position and size from an existing placeholder shape."""
        try:
            ph = slide.placeholders[placeholder_id]
        except KeyError:
            available = [p.placeholder_format.idx for p in slide.placeholders]
            raise SlideLayoutError(
                f"Placeholder {placeholder_id} not found on slide. "
                f"Available: {available}"
            )
        return Rect(
            x=int(ph.left),
            y=int(ph.top),
            width=int(ph.width),
            height=int(ph.height),
        )

    # ------------------------------------------------------------------
    # Single-image placement
    # ------------------------------------------------------------------
    def _place_image(
        self,
        slide,
        image_path: Path,
        dest_rect: Rect,
        label: str,
        slide_index: int,
    ) -> PlacementRecord:
        """Compute the contain transform and add the image to the slide."""
        w_emu, h_emu = self._get_image_dimensions_emu(image_path)
        source = Rect(x=0, y=0, width=w_emu, height=h_emu)

        transform = compute_contain_transform(source, dest_rect)

        # Apply to PowerPoint via add_picture
        slide.shapes.add_picture(
            str(image_path),
            Emu(transform.dest_x),
            Emu(transform.dest_y),
            Emu(transform.dest_width),
            Emu(transform.dest_height),
        )

        record = PlacementRecord(
            label=label,
            slide_index=slide_index,
            image_path=str(image_path),
            transform=transform,
            status="ok",
            timestamp=time.time(),
        )
        logger.info(
            "Placed '%s' on slide %d  → (%d, %d) %d×%d  [scale=%.4f, axis=%s]",
            label, slide_index,
            transform.dest_x, transform.dest_y,
            transform.dest_width, transform.dest_height,
            transform.scale_factor, transform.scale_axis,
        )
        return record

    # ------------------------------------------------------------------
    # Ensure enough slides exist
    # ------------------------------------------------------------------
    @staticmethod
    def _ensure_slide_count(prs: Presentation, required: int):
        """Add blank slides until the presentation has at least *required*."""
        while len(prs.slides) < required:
            layout = prs.slide_layouts[6]  # blank layout
            prs.slides.add_slide(layout)

    # ------------------------------------------------------------------
    # Main generation pipeline
    # ------------------------------------------------------------------
    def generate(self, manifest_path: str | Path) -> Path:
        """Execute the full report generation pipeline.

        1. Load the manifest.
        2. Open / create the presentation.
        3. For each row, resolve the image, compute the transform, and place it.
        4. Save the PPTX and the audit trail.

        Returns the path to the generated .pptx file.
        """
        df = self.load_manifest(manifest_path)

        # Open template or create blank
        if self.template_path and self.template_path.exists():
            prs = Presentation(str(self.template_path))
            logger.info("Opened template: %s", self.template_path)
        else:
            prs = Presentation()
            logger.info("Created blank presentation (no template).")

        self.audit_trail.clear()

        if len(df) == 0:
            # Template-only: no manifest image placements (e.g. content from Drive)
            logger.info("Manifest is empty; using template as-is.")
        else:
            max_slide = int(df["slide_index"].max())
            self._ensure_slide_count(prs, max_slide + 1)

        for _, row in df.iterrows():
            slide_idx = int(row["slide_index"])
            placeholder_id = int(row["placeholder_id"])
            label = str(row["label"])

            # --- Resolve image (with Legacy Bridge error handling) ---
            try:
                image_path = self._resolve_image(str(row["image_path"]))
            except ImageSourceError as exc:
                logger.warning("SKIP '%s': %s", label, exc)
                self.audit_trail.append(PlacementRecord(
                    label=label,
                    slide_index=slide_idx,
                    image_path=str(row["image_path"]),
                    transform=None,
                    status="error",
                    error_message=str(exc),
                    timestamp=time.time(),
                ))
                continue

            slide = prs.slides[slide_idx]

            # --- Determine destination rectangle ---
            if placeholder_id >= 0:
                try:
                    dest_rect = self._get_placeholder_rect(slide, placeholder_id)
                except SlideLayoutError as exc:
                    logger.warning("SKIP '%s': %s", label, exc)
                    self.audit_trail.append(PlacementRecord(
                        label=label,
                        slide_index=slide_idx,
                        image_path=str(image_path),
                        transform=None,
                        status="error",
                        error_message=str(exc),
                        timestamp=time.time(),
                    ))
                    continue
            else:
                # Free-placement using explicit coordinates from manifest
                dest_rect = Rect(
                    x=int(row["box_x_emu"]),
                    y=int(row["box_y_emu"]),
                    width=int(row["box_w_emu"]),
                    height=int(row["box_h_emu"]),
                )

            # --- Place the image ---
            try:
                record = self._place_image(slide, image_path, dest_rect, label, slide_idx)
                self.audit_trail.append(record)
            except Exception as exc:
                logger.error("ERROR placing '%s': %s", label, exc, exc_info=True)
                self.audit_trail.append(PlacementRecord(
                    label=label,
                    slide_index=slide_idx,
                    image_path=str(image_path),
                    transform=None,
                    status="error",
                    error_message=str(exc),
                    timestamp=time.time(),
                ))

        # --- Save outputs ---
        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(self.output_path))
        logger.info("Saved presentation → %s", self.output_path)

        audit_path = self.output_path.with_suffix(".audit.json")
        audit_data = {
            "generated_at": time.time(),
            "template": str(self.template_path) if self.template_path else None,
            "manifest": str(manifest_path),
            "output": str(self.output_path),
            "placements": [r.to_dict() for r in self.audit_trail],
            "summary": {
                "total": len(self.audit_trail),
                "ok": sum(1 for r in self.audit_trail if r.status == "ok"),
                "errors": sum(1 for r in self.audit_trail if r.status == "error"),
            },
        }
        audit_path.write_text(json.dumps(audit_data, indent=2), encoding="utf-8")
        logger.info("Saved audit trail → %s", audit_path)

        return self.output_path
