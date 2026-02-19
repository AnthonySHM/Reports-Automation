"""
Patch Trend Chart Generator

Generates a trend chart showing Microsoft KB patches and 3rd party software patches
over time using PIL/Pillow (no matplotlib required).
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger("shm.patch_chart")

# Chart styling to match reference
CHART_COLORS = {
    'microsoft': (231, 76, 60),      # Red for Microsoft KB patches
    'software': (82, 190, 128),      # Green for 3rd party software
    'grid': (213, 216, 220),         # Light gray for grid lines
    'text': (44, 62, 80),            # Dark text
    'bg': (255, 255, 255),           # White background
    'axis': (0, 0, 0),               # Black for axes
}


def generate_patch_trend_chart(
    dates: list[str],
    microsoft_counts: list[int],
    software_counts: list[int],
    output_path: str | Path,
    *,
    width: int = 1500,
    height: int = 675,
) -> Path:
    """Generate a patch trend chart image using PIL.
    
    Parameters
    ----------
    dates : list[str]
        List of date strings (e.g., ["November 10", "December 8", "January 5"])
    microsoft_counts : list[int]
        Microsoft KB patch counts for each date
    software_counts : list[int]
        3rd party software patch counts for each date
    output_path : str | Path
        Path where the chart image will be saved
    width : int
        Image width in pixels
    height : int
        Image height in pixels
    
    Returns
    -------
    Path
        Path to the saved chart image
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    # Create image
    img = Image.new('RGB', (width, height), CHART_COLORS['bg'])
    draw = ImageDraw.Draw(img)
    
    # Try to use a nice font, fallback to default
    try:
        title_font = ImageFont.truetype("arial.ttf", 42)
        label_font = ImageFont.truetype("arial.ttf", 32)
        tick_font = ImageFont.truetype("arial.ttf", 28)
        legend_font = ImageFont.truetype("arial.ttf", 28)
    except:
        title_font = ImageFont.load_default()
        label_font = ImageFont.load_default()
        tick_font = ImageFont.load_default()
        legend_font = ImageFont.load_default()
    
    # Chart margins — generous enough so labels never overlap data
    margin_top = 100
    margin_bottom = 175
    margin_left = 160
    margin_right = 100

    # Plot area
    plot_left = margin_left
    plot_right = width - margin_right
    plot_top = margin_top
    plot_bottom = height - margin_bottom
    plot_width = plot_right - plot_left
    plot_height = plot_bottom - plot_top

    # Title (centered above plot)
    title = "Missing Patch Quantity"
    title_bbox = draw.textbbox((0, 0), title, font=title_font)
    title_width = title_bbox[2] - title_bbox[0]
    title_x = (width - title_width) // 2
    draw.text((title_x, 20), title, fill=CHART_COLORS['text'], font=title_font)

    # Determine Y-axis range
    max_val = max(max(microsoft_counts), max(software_counts))
    y_max = ((max_val + 1) // 2 + 1) * 2  # Round up to next even number
    y_ticks = list(range(0, y_max + 1, 2))

    # Draw Y-axis gridlines and tick labels (right-aligned before the axis)
    for y_val in y_ticks:
        y_pos = plot_bottom - (y_val / y_max) * plot_height

        # Gridline
        draw.line(
            [(plot_left, y_pos), (plot_right, y_pos)],
            fill=CHART_COLORS['grid'], width=2,
        )

        # Y-axis tick label — right-aligned to sit just left of the axis
        label = str(y_val)
        lbl_bbox = draw.textbbox((0, 0), label, font=tick_font)
        lbl_w = lbl_bbox[2] - lbl_bbox[0]
        lbl_h = lbl_bbox[3] - lbl_bbox[1]
        draw.text(
            (plot_left - 15 - lbl_w, y_pos - lbl_h // 2),
            label, fill=CHART_COLORS['text'], font=tick_font,
        )

    # Draw axes
    draw.line([(plot_left, plot_top), (plot_left, plot_bottom)],
              fill=CHART_COLORS['axis'], width=3)
    draw.line([(plot_left, plot_bottom), (plot_right, plot_bottom)],
              fill=CHART_COLORS['axis'], width=3)

    # Y-axis label — rotated 90° so it never overlaps tick numbers
    y_label = "# of Patches"
    y_lbl_bbox = draw.textbbox((0, 0), y_label, font=label_font)
    y_lbl_w = y_lbl_bbox[2] - y_lbl_bbox[0]
    y_lbl_h = y_lbl_bbox[3] - y_lbl_bbox[1]
    txt_img = Image.new('RGBA', (y_lbl_w + 4, y_lbl_h + 4), (255, 255, 255, 0))
    ImageDraw.Draw(txt_img).text((0, 0), y_label,
                                  fill=CHART_COLORS['text'], font=label_font)
    txt_img = txt_img.rotate(90, expand=True)
    paste_x = 12
    paste_y = plot_top + (plot_height - txt_img.height) // 2
    img.paste(txt_img, (paste_x, paste_y), txt_img)

    # Calculate X positions
    n_points = len(dates)
    if n_points > 1:
        x_step = plot_width / (n_points - 1)
        x_positions = [plot_left + i * x_step for i in range(n_points)]
    else:
        x_positions = [plot_left + plot_width / 2]

    # Helper: data value → Y pixel
    def data_to_y(value):
        return plot_bottom - (value / y_max) * plot_height

    # Draw lines — Microsoft KB patches (red)
    ms_points = [(x_positions[i], data_to_y(microsoft_counts[i]))
                 for i in range(n_points)]
    for i in range(len(ms_points) - 1):
        draw.line([ms_points[i], ms_points[i + 1]],
                  fill=CHART_COLORS['microsoft'], width=5)

    # Draw lines — 3rd Party Software patches (green)
    sw_points = [(x_positions[i], data_to_y(software_counts[i]))
                 for i in range(n_points)]
    for i in range(len(sw_points) - 1):
        draw.line([sw_points[i], sw_points[i + 1]],
                  fill=CHART_COLORS['software'], width=5)

    # Draw markers
    marker_radius = 12
    for i in range(n_points):
        x, y = ms_points[i]
        draw.ellipse(
            [(x - marker_radius, y - marker_radius),
             (x + marker_radius, y + marker_radius)],
            fill=CHART_COLORS['microsoft'],
        )
        x, y = sw_points[i]
        draw.ellipse(
            [(x - marker_radius, y - marker_radius),
             (x + marker_radius, y + marker_radius)],
            fill=CHART_COLORS['software'],
        )

    # --- Bottom labels — spaced vertically so nothing overlaps ---

    # Row 1: X-axis date tick labels (just below the axis)
    date_labels_y = plot_bottom + 15
    for i, date in enumerate(dates):
        d_bbox = draw.textbbox((0, 0), date, font=tick_font)
        d_w = d_bbox[2] - d_bbox[0]
        x_center = x_positions[i]
        draw.text((x_center - d_w // 2, date_labels_y),
                  date, fill=CHART_COLORS['text'], font=tick_font)

    # Row 2: "Date" axis label (below the tick labels)
    x_label = "Date"
    xl_bbox = draw.textbbox((0, 0), x_label, font=label_font)
    xl_w = xl_bbox[2] - xl_bbox[0]
    x_label_y = date_labels_y + 40
    draw.text((plot_left + plot_width // 2 - xl_w // 2, x_label_y),
              x_label, fill=CHART_COLORS['text'], font=label_font)

    # Row 3: Legend (at the very bottom, well below "Date")
    legend_y = x_label_y + 48
    marker_size = 18

    ms_legend_text = "Microsoft KB Patches"
    sw_legend_text = "3rd Party Software Patches"

    ms_tw = draw.textbbox((0, 0), ms_legend_text, font=legend_font)[2]
    sw_tw = draw.textbbox((0, 0), sw_legend_text, font=legend_font)[2]
    legend_spacing = 60
    total_legend_w = (marker_size + 10 + ms_tw
                      + legend_spacing
                      + marker_size + 10 + sw_tw)
    legend_start_x = (width - total_legend_w) // 2

    # Microsoft legend item
    ms_lx = legend_start_x
    draw.ellipse(
        [(ms_lx, legend_y - marker_size // 2),
         (ms_lx + marker_size, legend_y + marker_size // 2)],
        fill=CHART_COLORS['microsoft'],
    )
    draw.text((ms_lx + marker_size + 10, legend_y - 12),
              ms_legend_text, fill=CHART_COLORS['text'], font=legend_font)

    # Software legend item
    sw_lx = ms_lx + marker_size + 10 + ms_tw + legend_spacing
    draw.ellipse(
        [(sw_lx, legend_y - marker_size // 2),
         (sw_lx + marker_size, legend_y + marker_size // 2)],
        fill=CHART_COLORS['software'],
    )
    draw.text((sw_lx + marker_size + 10, legend_y - 12),
              sw_legend_text, fill=CHART_COLORS['text'], font=legend_font)
    
    # Save image
    img.save(output_path, 'PNG', dpi=(150, 150))
    
    logger.info(
        "Generated patch trend chart → %s (%d dates, MS: %s, SW: %s)",
        output_path, len(dates), microsoft_counts, software_counts
    )
    
    return output_path


def find_chart_slide_index(pptx_path: str | Path) -> tuple[int, object] | None:
    """Find the first slide containing an embedded chart in a PPTX file.

    Parameters
    ----------
    pptx_path : str | Path
        Path to the PPTX file to inspect.

    Returns
    -------
    tuple[int, object] | None
        ``(slide_index, chart_shape)`` for the first chart found, or ``None``
        if no charts exist in the presentation.
    """
    from pptx import Presentation as PptxPresentation

    prs = PptxPresentation(str(pptx_path))
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_chart:
                logger.info(
                    "Found chart on slide %d (shape '%s', chart type: %s)",
                    idx, shape.name, shape.chart.chart_type,
                )
                return idx, shape
    logger.info("No embedded charts found in %s", pptx_path)
    return None


def update_chart_in_pptx(
    pptx_path: str | Path,
    slide_index: int,
    chart_shape_name: str,
    dates: list[str],
    microsoft_counts: list[int],
    software_counts: list[int],
) -> None:
    """Update chart data in a PPTX while preserving its visual style.

    Parameters
    ----------
    pptx_path : str | Path
        Path to the PPTX file containing the chart.
    slide_index : int
        0-based index of the slide containing the chart.
    chart_shape_name : str
        Name of the chart shape to update.
    dates : list[str]
        Category labels (formatted date strings).
    microsoft_counts : list[int]
        Microsoft KB patch counts per date.
    software_counts : list[int]
        3rd-party software patch counts per date.
    """
    from pptx import Presentation as PptxPresentation
    from pptx.chart.data import CategoryChartData

    prs = PptxPresentation(str(pptx_path))
    slide = prs.slides[slide_index]

    chart_shape = None
    for shape in slide.shapes:
        if shape.name == chart_shape_name and shape.has_chart:
            chart_shape = shape
            break

    if chart_shape is None:
        logger.warning(
            "Chart shape '%s' not found on slide %d of %s",
            chart_shape_name, slide_index, pptx_path,
        )
        return

    chart_data = CategoryChartData()
    chart_data.categories = dates
    chart_data.add_series("Microsoft KB", microsoft_counts)
    chart_data.add_series("Software Packages", software_counts)

    chart_shape.chart.replace_data(chart_data)
    prs.save(str(pptx_path))
    logger.info(
        "Updated chart '%s' on slide %d with %d data points",
        chart_shape_name, slide_index, len(dates),
    )
