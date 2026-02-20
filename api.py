"""
SHM Hybrid Architecture — Flask API

Provides HTTP endpoints that Google Apps Script (or any RPA frontend)
can call to trigger report generation.

Endpoints
---------
POST /api/v1/generate
    Trigger a report build from a layout manifest.

GET  /api/v1/health
    Liveness probe.

GET  /api/v1/audit/<report_id>
    Retrieve the audit trail for a completed report.
"""

from __future__ import annotations

import functools
import json
import logging
import os
import re
import secrets
import uuid
from datetime import datetime
from pathlib import Path

from flask import Flask, jsonify, request, send_file, session, redirect
from flask_cors import CORS
from werkzeug.security import generate_password_hash, check_password_hash
from pptx import Presentation as PptxPresentation

from core.report_generator import ReportGenerator, ImageSourceError, SlideLayoutError
from core.drive_agent import (
    DriveAgent, DriveAgentError, ClientFolderNotFound, NDRFolderNotFound,
    VNFolderNotFound, CSVTableAsset, NDRAsset,
    place_ndr_images, place_csv_tables, place_csv_tables_paginated,
    place_ndr_images_multi_sensor, place_csv_tables_multi_sensor,
    count_unique_kb_patches, count_unique_software_packages, count_inventory_systems,
    count_siem_installations, count_dfir_installations, count_xdr_installations,
    build_software_patches_csv, build_ms_patches_csv, replace_kb_count_in_slides,
    replace_ai_insight_in_slide, replace_endpoint_count_in_slide,
    copy_service_coverage_from_reference,
)
from core.patch_history import generate_patch_trend_chart_for_report
from core.table_optimizer import optimize_tables, cleanup_placeholder_tables
from core.slide_utils import delete_slide
from core.ai_insights import generate_cyber_insight
from core.nuclei_parser import (
    parse_nuclei_file, populate_vulnerability_slides,
    update_all_page_numbers, _find_slide_by_shape_name,
)
from create_template import build_multi_sensor_template

# ---------------------------------------------------------------------------
# App setup
# ---------------------------------------------------------------------------
app = Flask(__name__, static_folder="frontend", static_url_path="/static")
app.secret_key = os.getenv("SHM_SECRET_KEY", secrets.token_hex(32))
app.config["SESSION_COOKIE_HTTPONLY"] = True
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"
CORS(app, supports_credentials=True)

# ---------------------------------------------------------------------------
# Authentication
# ---------------------------------------------------------------------------
_ADMIN_USER = "admin"
_ADMIN_PASSWORD_HASH = generate_password_hash("Kx9#mP2vL$8nQw4jR6tY!")

# ---------------------------------------------------------------------------
# User store (JSON file)
# ---------------------------------------------------------------------------
_USERS_FILE = Path(os.getenv("SHM_BASE_DIR", ".")) / "users.json"
_CLEARANCE_CODE = "SHM"  # only accepted security clearance


def _load_users() -> dict:
    """Load registered users from JSON file."""
    if _USERS_FILE.exists():
        try:
            return json.loads(_USERS_FILE.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            return {}
    return {}


def _save_users(users: dict) -> None:
    """Persist registered users to JSON file."""
    _USERS_FILE.write_text(json.dumps(users, indent=2), encoding="utf-8")


def login_required(f):
    """Decorator that blocks unauthenticated access to protected routes."""
    @functools.wraps(f)
    def decorated(*args, **kwargs):
        if not session.get("authenticated"):
            # API requests get 401 JSON; browser requests get redirected
            if request.path.startswith("/api/"):
                return jsonify({"error": "Authentication required"}), 401
            return redirect("/login")
        return f(*args, **kwargs)
    return decorated

# Add response headers to handle downloads over HTTP
@app.after_request
def add_headers(response):
    """Add headers to allow downloads and reduce security warnings over HTTP."""
    # Allow downloads from same origin over HTTP
    if not response.headers.get('Content-Security-Policy'):
        response.headers['Content-Security-Policy'] = "default-src 'self' 'unsafe-inline'; connect-src 'self' blob:;"
    # Prevent caching of dynamic content
    if request.path.startswith('/api/'):
        response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '0'
    return response

LOG_LEVEL = os.getenv("SHM_LOG_LEVEL", "INFO").upper()
logging.basicConfig(
    level=LOG_LEVEL,
    format="%(asctime)s  %(name)-24s  %(levelname)-8s  %(message)s",
)
logger = logging.getLogger("shm.api")

BASE_DIR = Path(os.getenv("SHM_BASE_DIR", "."))
ASSETS_DIR = BASE_DIR / os.getenv("SHM_ASSETS_DIR", "assets")
OUTPUT_DIR = BASE_DIR / os.getenv("SHM_OUTPUT_DIR", "output")
TEMPLATE_DIR = BASE_DIR / os.getenv("SHM_TEMPLATE_DIR", "templates")
MANIFESTS_DIR = BASE_DIR / os.getenv("SHM_MANIFESTS_DIR", "manifests")

NDR_CACHE_DIR = BASE_DIR / os.getenv("SHM_NDR_CACHE", "assets/ndr_cache")

DEFAULT_TEMPLATE = "mantix4_report_template.pptx"
# Use template_only so report content comes from Drive (NDR images + CSV tables).
DEFAULT_MANIFEST = "template_only.csv"

# Client name mapping for filename generation
CLIENT_PREFIX_MAP = {
    "montereymech": ("WT", "Monterey"),
    "elephant": ("WT", "Elephant"),
    "cic_xerox": ("", "CIC"),
    "wiegmann": ("WT", "Wiegmann"),
    "saints": ("WT", "Saints"),
    "rdlr": ("", "RDLR"),
    "jorgensen_labs": ("", "Jorgensen"),
    "gatedrentals": ("", "Gated Rentals"),
    "cjfig": ("", "CJFIG"),
    "packard": ("", "Packard"),
}


def _format_filename(client_slug: str, start_date: str, end_date: str) -> str:
    """Format download filename according to pattern:
    'WT - M4 - Monterey Vuln Report - 8 December to 5 January - 2026'
    
    Args:
        client_slug: Client identifier (e.g. 'montereymech')
        start_date: Start date formatted as 'DD Month' or 'DD Month YYYY'
        end_date: End date formatted as 'DD Month' or 'DD Month YYYY'
    
    Returns:
        Formatted filename (without .pptx extension)
    """
    # Get client prefix and short name
    prefix, short_name = CLIENT_PREFIX_MAP.get(client_slug, ("", client_slug.title()))
    
    # Extract year from end_date (e.g. "5 January 2026" -> "2026")
    year_match = re.search(r'\b(20\d{2})\b', end_date)
    year = year_match.group(1) if year_match else ""
    
    # Clean dates: remove year from individual dates since we append it at end
    start_clean = re.sub(r'\s+20\d{2}', '', start_date).strip()
    end_clean = re.sub(r'\s+20\d{2}', '', end_date).strip()
    
    # Build filename parts
    parts = []
    if prefix:
        parts.append(prefix)
    parts.append("M4")
    parts.append(f"{short_name} Vuln Report")
    parts.append(f"{start_clean} to {end_clean}")
    if year:
        parts.append(year)
    
    return " - ".join(parts)


def _save_report_metadata(report_id: str, client_slug: str, client_label: str, 
                          start_date: str, end_date: str) -> None:
    """Save report metadata for later retrieval during download."""
    metadata = {
        "report_id": report_id,
        "client_slug": client_slug,
        "client_label": client_label,
        "start_date": start_date,
        "end_date": end_date,
        "generated_at": datetime.now().isoformat(),
    }
    metadata_path = OUTPUT_DIR / f"{report_id}.meta.json"
    with open(metadata_path, 'w', encoding='utf-8') as f:
        json.dump(metadata, f, indent=2)


def _load_report_metadata(report_id: str) -> dict | None:
    """Load report metadata from file."""
    metadata_path = OUTPUT_DIR / f"{report_id}.meta.json"
    if not metadata_path.exists():
        return None
    try:
        with open(metadata_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return None


def _resolve_credentials_path() -> Path | None:
    """Resolve Google service-account credentials path.

    Uses SHM_GOOGLE_CREDENTIALS if set; otherwise looks for any .json
    in the project's Credentials/ folder.
    """
    env_path = os.getenv("SHM_GOOGLE_CREDENTIALS")
    if env_path:
        p = Path(env_path)
        if p.exists():
            return p.resolve()
        return None
    creds_dir = BASE_DIR / "Credentials"
    if not creds_dir.is_dir():
        return None
    for f in sorted(creds_dir.glob("*.json")):
        return f.resolve()
    return None


# ---------------------------------------------------------------------------
# Google Drive agent (optional — skipped if no credentials configured)
# ---------------------------------------------------------------------------
_drive_agent: DriveAgent | None = None
_creds_path = _resolve_credentials_path()
if _creds_path:
    try:
        _drive_agent = DriveAgent(credentials_path=_creds_path)
        logger.info("Google Drive agent ready (credentials: %s)", _creds_path.name)
    except DriveAgentError as _e:
        logger.warning("Drive agent disabled: %s", _e)
else:
    logger.warning(
        "Drive agent disabled: no credentials. Set SHM_GOOGLE_CREDENTIALS or "
        "place a service-account .json in the Credentials/ folder."
    )


# ---------------------------------------------------------------------------
# Cover-slide personalisation
# ---------------------------------------------------------------------------
def _personalise_cover(pptx_path: Path, client_name: str,
                       start_date: str, end_date: str) -> None:
    """Open the generated PPTX and inject client/period into the cover slide.

    Targets shapes by name:
        client_name       → "[Client Name]"  → actual client name (before hyphen)
        reporting_period  → placeholder text  → "Reporting Period: X to Y"
        
    Year appears only once after the end date (e.g., "2 February to 7 February 2026").
    Client name displays only the part before the first hyphen (e.g., "Elephant - 1st of the month" → "Elephant").
    """
    prs = PptxPresentation(str(pptx_path))
    cover = prs.slides[0]
    
    # Format period: remove year from start date if present, keep year on end date
    # Input: "2 February 2026" to "7 February 2026"
    # Output: "Reporting Period: 2 February to 7 February 2026"
    start_clean = re.sub(r'\s+\d{4}', '', start_date).strip()
    
    # Ensure end_date has year, if not present extract from start_date or use current year
    if not re.search(r'\d{4}', end_date):
        year_match = re.search(r'\d{4}', start_date)
        if year_match:
            end_date = f"{end_date} {year_match.group()}"
        else:
            from datetime import datetime
            end_date = f"{end_date} {datetime.now().year}"
    
    period = f"Reporting Period: {start_clean} to {end_date}"
    
    # Extract only the part before the first hyphen for the cover slide
    # e.g., "Elephant - 1st of the month" → "Elephant"
    cover_client_name = client_name.split('-')[0].strip() if '-' in client_name else client_name

    for shape in cover.shapes:
        if shape.name == "client_name" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = cover_client_name
        elif shape.name == "reporting_period" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = period

    prs.save(str(pptx_path))


# ---------------------------------------------------------------------------
# Login / Logout
# ---------------------------------------------------------------------------
@app.route("/login")
def login_page():
    """Serve the login page. If already authenticated, redirect to app."""
    if session.get("authenticated"):
        return redirect("/")
    return send_file("frontend/login.html")


@app.route("/api/v1/login", methods=["POST"])
def api_login():
    """Authenticate user and create session."""
    data = request.get_json(silent=True) or {}
    username = data.get("username", "").strip()
    password = data.get("password", "")

    # Check hardcoded admin
    if username == _ADMIN_USER and check_password_hash(_ADMIN_PASSWORD_HASH, password):
        session["authenticated"] = True
        session["user"] = username
        logger.info("Login successful for user '%s' from %s", username, request.remote_addr)
        return jsonify({"status": "ok", "message": "Login successful"})

    # Check registered users
    users = _load_users()
    if username in users and check_password_hash(users[username]["password_hash"], password):
        session["authenticated"] = True
        session["user"] = username
        logger.info("Login successful for user '%s' from %s", username, request.remote_addr)
        return jsonify({"status": "ok", "message": "Login successful"})

    logger.warning("Failed login attempt for user '%s' from %s", username, request.remote_addr)
    return jsonify({"error": "Invalid username or password"}), 401


@app.route("/api/v1/logout", methods=["POST"])
def api_logout():
    """Clear session and log out."""
    user = session.get("user", "unknown")
    session.clear()
    logger.info("User '%s' logged out", user)
    return jsonify({"status": "ok", "message": "Logged out"})


@app.route("/register")
def register_page():
    """Serve the registration page. If already authenticated, redirect to app."""
    if session.get("authenticated"):
        return redirect("/")
    return send_file("frontend/register.html")


@app.route("/api/v1/register", methods=["POST"])
def api_register():
    """Create a new user account after clearance verification."""
    data = request.get_json(silent=True) or {}
    username = data.get("username", "").strip()
    password = data.get("password", "")
    clearance = data.get("clearance", "").strip()

    # Validate fields
    if not username or len(username) < 3:
        return jsonify({"error": "Username must be at least 3 characters"}), 400
    if not password or len(password) < 6:
        return jsonify({"error": "Password must be at least 6 characters"}), 400

    # Block reserved admin username
    if username.lower() == _ADMIN_USER:
        return jsonify({"error": "Username is not available"}), 409

    # Check if username already taken
    users = _load_users()
    if username in users:
        return jsonify({"error": "Username is already taken"}), 409

    # Verify security clearance — only "SHM" is accepted
    if clearance != _CLEARANCE_CODE:
        logger.warning(
            "Registration denied for '%s' — invalid clearance from %s",
            username, request.remote_addr,
        )
        return jsonify({"error": "Invalid clearance code. Access denied."}), 403

    # All checks passed — create user
    users[username] = {
        "password_hash": generate_password_hash(password),
        "created_at": datetime.now().isoformat(),
    }
    _save_users(users)
    logger.info("New user '%s' registered from %s", username, request.remote_addr)
    return jsonify({"status": "ok", "message": "Account created successfully"}), 201


@app.route("/api/v1/auth/check", methods=["GET"])
def auth_check():
    """Check if current session is authenticated."""
    if session.get("authenticated"):
        return jsonify({"authenticated": True, "user": session.get("user")})
    return jsonify({"authenticated": False}), 401


# ---------------------------------------------------------------------------
# Frontend
# ---------------------------------------------------------------------------
@app.route("/")
@login_required
def index():
    return send_file("frontend/index.html")


@app.route("/favicon.ico")
@login_required
def favicon():
    """Return 204 No Content for favicon to prevent 404 errors."""
    return "", 204


# ---------------------------------------------------------------------------
# Health
# ---------------------------------------------------------------------------
@app.route("/api/v1/health", methods=["GET"])
def health():
    return jsonify({"status": "ok", "engine": "SHM Slide Layout Engine v1.0"})


# ---------------------------------------------------------------------------
# Generate report
# ---------------------------------------------------------------------------
@app.route("/api/v1/generate", methods=["POST"])
@login_required
def generate_report():
    """Trigger report generation.

    Expects JSON body::

        {
            "client_name": "Acme Corp",
            "start_date":  "22 December",
            "end_date":    "26 January 2026",
            "manifest":    "sample_security_report.csv",  // optional
            "template":    "mantix4_report_template.pptx", // optional
            "report_id":   "custom-id"                     // optional
        }

    ``client_name``, ``start_date`` and ``end_date`` are required and will be
    injected into the cover slide of the generated report.  ``manifest`` and
    ``template`` default to the sample manifest and Mantix4 template.
    """
    data = request.get_json(silent=True) or {}

    # --- Required fields: client & period ----------------------------------
    client_name = data.get("client_name", "").strip()   # slug for Drive
    client_label = data.get("client_label", "").strip() or client_name  # display name for cover
    start_date = data.get("start_date", "").strip()
    end_date = data.get("end_date", "").strip()

    missing = []
    if not client_name:
        missing.append("client_name")
    if not start_date:
        missing.append("start_date")
    if not end_date:
        missing.append("end_date")
    if missing:
        return jsonify({
            "error": f"Missing required fields: {', '.join(missing)}"
        }), 400

    # --- Manifest (defaults to sample) -------------------------------------
    manifest_name = data.get("manifest", DEFAULT_MANIFEST)
    manifest_path = MANIFESTS_DIR / manifest_name
    if not manifest_path.exists():
        return jsonify({"error": f"Manifest not found: {manifest_name}"}), 404

    # --- Sensor detection & template selection --------------------------------
    template_name = data.get("template", DEFAULT_TEMPLATE)
    detected_sensors: list[str] = ["DEFAULT"]
    is_multi_sensor = False

    if _drive_agent:
        try:
            detected_sensors = _drive_agent.detect_sensors(client_slug=client_name)
            logger.info(
                "Sensor detection for '%s': %d sensor(s) found — %s",
                client_name, len(detected_sensors), detected_sensors,
            )
        except Exception as exc:
            logger.warning(
                "Sensor detection failed for '%s', defaulting to single sensor: %s",
                client_name, exc,
            )
            detected_sensors = ["DEFAULT"]

    # Determine if multi-sensor template is needed
    is_multi_sensor = (
        len(detected_sensors) > 1
        and detected_sensors != ["DEFAULT"]
    )

    # Template selection: generate custom template for multi-sensor, else use standard
    if is_multi_sensor:
        try:
            template_path = build_multi_sensor_template(detected_sensors)
            logger.info(
                "Generated multi-sensor template for %d sensors: %s",
                len(detected_sensors), template_path,
            )
        except Exception as exc:
            logger.error(
                "Multi-sensor template generation failed, falling back to standard: %s",
                exc, exc_info=True,
            )
            is_multi_sensor = False
            template_path = TEMPLATE_DIR / template_name
            if not template_path.exists():
                return jsonify({"error": f"Template not found: {template_name}"}), 404
    else:
        template_path = TEMPLATE_DIR / template_name
        if not template_path.exists():
            return jsonify({"error": f"Template not found: {template_name}"}), 404
        logger.info("Using standard template: %s", template_path)

    report_id = data.get("report_id", uuid.uuid4().hex[:12])
    output_path = OUTPUT_DIR / f"{report_id}.pptx"

    try:
        generator = ReportGenerator(
            template_path=template_path,
            output_path=output_path,
            assets_dir=ASSETS_DIR,
        )
        result_path = generator.generate(manifest_path)

        # Inject client display name and reporting period into cover slide
        _personalise_cover(result_path, client_label, start_date, end_date)

        # Auto-populate patch counts from VN folder CSVs and build
        # the deduplicated software-patches table CSV for slide 14.
        kb_count_populated = False
        kb_warning = None
        software_patches_asset: CSVTableAsset | None = None
        ms_patches_asset: CSVTableAsset | None = None
        patch_chart_asset: NDRAsset | None = None

        if _drive_agent:
            try:
                # KB count from remediation plan CSV
                kb_count = None
                remediation_csv = _drive_agent.fetch_remediation_csv(
                    client_name, NDR_CACHE_DIR,
                )
                if remediation_csv:
                    kb_count = count_unique_kb_patches(remediation_csv)
                    
                    # Build MS patches table from remediation plan
                    ms_csv, ms_kb_count = build_ms_patches_csv(
                        remediation_csv, NDR_CACHE_DIR / client_name,
                    )
                    if ms_csv:
                        ms_patches_asset = CSVTableAsset(
                            local_path=ms_csv,
                            slide_index=15,
                            shape_name="ms_patches_table",
                            original_filename="top_ms_patches.csv",
                        )
                        logger.info("Created MS patches table asset with %d unique KBs", ms_kb_count)

                # Software package count + table from vulnaggregateother / non_windows_updates CSV
                sw_count = None
                software_csv = _drive_agent.fetch_software_csv(
                    client_name, NDR_CACHE_DIR,
                )
                if software_csv:
                    clean_csv, sw_count = build_software_patches_csv(
                        software_csv, NDR_CACHE_DIR / client_name,
                    )
                    if clean_csv:
                        software_patches_asset = CSVTableAsset(
                            local_path=clean_csv,
                            slide_index=14,
                            shape_name="software_patches_table",
                            original_filename="top_software_patches.csv",
                        )

                if kb_count is not None:
                    replacements = replace_kb_count_in_slides(
                        result_path, kb_count, sw_count,
                    )
                    kb_count_populated = replacements > 0
                elif sw_count is not None:
                    replacements = replace_kb_count_in_slides(
                        result_path, 0, sw_count,
                    )
                    kb_count_populated = replacements > 0

                # Generate patch trend chart image
                try:
                    chart_ms = kb_count if kb_count is not None else 0
                    chart_sw = sw_count if sw_count is not None else 0
                    chart_output = NDR_CACHE_DIR / client_name / "patch_trend.png"
                    # Convert display date (e.g. "26 January 2026") to YYYY-MM-DD
                    try:
                        _end_dt = datetime.strptime(end_date, "%d %B %Y")
                        chart_end_date = _end_dt.strftime("%Y-%m-%d")
                    except ValueError:
                        chart_end_date = end_date  # fallback: let history manager handle it
                    chart_path = generate_patch_trend_chart_for_report(
                        client_slug=client_name,
                        end_date=chart_end_date,
                        microsoft_count=chart_ms,
                        software_count=chart_sw,
                        output_path=chart_output,
                        drive_agent=_drive_agent,
                    )
                    if chart_path:
                        # Slide 14 = index 13 in single-sensor template
                        patch_chart_asset = NDRAsset(
                            local_path=Path(chart_path),
                            slide_index=13,
                            shape_name="required_patches_chart",
                            original_filename="patch_trend.png",
                        )
                        logger.info(
                            "Patch trend chart generated: %s (MS=%d, SW=%d)",
                            chart_path, chart_ms, chart_sw,
                        )
                except Exception as exc:
                    logger.warning("Failed to generate patch trend chart: %s", exc)

                # Generate AI-powered cyber insight for slide 3
                try:
                    insight = generate_cyber_insight(
                        microsoft_kb_count=kb_count if kb_count is not None else 0,
                        software_count=sw_count if sw_count is not None else 0,
                        key_detections=0,  # Not yet extracted from data
                        soc_tickets=0,     # Not yet extracted from data
                        analyst_investigations=0,  # Not yet extracted from data
                        nuclei_vulns_count=0,  # Will be populated later if available
                    )
                    replaced = replace_ai_insight_in_slide(result_path, insight)
                    if replaced:
                        logger.info("AI insight generated and inserted on slide 3")
                except Exception as exc:
                    logger.warning("Failed to generate AI insight: %s", exc)
                
                # Copy service coverage slide content from reference PPTX
                try:
                    ref_pptx = _drive_agent.fetch_reference_pptx(client_name)
                    if ref_pptx:
                        copied = copy_service_coverage_from_reference(
                            result_path, ref_pptx,
                        )
                        if copied:
                            logger.info(
                                "Service coverage copied from reference PPTX: %s",
                                ref_pptx,
                            )
                        else:
                            logger.warning(
                                "Could not copy service coverage from reference PPTX"
                            )
                    else:
                        logger.info(
                            "No reference PPTX found for '%s'; using template content",
                            client_name,
                        )
                except Exception as exc:
                    logger.warning(
                        "Failed to copy service coverage from reference: %s", exc,
                    )

                # Fetch inventory report and populate endpoint count on slide 4
                try:
                    inventory_csv = _drive_agent.fetch_inventory_csv(
                        client_name, NDR_CACHE_DIR,
                    )
                    if inventory_csv:
                        endpoint_count = count_inventory_systems(inventory_csv)
                        siem_count = count_siem_installations(inventory_csv)
                        dfir_count = count_dfir_installations(inventory_csv)
                        xdr_count = count_xdr_installations(inventory_csv)
                        
                        if endpoint_count is not None:
                            replacements = replace_endpoint_count_in_slide(
                                result_path, endpoint_count, siem_count, dfir_count, xdr_count
                            )
                            if replacements > 0:
                                msg_parts = [f"{endpoint_count} systems"]
                                if siem_count is not None:
                                    if siem_count == 0:
                                        msg_parts.append("SIEM: Client Managed")
                                    else:
                                        msg_parts.append(f"{siem_count} SIEM agents")
                                if xdr_count is not None:
                                    if xdr_count == 0:
                                        msg_parts.append("XDR: Client Managed")
                                    else:
                                        msg_parts.append(f"{xdr_count} XDR agents")
                                if dfir_count is not None:
                                    if dfir_count == 0:
                                        msg_parts.append("DFIR: Client Managed")
                                    else:
                                        msg_parts.append(f"{dfir_count} DFIR agents")
                                logger.info(
                                    "Service coverage counts inserted on slide 4: %s",
                                    ", ".join(msg_parts),
                                )
                except Exception as exc:
                    logger.warning("Failed to populate service coverage counts: %s", exc)
            except (ClientFolderNotFound, VNFolderNotFound) as exc:
                kb_warning = str(exc)
            except Exception as exc:
                kb_warning = f"KB count error: {exc}"

        # Fetch NDR images from Google Drive and place on slides
        ndr_placed = 0
        ndr_warning = None
        csv_tables_placed = 0
        tables_optimized = 0
        csv_warning = None
        nuclei_vulns_populated = 0
        nuclei_warning = None

        if _drive_agent:
            if is_multi_sensor:
                # --- Multi-sensor path ---
                # Build sensor-to-slide mapping: each sensor gets 10 consecutive slides
                # (8 NDR data slides + 2 vulnerability/mitigation text slides)
                # But for image/CSV placement, we only map the first 8 NDR slides
                sensor_slide_map: dict[str, list[int]] = {}
                for i, sensor_id in enumerate(detected_sensors):
                    base = 5 + i * 10  # 0-based slide indices (10 slides per sensor)
                    # Only map the first 8 slides (NDR data slides) for placement
                    sensor_slide_map[sensor_id] = list(range(base, base + 8))
                logger.info("Sensor-to-slide mapping (NDR slides only): %s", sensor_slide_map)

                n_sensors = len(detected_sensors)

                # Fetch and place NDR images (multi-sensor)
                sensor_ndr_assets = {}  # populated below; used by post-processing
                try:
                    sensor_ndr_assets = _drive_agent.fetch_ndr_assets_multi_sensor(
                        client_slug=client_name,
                        sensor_ids=detected_sensors,
                        dest_dir=NDR_CACHE_DIR,
                    )
                    ndr_placed_map = place_ndr_images_multi_sensor(
                        result_path, sensor_ndr_assets, sensor_slide_map,
                    )
                    ndr_placed = sum(ndr_placed_map.values())
                    for sid, cnt in ndr_placed_map.items():
                        logger.info("Sensor '%s': placed %d NDR images", sid, cnt)
                    logger.info(
                        "Total NDR images placed (multi-sensor): %d for '%s'",
                        ndr_placed, client_name,
                    )
                except (ClientFolderNotFound, NDRFolderNotFound) as exc:
                    ndr_warning = str(exc)
                    logger.warning("NDR fetch skipped: %s", exc)
                except Exception as exc:
                    ndr_warning = f"NDR fetch error: {exc}"
                    logger.error("NDR fetch failed: %s", exc, exc_info=True)

                # Place patch trend chart on the Required Patches slide (multi-sensor)
                if patch_chart_asset:
                    try:
                        ms_patches_slide_idx = 5 + n_sensors * 10  # Required Patches slide
                        chart_asset_ms = NDRAsset(
                            local_path=patch_chart_asset.local_path,
                            slide_index=ms_patches_slide_idx,
                            shape_name="required_patches_chart",
                            original_filename="patch_trend.png",
                        )
                        chart_placed = place_ndr_images(result_path, [chart_asset_ms])
                        if chart_placed:
                            ndr_placed += chart_placed
                            logger.info(
                                "Placed patch trend chart on slide %d (multi-sensor)",
                                ms_patches_slide_idx,
                            )
                    except Exception as exc:
                        logger.warning("Failed to place patch trend chart (multi-sensor): %s", exc)

                # Fetch and place CSV tables (multi-sensor)
                try:
                    sensor_csv_assets = _drive_agent.fetch_csv_assets_multi_sensor(
                        client_slug=client_name,
                        sensor_ids=detected_sensors,
                        dest_dir=NDR_CACHE_DIR,
                    )

                    # Separate sensor-specific NDR CSVs from non-sensor patches CSVs.
                    # Patches CSVs are grouped under "DEFAULT" by the fetch method.
                    ndr_csv_assets: dict[str, list] = {}
                    patches_csv_assets: list = []

                    for sid, assets_list in sensor_csv_assets.items():
                        if sid == "DEFAULT":
                            patches_csv_assets.extend(assets_list)
                        else:
                            ndr_csv_assets[sid] = assets_list

                    # Inject deduplicated software-patches CSV from VN folder
                    # if the NDR folder didn't already provide one for slide 14.
                    if software_patches_asset:
                        has_sw_from_ndr = any(
                            a.slide_index == 14 for a in patches_csv_assets
                        )
                        if not has_sw_from_ndr:
                            patches_csv_assets.append(software_patches_asset)
                            logger.info(
                                "Injected VN-folder software patches CSV for slide 14 (multi-sensor)",
                            )
                    
                    # Inject MS patches CSV from VN folder (slide 15)
                    if ms_patches_asset:
                        has_ms_from_ndr = any(
                            a.slide_index == 15 for a in patches_csv_assets
                        )
                        if not has_ms_from_ndr:
                            patches_csv_assets.append(ms_patches_asset)
                            logger.info(
                                "Injected VN-folder MS patches CSV for slide 15 (multi-sensor)",
                            )

                    # Place sensor-specific NDR CSV tables
                    csv_placed_map = place_csv_tables_multi_sensor(
                        result_path, ndr_csv_assets, sensor_slide_map,
                    )
                    csv_tables_placed = sum(csv_placed_map.values())
                    for sid, cnt in csv_placed_map.items():
                        logger.info("Sensor '%s': placed %d CSV tables", sid, cnt)

                    # Place patches tables with adjusted slide indices.
                    # In multi-sensor template: patches start after all NDR sections.
                    # Each sensor has 10 slides (8 NDR + 2 vuln/mitigation)
                    # Required Patches text slide = 5 + N*10
                    # Software Patches table     = 5 + N*10 + 1
                    # MS Patches table           = 5 + N*10 + 2
                    patches_base = 5 + n_sensors * 10
                    idx_remap = {
                        14: patches_base + 1,
                        15: patches_base + 2,
                    }
                    for asset in patches_csv_assets:
                        old_idx = asset.slide_index
                        if old_idx in idx_remap:
                            asset.slide_index = idx_remap[old_idx]
                            logger.info(
                                "Remapped patches CSV '%s' slide %d -> %d",
                                asset.original_filename, old_idx, asset.slide_index,
                            )

                    if patches_csv_assets:
                        patches_placed, _ = place_csv_tables_paginated(
                            result_path, patches_csv_assets,
                        )
                        csv_tables_placed += patches_placed
                        logger.info("Placed %d patches CSV tables", patches_placed)

                    # Build correctly-adjusted assets for table optimizer.
                    # NDR CSV assets still carry base slide indices (5-12) and
                    # base shape names.  The optimizer needs the ACTUAL indices
                    # and sensor-suffixed shape names to find the tables.
                    from core.drive_agent import CSVTableAsset as _CSVTableAsset
                    optimizer_assets: list[_CSVTableAsset] = list(patches_csv_assets)

                    base_slides = [5, 6, 7, 8, 9, 10, 11, 12]
                    for sid, assets_list in ndr_csv_assets.items():
                        slide_indices = sensor_slide_map.get(sid, [])
                        for asset in assets_list:
                            if asset.slide_index in base_slides and slide_indices:
                                offset = base_slides.index(asset.slide_index)
                                if offset < len(slide_indices):
                                    actual_idx = slide_indices[offset]
                                else:
                                    actual_idx = asset.slide_index
                            else:
                                actual_idx = asset.slide_index

                            optimizer_assets.append(_CSVTableAsset(
                                local_path=asset.local_path,
                                slide_index=actual_idx,
                                shape_name=f"{asset.shape_name}_{sid.lower()}",
                                original_filename=asset.original_filename,
                                sensor_id=sid,
                            ))

                    tables_optimized = optimize_tables(
                        result_path, optimizer_assets,
                        ndr_slide_indices=set(
                            idx for indices in sensor_slide_map.values()
                            for idx in indices
                        ),
                        report_date=end_date,
                    )

                    logger.info(
                        "Total CSV tables placed (multi-sensor): %d (%d optimised) for '%s'",
                        csv_tables_placed, tables_optimized, client_name,
                    )
                except (ClientFolderNotFound, NDRFolderNotFound) as exc:
                    csv_warning = str(exc)
                    logger.warning("CSV table fetch skipped: %s", exc)
                except Exception as exc:
                    csv_warning = f"CSV table fetch error: {exc}"
                    logger.error("CSV table fetch failed: %s", exc, exc_info=True)

                # --- Post-processing: cleanup & heatmap deletion (multi-sensor) ---
                try:
                    # 1. Cleanup leftover template tables
                    pp_patches_base = 5 + n_sensors * 10
                    all_table_slides: set[int] = set()
                    for indices in sensor_slide_map.values():
                        # Table slides are offsets 2-7 within the 8 NDR slides (indices[2:8])
                        all_table_slides.update(indices[2:8])
                    all_table_slides.update({pp_patches_base + 1, pp_patches_base + 2})
                    cleanup_placeholder_tables(result_path, all_table_slides)

                    # 2. Delete unused heatmap slides (2nd heatmap per sensor)
                    slides_to_delete: list[int] = []
                    for sid, assets_list in sensor_ndr_assets.items():
                        has_heatmap_2 = any(
                            a.slide_index == 6 for a in assets_list
                        )
                        if not has_heatmap_2:
                            sensor_base = sensor_slide_map[sid][0]
                            slides_to_delete.append(sensor_base + 1)

                    if slides_to_delete:
                        slides_to_delete.sort(reverse=True)
                        prs = PptxPresentation(str(result_path))
                        for idx in slides_to_delete:
                            delete_slide(prs, idx)
                        prs.save(str(result_path))
                        logger.info(
                            "Deleted %d unused heatmap slides: %s",
                            len(slides_to_delete), slides_to_delete,
                        )
                except Exception as exc:
                    logger.error(
                        "Post-processing failed: %s", exc, exc_info=True,
                    )

                # --- Nuclei vulnerability population (multi-sensor) ---
                try:
                    nuclei_files = _drive_agent.fetch_nuclei_files_multi_sensor(
                        client_slug=client_name,
                        sensor_ids=detected_sensors,
                        dest_dir=NDR_CACHE_DIR,
                    )
                    if nuclei_files:
                        prs = PptxPresentation(str(result_path))
                        total_nuclei_inserted = 0

                        for sensor_id in detected_sensors:
                            suffix = sensor_id.lower()
                            nuclei_path = nuclei_files.get(sensor_id)
                            vulns = parse_nuclei_file(nuclei_path) if nuclei_path else []

                            inserted = populate_vulnerability_slides(
                                prs,
                                vulns,
                                vuln_shape_name=f"internal_vulns_content_{suffix}",
                                mit_shape_name=f"internal_mitigation_content_{suffix}",
                                vuln_intro_shape_name=f"internal_vulns_intro_{suffix}",
                                mit_intro_shape_name=f"internal_mitigation_intro_{suffix}",
                                sensor_suffix=sensor_id,
                            )
                            total_nuclei_inserted += inserted
                            if vulns:
                                nuclei_vulns_populated += 1

                        # Delete aggregate vuln/mitigation slides (non-sensor-specific)
                        # In multi-sensor reports, we only want per-sensor slides
                        agg_vuln_idx = _find_slide_by_shape_name(prs, "internal_vulns_content")
                        agg_mit_idx = _find_slide_by_shape_name(prs, "internal_mitigation_content")
                        
                        slides_to_delete = []
                        if agg_vuln_idx is not None:
                            slides_to_delete.append(agg_vuln_idx)
                            logger.info("Marked aggregate vulnerability slide (index %d) for deletion", agg_vuln_idx)
                        if agg_mit_idx is not None:
                            slides_to_delete.append(agg_mit_idx)
                            logger.info("Marked aggregate mitigation slide (index %d) for deletion", agg_mit_idx)
                        
                        # Delete in reverse order to maintain indices
                        if slides_to_delete:
                            for idx in sorted(slides_to_delete, reverse=True):
                                delete_slide(prs, idx)
                                logger.info("Deleted aggregate slide at index %d", idx)

                        update_all_page_numbers(prs)
                        prs.save(str(result_path))
                        logger.info(
                            "Nuclei population complete: %d sensors populated, "
                            "%d extra slides inserted, %d aggregate slides deleted for '%s'",
                            nuclei_vulns_populated, total_nuclei_inserted, len(slides_to_delete),
                            client_name,
                        )
                    else:
                        logger.info("No nuclei files found for any sensor of '%s'", client_name)
                except (ClientFolderNotFound, NDRFolderNotFound) as exc:
                    nuclei_warning = f"Nuclei fetch skipped: {exc}"
                    logger.warning("Nuclei fetch skipped: %s", exc)
                except Exception as exc:
                    nuclei_warning = f"Nuclei processing error: {exc}"
                    logger.error("Nuclei processing failed: %s", exc, exc_info=True)

            else:
                # --- Single-sensor path (original behaviour) ---
                assets = []  # populated below; used by post-processing
                try:
                    assets = _drive_agent.fetch_ndr_assets(
                        client_slug=client_name,
                        dest_dir=NDR_CACHE_DIR,
                    )
                except (ClientFolderNotFound, NDRFolderNotFound) as exc:
                    ndr_warning = str(exc)
                    logger.warning("NDR fetch skipped: %s", exc)
                except Exception as exc:
                    ndr_warning = f"NDR fetch error: {exc}"
                    logger.error("NDR fetch failed: %s", exc, exc_info=True)

                # Include patch trend chart alongside NDR images
                if patch_chart_asset:
                    assets.append(patch_chart_asset)

                if assets:
                    ndr_placed = place_ndr_images(result_path, assets)
                    logger.info("Placed %d NDR images for '%s'", ndr_placed, client_name)

                # Fetch CSV files and populate table slides
                try:
                    csv_assets = _drive_agent.fetch_csv_assets(
                        client_slug=client_name,
                        dest_dir=NDR_CACHE_DIR,
                    )

                    # Inject deduplicated software-patches CSV from VN folder
                    # if the NDR folder didn't already provide one for slide 14.
                    if software_patches_asset:
                        has_sw_from_ndr = any(
                            a.slide_index == 14 for a in csv_assets
                        )
                        if not has_sw_from_ndr:
                            csv_assets.append(software_patches_asset)
                            logger.info(
                                "Injected VN-folder software patches CSV for slide 14",
                            )
                    
                    # Inject MS patches CSV from VN folder (slide 15)
                    if ms_patches_asset:
                        has_ms_from_ndr = any(
                            a.slide_index == 15 for a in csv_assets
                        )
                        if not has_ms_from_ndr:
                            csv_assets.append(ms_patches_asset)
                            logger.info(
                                "Injected VN-folder MS patches CSV for slide 15",
                            )

                    csv_tables_placed, _ = place_csv_tables_paginated(
                        result_path, csv_assets
                    )
                    tables_optimized = optimize_tables(
                        result_path, csv_assets,
                        report_date=end_date,
                    )
                    logger.info(
                        "Populated %d CSV tables (%d optimised) for '%s'",
                        csv_tables_placed, tables_optimized, client_name,
                    )
                except (ClientFolderNotFound, NDRFolderNotFound) as exc:
                    csv_warning = str(exc)
                    logger.warning("CSV table fetch skipped: %s", exc)
                except Exception as exc:
                    csv_warning = f"CSV table fetch error: {exc}"
                    logger.error("CSV table fetch failed: %s", exc, exc_info=True)

                # --- Post-processing: cleanup & heatmap deletion (single-sensor) ---
                try:
                    # 1. Delete unused heatmap slide FIRST (slide 6 is the 2nd heatmap)
                    #    This must happen before cleanup, as deleting shifts all subsequent indices down by 1
                    slide_deleted = False
                    has_heatmap_2 = any(
                        a.slide_index == 6 for a in assets
                    )
                    if not has_heatmap_2:
                        prs = PptxPresentation(str(result_path))
                        delete_slide(prs, 6)
                        prs.save(str(result_path))
                        slide_deleted = True
                        logger.info("Deleted unused 2nd heatmap slide (index 6)")

                    # 2. Cleanup leftover template tables
                    #    If slide 6 was deleted, adjust indices down by 1 for slides after it
                    if slide_deleted:
                        cleanup_indices = {6, 7, 8, 9, 10, 11, 13, 14}  # Adjusted for slide deletion
                    else:
                        cleanup_indices = {7, 8, 9, 10, 11, 12, 14, 15}
                    cleanup_placeholder_tables(result_path, cleanup_indices)
                except Exception as exc:
                    logger.error(
                        "Post-processing failed: %s", exc, exc_info=True,
                    )

                # --- Nuclei vulnerability population (single-sensor) ---
                try:
                    nuclei_path = _drive_agent.fetch_nuclei_file(
                        client_slug=client_name,
                        dest_dir=NDR_CACHE_DIR,
                    )
                    prs = PptxPresentation(str(result_path))
                    vulns = parse_nuclei_file(nuclei_path) if nuclei_path else []

                    inserted = populate_vulnerability_slides(
                        prs,
                        vulns,
                        vuln_shape_name="internal_vulns_content",
                        mit_shape_name="internal_mitigation_content",
                        vuln_intro_shape_name="internal_vulns_intro",
                        mit_intro_shape_name="internal_mitigation_intro",
                    )
                    if vulns:
                        nuclei_vulns_populated = 1

                    update_all_page_numbers(prs)
                    prs.save(str(result_path))
                    logger.info(
                        "Nuclei population complete: %d vulns, %d extra slides for '%s'",
                        len(vulns), inserted, client_name,
                    )
                except (ClientFolderNotFound, NDRFolderNotFound) as exc:
                    nuclei_warning = f"Nuclei fetch skipped: {exc}"
                    logger.warning("Nuclei fetch skipped: %s", exc)
                except Exception as exc:
                    nuclei_warning = f"Nuclei processing error: {exc}"
                    logger.error("Nuclei processing failed: %s", exc, exc_info=True)
        else:
            ndr_warning = (
                "Drive integration not available: Google credentials not configured. "
                "NDR images were not fetched."
            )
            csv_warning = (
                "Drive integration not available: Google credentials not configured. "
                "CSV table data was not fetched."
            )

        # Save metadata for download filename generation
        _save_report_metadata(report_id, client_name, client_label, start_date, end_date)

        summary = {
            "total": len(generator.audit_trail),
            "ok": sum(1 for r in generator.audit_trail if r.status == "ok"),
            "errors": sum(1 for r in generator.audit_trail if r.status == "error"),
            "ndr_images": ndr_placed,
            "table_pages": csv_tables_placed,
            "tables_optimized": tables_optimized,
            "sensors_detected": detected_sensors,
            "multi_sensor": is_multi_sensor,
            "nuclei_vulns_populated": nuclei_vulns_populated,
            "kb_count_populated": kb_count_populated,
        }

        response = {
            "status": "completed",
            "report_id": report_id,
            "output_file": str(result_path),
            "audit_file": str(result_path.with_suffix(".audit.json")),
            "summary": summary,
        }
        if ndr_warning:
            response["ndr_warning"] = ndr_warning
        if csv_warning:
            response["csv_warning"] = csv_warning
        if nuclei_warning:
            response["nuclei_warning"] = nuclei_warning
        if kb_warning:
            response["kb_warning"] = kb_warning

        return jsonify(response)

    except FileNotFoundError as exc:
        logger.error("Generation failed: %s", exc)
        return jsonify({"error": str(exc)}), 404

    except (ImageSourceError, SlideLayoutError) as exc:
        logger.error("Generation failed: %s", exc)
        return jsonify({"error": str(exc)}), 422

    except Exception as exc:
        logger.exception("Unexpected error during generation")
        return jsonify({"error": f"Internal error: {exc}"}), 500


# ---------------------------------------------------------------------------
# Download report
# ---------------------------------------------------------------------------
@app.route("/api/v1/download/<report_id>", methods=["GET", "HEAD"])
@login_required
def download_report(report_id: str):
    """Download a generated .pptx by report ID with formatted filename."""
    output_path = OUTPUT_DIR / f"{report_id}.pptx"
    if not output_path.exists():
        return jsonify({"error": f"Report '{report_id}' not found."}), 404
    
    # HEAD request for checking if file exists
    if request.method == "HEAD":
        return "", 200
    
    # Load metadata to generate proper filename
    metadata = _load_report_metadata(report_id)
    if metadata:
        # Format: "WT - M4 - Monterey Vuln Report - 8 December to 5 January - 2026"
        download_name = _format_filename(
            metadata['client_slug'],
            metadata['start_date'],
            metadata['end_date']
        ) + ".pptx"
    else:
        # Fallback if metadata not found
        download_name = f"report_{report_id}.pptx"
    
    return send_file(
        str(output_path.resolve()),
        as_attachment=True,
        download_name=download_name,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )


# ---------------------------------------------------------------------------
# Audit trail
# ---------------------------------------------------------------------------
@app.route("/api/v1/audit/<report_id>", methods=["GET"])
@login_required
def get_audit(report_id: str):
    """Retrieve the JSON audit trail for a completed report."""
    audit_path = OUTPUT_DIR / f"{report_id}.audit.json"
    if not audit_path.exists():
        return jsonify({"error": f"Audit trail for '{report_id}' not found."}), 404
    return send_file(str(audit_path.resolve()), mimetype="application/json")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    port = int(os.getenv("SHM_PORT", "5000"))
    debug = os.getenv("SHM_DEBUG", "false").lower() == "true"
    app.run(host="0.0.0.0", port=port, debug=debug)
